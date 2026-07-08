from __future__ import annotations

import base64
import io
import os
import re
from typing import Iterable, Sequence

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    from services.slide_render_plan import FontProfile, SlideRenderPlan
except Exception:  # pragma: no cover - fallback when imported as package module
    from ..slide_render_plan import FontProfile, SlideRenderPlan


SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _pt(value: int) -> Pt:
    return Pt(value)


def _rgb(hex_color: str, default: str = "#1E293B") -> RGBColor:
    color = (hex_color or default).strip()
    if not re.match(r"^#[0-9a-fA-F]{6}$", color):
        color = default
    return RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))


def _component_spec(plan: SlideRenderPlan) -> dict:
    spec = getattr(plan, "component_spec", {}) or {}
    return spec if isinstance(spec, dict) else {}


def _render_hint(plan: SlideRenderPlan, key: str, default=True):
    hints = _component_spec(plan).get("render_hints") or {}
    if not isinstance(hints, dict):
        return default
    return hints.get(key, default)


def _mix_hex(hex_color: str, amount: float = 0.82) -> str:
    color = (hex_color or "#1E293B").strip()
    if not re.match(r"^#[0-9a-fA-F]{6}$", color):
        color = "#1E293B"
    r = int(color[1:3], 16)
    g = int(color[3:5], 16)
    b = int(color[5:7], 16)
    mixed = tuple(int(channel + (255 - channel) * amount) for channel in (r, g, b))
    return "#{:02X}{:02X}{:02X}".format(*mixed)


def _set_drawingml_font(run, font_profile: FontProfile | None, role: str = "body") -> None:
    if not font_profile:
        return
    latin = font_profile.latin_font or font_profile.body_font
    east_asian = font_profile.east_asian_font or font_profile.body_font
    if role == "title":
        latin = font_profile.title_font or latin
        east_asian = font_profile.title_font or east_asian

    run.font.name = latin
    r_pr = run._r.get_or_add_rPr()
    for tag, face in (("a:latin", latin), ("a:ea", east_asian), ("a:cs", latin)):
        element = r_pr.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            r_pr.append(element)
        element.set("typeface", face or font_profile.fallback_font)


def _apply_paragraph_font(paragraph, font_profile: FontProfile | None, role: str = "body") -> None:
    if not font_profile:
        return
    for run in paragraph.runs:
        _set_drawingml_font(run, font_profile, role)


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int,
    color="#1E293B",
    bold=False,
    align=None,
    font_profile: FontProfile | None = None,
    font_role: str = "body",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = _pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    if align is not None:
        p.alignment = align
    _apply_paragraph_font(p, font_profile, font_role)
    return box


def _add_body_text(
    slide,
    box_spec: Sequence[float],
    lines: Iterable[str],
    font_size: int,
    color: str = "#334155",
    font_profile: FontProfile | None = None,
):
    x, y, w, h = box_spec
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size if idx < 3 else max(12, font_size - 2))
        p.font.color.rgb = _rgb(color)
        p.space_after = Pt(8)
        _apply_paragraph_font(p, font_profile, "body")
    return box


def _add_card(slide, x, y, w, h, title: str, body: str, accent: str, soft: str = "#FFFFFF", font_profile: FontProfile | None = None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(soft, "#FFFFFF")
    card.line.color.rgb = _rgb(accent)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()

    _add_textbox(slide, x + 0.22, y + 0.16, w - 0.42, 0.34, title, 13, accent, True, font_profile=font_profile, font_role="title")
    _add_textbox(slide, x + 0.22, y + 0.58, w - 0.42, h - 0.72, body, 11, "#475569", font_profile=font_profile)
    return card


def _add_metric_card(slide, x, y, w, h, label: str, value: str, accent: str, soft: str = "#FFFFFF", font_profile: FontProfile | None = None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(soft, "#FFFFFF")
    card.line.color.rgb = _rgb(accent)
    _add_textbox(slide, x + 0.16, y + 0.18, w - 0.32, 0.36, label, 11, "#475569", True, font_profile=font_profile, font_role="title")
    _add_textbox(slide, x + 0.16, y + 0.78, w - 0.32, max(0.6, h - 1.0), value, 28, accent, True, font_profile=font_profile, font_role="title")
    return card


def _extract_metric_value(text: str) -> str:
    match = re.search(r"([+-]?\d+(?:\.\d+)?\s*%|[+-]?\d+(?:\.\d+)?\s*(?:万|亿|k|K|w|W)?)", text or "")
    return match.group(1).replace(" ", "") if match else (text or "")[:14]


def _add_video_link(slide, x, y, w, h, text: str, url: str, accent: str, font_profile: FontProfile | None = None):
    box = _add_textbox(slide, x, y, w, h, text, 12, accent, True, font_profile=font_profile, font_role="title")
    if url:
        run = box.text_frame.paragraphs[0].runs[0]
        run.hyperlink.address = url
    return box


def _placeholder_image_bytes(label: str) -> io.BytesIO:
    img = Image.new("RGB", (960, 540), "#E2E8F0")
    draw = ImageDraw.Draw(img)
    draw.rectangle((28, 28, 932, 512), outline="#94A3B8", width=6)
    draw.line((28, 512, 932, 28), fill="#CBD5E1", width=5)
    draw.line((28, 28, 932, 512), fill="#CBD5E1", width=5)
    draw.text((48, 46), label[:48] or "IMAGE PLACEHOLDER", fill="#334155")
    stream = io.BytesIO()
    img.save(stream, format="PNG")
    stream.seek(0)
    return stream


def _image_stream_from_source(image_source: str, output_path: str) -> io.BytesIO | str | None:
    if not image_source:
        return None

    source = image_source.strip()
    if source.startswith("data:image") and "," in source:
        try:
            return io.BytesIO(base64.b64decode(source.split(",", 1)[1]))
        except Exception:
            return None

    candidates = [source]
    if source.startswith("/api/generated/"):
        rel = source.split("/api/generated/", 1)[1]
        candidates.append(os.path.join(os.path.dirname(output_path), "generated", rel))
        candidates.append(os.path.join(os.path.dirname(output_path), rel))
        candidates.append(os.path.join(os.path.dirname(output_path), source.rsplit("/", 1)[-1]))
    if not os.path.isabs(source):
        candidates.append(os.path.join(os.path.dirname(output_path), source))
        candidates.append(os.path.abspath(source))

    for candidate in candidates:
        if candidate and os.path.exists(candidate):
            return candidate
    return None


def _has_resolved_image(image_source: str, output_path: str) -> bool:
    return _image_stream_from_source(image_source, output_path) is not None


def _add_asset_picture(slide, image_source: str, box_spec: Sequence[float], label: str, output_path: str):
    x, y, w, h = box_spec
    image_stream = _image_stream_from_source(image_source, output_path)
    if image_stream:
        return slide.shapes.add_picture(image_stream, Inches(x), Inches(y), Inches(w), Inches(h))
    return slide.shapes.add_picture(_placeholder_image_bytes(label), Inches(x), Inches(y), Inches(w), Inches(h))


def _add_rect(slide, x, y, w, h, fill: str, line: str | None = None, transparency: int = 0, rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = _rgb(line)
    else:
        shape.line.fill.background()
    return shape


def _add_numbered_dot(slide, x, y, number: int, accent: str, font_profile: FontProfile | None = None):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.34), Inches(0.34))
    dot.fill.solid()
    dot.fill.fore_color.rgb = _rgb(accent)
    dot.line.fill.background()
    _add_textbox(slide, x + 0.04, y + 0.06, 0.26, 0.18, str(number), 8, "#FFFFFF", True, PP_ALIGN.CENTER, font_profile, "title")
    return dot


def _add_visual_photo_field(
    slide,
    plan: SlideRenderPlan,
    box_spec: Sequence[float],
    label: str,
    output_path: str,
    evidence_grid: bool = False,
):
    x, y, w, h = box_spec
    accent = plan.accent
    soft = plan.soft or _mix_hex(accent, 0.82)
    image_source = plan.image_sources[0] if plan.image_sources else ""
    has_real_image = _has_resolved_image(image_source, output_path)

    _add_rect(slide, x, y, w, h, _mix_hex(accent, 0.86), accent, 8, rounded=True)
    for idx in range(5):
        band_y = y + 0.22 + idx * (h - 0.44) / 5
        _add_rect(slide, x + 0.12, band_y, max(0.2, w - 0.24), 0.12, soft, None, 46)

    _add_asset_picture(slide, image_source, (x + 0.1, y + 0.1, max(0.2, w - 0.2), max(0.2, h - 0.2)), label, output_path)
    _add_rect(slide, x, y, w, h, "#FFFFFF", accent, 100, rounded=True)

    if evidence_grid and not has_real_image:
        card_w = max(0.7, (w - 0.64) / 2)
        card_h = max(0.48, (h - 0.74) / 2)
        for row in range(2):
            for col in range(2):
                cx = x + 0.22 + col * (card_w + 0.2)
                cy = y + 0.24 + row * (card_h + 0.18)
                _add_rect(slide, cx, cy, card_w, card_h, "#FFFFFF", "#CBD5E1", 0, rounded=True)
                _add_rect(slide, cx + 0.12, cy + 0.13, max(0.1, card_w - 0.24), 0.12, _mix_hex(accent, 0.64))
                _add_rect(slide, cx + 0.12, cy + card_h - 0.2, max(0.1, card_w - 0.24), 0.04, "#CBD5E1")

    if not has_real_image and not evidence_grid:
        _add_textbox(slide, x + 0.22, y + 0.22, max(0.4, w - 0.44), 0.32, label, 12, "#334155", True, font_profile=plan.intent.font_profile, font_role="title")


def _add_step_flow_cards(slide, plan: SlideRenderPlan) -> None:
    spec = plan.layout_spec
    font_profile = plan.intent.font_profile
    centers = []
    for idx, box in enumerate(spec.card_boxes):
        x, y, w, h = box
        text = plan.card_texts[idx % len(plan.card_texts)] if plan.card_texts else spec.label
        _add_rect(slide, x, y, w, h, "#FFFFFF", plan.accent, 0, rounded=True)
        _add_numbered_dot(slide, x + 0.18, y + 0.18, idx + 1, plan.accent, font_profile)
        _add_textbox(slide, x + 0.62, y + 0.18, max(0.2, w - 0.8), max(0.3, h - 0.34), str(text)[:120], 12, "#334155", idx == 0, font_profile=font_profile)
        centers.append((x + w / 2, y + h / 2))

    for first, second in zip(centers, centers[1:]):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(first[0] + 1.25), Inches(first[1]), Inches(second[0] - 1.25), Inches(second[1]))
        line.line.color.rgb = _rgb(plan.accent)


def _add_campaign_timeline(slide, plan: SlideRenderPlan) -> None:
    spec = plan.layout_spec
    font_profile = plan.intent.font_profile
    y = 5.92
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.25), Inches(y), Inches(11.75), Inches(y))
    line.line.color.rgb = _rgb(plan.accent)
    line.line.width = Pt(2)

    for idx, box in enumerate(spec.card_boxes):
        x, card_y, w, h = box
        cx = x + w / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.08), Inches(y - 0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(plan.accent)
        dot.line.fill.background()
        text = plan.card_texts[idx % len(plan.card_texts)] if plan.card_texts else spec.label
        _add_rect(slide, x, card_y, w, h, "#FFFFFF", plan.accent, 0, rounded=True)
        _add_textbox(slide, x + 0.18, card_y + 0.18, w - 0.36, 0.28, f"阶段 {idx + 1}", 11, plan.accent, True, font_profile=font_profile, font_role="title")
        _add_textbox(slide, x + 0.18, card_y + 0.56, w - 0.36, h - 0.72, str(text)[:110], 11, "#334155", font_profile=font_profile)


def _add_section_marks(slide, plan: SlideRenderPlan) -> None:
    for idx, box in enumerate(plan.layout_spec.card_boxes):
        x, y, w, h = box
        fill = plan.accent if idx == 0 else plan.soft
        line = plan.accent
        _add_rect(slide, x, y, w, h, fill, line, 14 if idx else 0, rounded=True)
        if idx == 0:
            _add_textbox(slide, x + 0.18, y + 0.24, max(0.2, w - 0.36), 0.86, plan.intent.layout_label[:28], 10, "#FFFFFF", True, font_profile=plan.intent.font_profile, font_role="title")


def _add_metric_dashboard_cards(slide, plan: SlideRenderPlan) -> None:
    spec = plan.layout_spec
    for idx, card_box in enumerate(spec.card_boxes):
        text = plan.card_texts[idx % len(plan.card_texts)] if plan.card_texts else spec.label
        label = re.split(r"[:：\s]+", str(text), maxsplit=1)[0][:12] or f"指标 {idx + 1}"
        value = _extract_metric_value(str(text))
        _add_metric_card(slide, *card_box, label, value, plan.accent, "#FFFFFF", plan.intent.font_profile)


def _draw_archetype_native_objects(slide, plan: SlideRenderPlan, output_path: str) -> tuple[bool, bool]:
    spec = plan.layout_spec
    archetype = plan.visual_profile.archetype
    handled_image = False
    handled_cards = False

    if archetype in {"hero_photo_claim", "strategy_claim_collage", "video_material_board"} and spec.image_box:
        label = "视频缩略图 / DEMO" if archetype == "video_material_board" else "主视觉素材"
        _add_visual_photo_field(slide, plan, spec.image_box, label, output_path)
        handled_image = True
        if archetype == "strategy_claim_collage":
            for idx, size in enumerate((0.42, 0.32, 0.24)):
                _add_rect(slide, 10.7 + idx * 0.48, 0.82 + idx * 0.34, size, size, plan.soft, plan.accent, 0)

    if archetype == "evidence_wall" and spec.image_box:
        _add_visual_photo_field(slide, plan, spec.image_box, "证据截图 / 素材占位", output_path, evidence_grid=True)
        handled_image = True

    if archetype == "metric_dashboard":
        _add_metric_dashboard_cards(slide, plan)
        handled_cards = True
    elif archetype == "step_flow":
        _add_step_flow_cards(slide, plan)
        handled_cards = True
    elif archetype == "campaign_timeline":
        _add_campaign_timeline(slide, plan)
        handled_cards = True
    elif archetype == "section_divider":
        _add_section_marks(slide, plan)
        handled_cards = True

    return handled_image, handled_cards


def _draw_background(
    slide,
    prs: Presentation,
    accent: str,
    paper: str = "#F8FAFC",
    grid: bool = False,
    soft: str = "#E2E8F0",
    background_source: str = "",
    output_path: str = "",
    show_left_rail: bool = True,
):
    if background_source:
        _add_asset_picture(slide, background_source, (0, 0, SLIDE_W_IN, SLIDE_H_IN), "背景素材图", output_path)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(paper, "#F8FAFC")
    bg.fill.transparency = 18 if background_source else 0
    bg.line.fill.background()

    if show_left_rail:
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), prs.slide_height)
        rail.fill.solid()
        rail.fill.fore_color.rgb = _rgb(accent)
        rail.line.fill.background()

    if grid:
        for x in [1.2, 2.4, 3.6, 4.8, 6.0, 7.2, 8.4, 9.6, 10.8, 12.0]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(0.24), Inches(x), Inches(7.2))
            line.line.color.rgb = _rgb(soft, "#E2E8F0")
            line.line.transparency = 68
        for y in [1.2, 2.4, 3.6, 4.8, 6.0]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.42), Inches(y), Inches(12.68), Inches(y))
            line.line.color.rgb = _rgb(soft, "#E2E8F0")
            line.line.transparency = 72

    wash = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.85), Inches(-0.8), Inches(3.4), Inches(2.3))
    wash.fill.solid()
    wash.fill.fore_color.rgb = _rgb(soft, "#E2E8F0")
    wash.fill.transparency = 38
    wash.line.fill.background()


def _add_skill_tag(slide, label: str, accent: str):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.74), Inches(0.28), Inches(1.55), Inches(0.36))
    tag.fill.solid()
    tag.fill.fore_color.rgb = _rgb(accent)
    tag.line.fill.background()
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = _rgb("#FFFFFF")
    p.alignment = PP_ALIGN.CENTER
    return tag


def render_plan_slide(prs: Presentation, plan: SlideRenderPlan, output_path: str):
    spec = plan.layout_spec
    accent = plan.accent
    soft = plan.soft
    font_profile = plan.intent.font_profile
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background_source = plan.background_sources[0] if plan.background_sources else ""
    _draw_background(
        slide,
        prs,
        accent,
        plan.paper,
        plan.grid,
        soft,
        background_source,
        output_path,
        show_left_rail=bool(_render_hint(plan, "show_left_rail", True)),
    )
    if _render_hint(plan, "show_fixed_tag", True):
        _add_skill_tag(slide, spec.label, accent)

    if _render_hint(plan, "show_generic_separator", True):
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(0.74),
            Inches(spec.title_box[1] + spec.title_box[3] + 0.16),
            Inches(12.45),
            Inches(spec.title_box[1] + spec.title_box[3] + 0.16),
        )
        line.line.color.rgb = _rgb(accent)

    _add_textbox(slide, *spec.title_box, plan.title, 30, "#0F172A", True, font_profile=font_profile, font_role="title")
    _add_body_text(slide, spec.body_box, plan.body_lines, spec.body_font_size, "#233044", font_profile=font_profile)

    handled_image, handled_cards = _draw_archetype_native_objects(slide, plan, output_path)

    if spec.image_box and not handled_image:
        image_source = plan.image_sources[0] if plan.image_sources else ""
        _add_asset_picture(slide, image_source, spec.image_box, plan.image_label, output_path)

    for card_idx, card_box in enumerate(spec.card_boxes):
        if handled_cards:
            break
        text = plan.card_texts[card_idx % len(plan.card_texts)] if plan.card_texts else spec.label
        if spec.label == "数据结果":
            bits = re.split(r"[:：\s]+", text, maxsplit=1)
            metric_label = bits[0][:12] if bits else f"指标 {card_idx + 1}"
            metric_value = bits[1][:18] if len(bits) > 1 else text[:18]
            _add_metric_card(slide, *card_box, metric_label, metric_value, accent, soft, font_profile)
        else:
            _add_card(slide, *card_box, f"{spec.label} {card_idx + 1}", text[:72], accent, soft, font_profile)

    if plan.video_links:
        _add_video_link(slide, 7.16 if spec.label == "视频素材" else 9.28, 6.48, 3.1, 0.32, "打开视频素材链接", plan.video_links[0], accent, font_profile)

    notes = slide.notes_slide.notes_text_frame
    notes.text = "\n".join(plan.notes).strip()
    return slide


def render_plans_to_pptx(plans: Sequence[SlideRenderPlan], output_path: str) -> str:
    if not plans:
        raise ValueError("outline 中没有 pages 数据")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    for plan in plans:
        render_plan_slide(prs, plan, output_path)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path
