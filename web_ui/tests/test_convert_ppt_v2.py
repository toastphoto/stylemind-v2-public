import os
import sys
import base64
import json
import shutil
import uuid
import zipfile
from io import BytesIO
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
HAS_NODE_MODULES = os.path.isdir(os.path.join(REPO_ROOT, "node_modules"))
REFERENCE_TEMPLATE_FIXTURES = [
    os.path.join(REPO_ROOT, "reference_samples", "curated_templates", "stylemind_cleaned_reference_templates.pptx"),
    os.path.join(REPO_ROOT, "reference_samples", "curated_templates", "stylemind_cleaned_reference_templates_registry.json"),
]
DASHIAI_SEED_REGISTRY = os.path.join(
    REPO_ROOT,
    "reference_samples",
    "dashiai_theme_seed",
    "stylemind_dashiai_theme_seed_registry.json",
)
DASHIAI_PPT_SKILL_ROOT = os.environ.get("DASHIAI_PPT_SKILL_ROOT", "")
HAS_DASHIAI_PPT_SKILL = bool(DASHIAI_PPT_SKILL_ROOT and os.path.isdir(DASHIAI_PPT_SKILL_ROOT))

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))
import api_server as api_server_module
from api_server import app
from services.structured_pptx_renderer import LAYOUT_SPECS, build_structured_pptx, resolve_page_skill
from services.campaign_style_catalog import campaign_style_render_tokens
from services.slide_render_plan import SlideRenderPlan, VisualProfile, build_deck_render_plan, deck_render_plan_to_dict
from services.renderers.python_pptx_renderer import render_plans_to_pptx
from services.renderers.pptxgenjs_renderer import render_outline_to_pptxgenjs
from services.renderers.reference_template_renderer import render_outline_to_reference_template
from services.renderers.html_dom_renderer import render_outline_to_html_dom

# 1x1 transparent png
PNG_1X1 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9WZ0kAAAAASUVORK5CYII="


def require_node_renderer():
    if not HAS_NODE_MODULES:
        pytest.skip("Node renderer dependencies are not installed; run npm install to enable this test.")


def require_reference_templates():
    require_node_renderer()
    if not all(os.path.exists(path) for path in REFERENCE_TEMPLATE_FIXTURES):
        pytest.skip("Private reference-template fixtures are not bundled in the public release.")


def require_html_dom_runtime():
    require_node_renderer()
    if not HAS_DASHIAI_PPT_SKILL:
        pytest.skip("Set DASHIAI_PPT_SKILL_ROOT to run the optional HTML DOM transcription test.")


def require_dashiai_seed_registry():
    if not os.path.exists(DASHIAI_SEED_REGISTRY):
        pytest.skip("DashiAI seed registry is optional and is not bundled in the public release.")


def pptx_xml_counts(path):
    with zipfile.ZipFile(path) as zf:
        slide_names = sorted(name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        root = ET.fromstring(zf.read("ppt/presentation.xml"))
        ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        slide_ids = root.findall(".//p:sldId", ns)
        text_nodes = 0
        picture_tags = 0
        for name in slide_names:
            raw = zf.read(name).decode("utf-8", errors="ignore")
            text_nodes += raw.count("<a:t>")
            picture_tags += raw.count("<p:pic>")
    return {
        "slide_xml_files": len(slide_names),
        "presentation_slide_ids": len(slide_ids),
        "text_nodes": text_nodes,
        "picture_tags": picture_tags,
    }


def pptx_media_blobs(path):
    with zipfile.ZipFile(path) as zf:
        return [zf.read(name) for name in zf.namelist() if name.startswith("ppt/media/")]


def pptx_text_values(path):
    texts = []
    with zipfile.ZipFile(path) as zf:
        slide_names = []
        try:
            presentation = zf.read("ppt/presentation.xml").decode("utf-8", errors="ignore")
            rels = zf.read("ppt/_rels/presentation.xml.rels").decode("utf-8", errors="ignore")
            rel_map = {}
            for rel in ET.fromstring(rels):
                rel_id = rel.attrib.get("Id")
                target = rel.attrib.get("Target")
                if rel_id and target:
                    slide_path = target.replace("../", "ppt/")
                    if not slide_path.startswith("ppt/"):
                        slide_path = f"ppt/{slide_path}"
                    rel_map[rel_id] = slide_path
            ns = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            }
            root = ET.fromstring(presentation)
            for slide_id in root.findall(".//p:sldId", ns):
                rel_id = slide_id.attrib.get(qn("r:id"))
                if rel_id in rel_map:
                    slide_names.append(rel_map[rel_id])
        except Exception:
            slide_names = []
        if not slide_names:
            slide_names = sorted(name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        for name in slide_names:
            root = ET.fromstring(zf.read(name))
            for node in root.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                if node.text and node.text.strip():
                    texts.append(node.text.strip())
    return texts


def test_convert_ppt_v2_endpoint_exists():
    payload = {
        "images": [f"data:image/png;base64,{PNG_1X1}"],
        "outline": {"pages": [{"index": 1, "title": "t", "content": "c"}]},
        "options": {"mode": "v2", "dry_run": True},
    }
    with app.test_client() as client:
        r = client.post("/api/convert-ppt-v2", json=payload)
    assert r.status_code == 200


def test_upload_page_exposes_structured_pptx_export():
    with app.test_client() as client:
        r = client.get("/upload")

    assert r.status_code == 200
    html = r.get_data(as_text=True)
    assert "exportPptxBtn" in html
    assert "生成可编辑 PPTX" in html
    assert "生成结构预览" in html
    assert "AI 生成PNG兜底" in html
    assert "image-2 仅用于素材图" in html
    assert "PPTX 转写策略" in html
    assert "PPTX renderer" not in html
    assert "right-rail" in html
    assert "progress-active" in html
    assert "/api/export-pptx-structured" in html
    assert "/api/render-plan-preview" in html
    assert "workbench-active" in html
    assert "导出与预览" in html
    assert "交付检查" in html
    assert "风格设置" in html
    assert "添加本页素材图" in html
    assert "添加背景图" in html
    assert "AI 生成素材图" in html
    assert "AI 生成背景图" in html
    assert "批量补背景图" in html
    assert "pptxRendererSelect" in html
    assert "参考页复刻转写（实验）" in html
    assert "HTML 同源转写（实验）" in html
    assert "结构化绘制（通用）" in html
    assert "/api/upload-page-asset" in html
    assert "/api/generate-page-asset" in html
    assert "/api/generate-deck-assets" in html


def test_export_pptx_structured_creates_editable_objects():
    payload = {
        "outline": {
            "title": "结构化导出测试",
            "pages": [
                {
                    "index": 1,
                    "title": "开场定调",
                    "type": "opening",
                    "layout": "开场定调-观点海报",
                    "brief": "用一句核心判断建立提案气质",
                    "content": "用户心智正在从种草转向体验验证。\n品牌需要把内容资产变成可执行证据。",
                    "fixed_images": ["hero-product.png"],
                    "video_links": ["https://example.com/video"],
                }
            ],
        }
    }
    with app.test_client() as client:
        r = client.post("/api/export-pptx-structured", json=payload)

    assert r.status_code == 200
    body = r.get_json()
    assert body["status"] == "success"
    assert body["renderer"] == "python-pptx"
    assert body["pipeline"] == "outline_to_render_plan_to_native_pptx"
    assert body["preview_pipeline"] == "outline_to_render_plan_to_preview_png"
    assert body["fallback_preview_pipeline"] == "png_pdf_preview_only"
    assert body["uses_full_page_image_model"] is False
    assert "素材图" in body["asset_image_model_role"]
    assert "不走整页 image-2" in body["message"]

    pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", body["filename"]))
    render_plan_path = None
    if body.get("render_plan_path"):
        render_plan_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", "generated", "render_plans", os.path.basename(body["render_plan_path"])))
    try:
        prs = Presentation(pptx_path)
        slide = prs.slides[0]
        shapes = list(slide.shapes)

        assert any(shape.has_text_frame and "开场定调" in shape.text for shape in shapes)
        assert any(shape.has_text_frame and "用户心智" in shape.text for shape in shapes)
        assert any(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for shape in shapes)
        assert not any(shape.shape_type == MSO_SHAPE_TYPE.LINE for shape in shapes)
        assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in shapes)

        full_slide_pictures = [
            shape
            for shape in shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.width >= prs.slide_width * 0.9 and shape.height >= prs.slide_height * 0.9
        ]
        assert not full_slide_pictures
    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)
        if render_plan_path and os.path.exists(render_plan_path):
            os.remove(render_plan_path)


def test_render_plan_preview_endpoint_uses_same_structured_plan():
    payload = {
        "outline": {
            "title": "结构预览测试",
            "pages": [
                {
                    "index": 1,
                    "title": "一、需求回顾【一页】",
                    "type": "content",
                    "layout": "满版图片-全屏背景",
                    "content": "搜索提升与品牌声量沉淀。",
                },
                {
                    "index": 2,
                    "title": "核心主题【2页】",
                    "type": "content",
                    "layout": "满版图片-左图右文",
                    "content": "core idea：#尽情享受美的夏天",
                },
                {
                    "index": 3,
                    "title": ". 平台热度正在快速上升：",
                    "type": "chart",
                    "layout": "模块化-表格",
                    "content": "+723% 新发布笔记数",
                },
            ],
        }
    }

    with app.test_client() as client:
        r = client.post("/api/render-plan-preview", json=payload)

    assert r.status_code == 200
    body = r.get_json()
    assert body["status"] == "success"
    assert body["pipeline"] == "outline_to_slide_layout_spec_to_render_plan"
    assert body["preview_pipeline"] == "outline_to_render_plan_to_preview_png"
    assert body["final_pptx_pipeline"] == "outline_to_render_plan_to_native_pptx"
    assert body["uses_full_page_image_model"] is False
    assert body["role_distribution"] == {"开场定调": 1, "创意主张": 1, "数据结果": 1}
    assert len(body["previews"]) == 3
    assert body["previews"][0]["page_role"] == "开场定调"
    assert body["previews"][1]["page_role"] == "创意主张"
    assert body["previews"][2]["page_role"] == "数据结果"
    assert body["previews"][1]["visual_profile"]["archetype"] == "strategy_claim_collage"
    assert body["previews"][2]["visual_profile"]["archetype"] == "metric_dashboard"
    image = body["previews"][0]["image"]
    assert image.startswith("data:image/png;base64,")
    assert base64.b64decode(image.split(",", 1)[1]).startswith(b"\x89PNG")


def test_upload_page_asset_returns_generated_asset_url():
    with app.test_client() as client:
        r = client.post(
            "/api/upload-page-asset",
            data={"file": (BytesIO(base64.b64decode(PNG_1X1)), "page-asset.png")},
            content_type="multipart/form-data",
        )

    assert r.status_code == 200
    body = r.get_json()
    assert body["status"] == "success"
    assert body["asset_url"].startswith("/api/generated/page_assets/")
    assert body["asset_role"] == "image"
    assert body["filename"].endswith(".png")

    asset_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", "generated", "page_assets", body["filename"]))
    try:
        assert os.path.exists(asset_path)
    finally:
        if os.path.exists(asset_path):
            os.remove(asset_path)


def test_upload_page_asset_background_role_updates_outline_session():
    session_id = f"test-bg-{uuid.uuid4().hex}"
    api_server_module._uploaded_outline_cache[session_id] = {"outline": {"pages": [{"title": "背景页", "content": "背景图应进入底层图片对象。"}]}}
    body = {}
    try:
        with app.test_client() as client:
            r = client.post(
                "/api/upload-page-asset",
                data={
                    "file": (BytesIO(base64.b64decode(PNG_1X1)), "page-bg.png"),
                    "session_id": session_id,
                    "page_index": "0",
                    "asset_role": "background",
                },
                content_type="multipart/form-data",
            )

        assert r.status_code == 200
        body = r.get_json()
        assert body["status"] == "success"
        assert body["asset_role"] == "background"
        assert body["updated_backgrounds"] == [body["asset_url"]]
        assert api_server_module._uploaded_outline_cache[session_id]["outline"]["pages"][0]["background_images"] == [body["asset_url"]]
        assert "fixed_images" not in api_server_module._uploaded_outline_cache[session_id]["outline"]["pages"][0]
    finally:
        if body.get("filename"):
            asset_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", "generated", "page_assets", body["filename"]))
            if os.path.exists(asset_path):
                os.remove(asset_path)
        session_path = api_server_module._get_outline_path(session_id)
        if os.path.exists(session_path):
            os.remove(session_path)
        api_server_module._uploaded_outline_cache.pop(session_id, None)


def test_generate_page_asset_dry_run_updates_background_asset_layer():
    session_id = f"test-ai-bg-{uuid.uuid4().hex}"
    api_server_module._uploaded_outline_cache[session_id] = {"outline": {"pages": [{"title": "AI背景页", "content": "用文章主题生成底层背景图。"}]}, "images": []}
    body = {}
    try:
        with app.test_client() as client:
            r = client.post(
                "/api/generate-page-asset",
                json={
                    "session_id": session_id,
                    "page_index": 0,
                    "asset_role": "background",
                    "dry_run": True,
                    "style": {"description": "小红书 campaign 提案风格"},
                },
            )

        assert r.status_code == 200
        body = r.get_json()
        assert body["status"] == "success"
        assert body["asset_role"] == "background"
        assert body["provider"] == "dry_run"
        assert body["asset_url"].startswith("/api/generated/page_assets/")
        assert body["updated_backgrounds"] == [body["asset_url"]]
        page = api_server_module._uploaded_outline_cache[session_id]["outline"]["pages"][0]
        assert page["background_images"] == [body["asset_url"]]
        assert "fixed_images" not in page
        assert api_server_module._uploaded_outline_cache[session_id]["images"] == []
        assert "Produce only the asset layer" in body["prompt"]
    finally:
        if body.get("filename"):
            asset_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", "generated", "page_assets", body["filename"]))
            if os.path.exists(asset_path):
                os.remove(asset_path)
        session_path = api_server_module._get_outline_path(session_id)
        if os.path.exists(session_path):
            os.remove(session_path)
        api_server_module._uploaded_outline_cache.pop(session_id, None)


def test_generate_deck_assets_dry_run_updates_missing_backgrounds():
    session_id = f"test-deck-bg-{uuid.uuid4().hex}"
    api_server_module._uploaded_outline_cache[session_id] = {
        "outline": {
            "pages": [
                {"title": "需要背景", "content": "用文章主题生成底层背景图。"},
                {"title": "已有背景", "content": "已有图不应重复生成。", "background_images": ["/api/generated/page_assets/existing.png"]},
                {"title": "也需要背景", "content": "补齐第二张缺失背景。"},
            ]
        },
        "images": [],
    }
    body = {}
    try:
        with app.test_client() as client:
            r = client.post(
                "/api/generate-deck-assets",
                json={
                    "session_id": session_id,
                    "asset_role": "background",
                    "dry_run": True,
                    "missing_only": True,
                    "max_pages": 12,
                    "style": {"description": "小红书 campaign 提案风格"},
                },
            )

        assert r.status_code == 200
        body = r.get_json()
        assert body["status"] == "success"
        assert body["asset_role"] == "background"
        assert body["generated_count"] == 2
        assert body["skipped_count"] == 1
        store = api_server_module._uploaded_outline_cache[session_id]
        pages = store["outline"]["pages"]
        assert pages[0]["background_images"][0].startswith("/api/generated/page_assets/")
        assert pages[1]["background_images"] == ["/api/generated/page_assets/existing.png"]
        assert pages[2]["background_images"][0].startswith("/api/generated/page_assets/")
        assert all(item["provider"] == "dry_run" for item in body["generated"])
    finally:
        for item in body.get("generated", []):
            filename = item.get("filename")
            if filename:
                asset_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", "generated", "page_assets", filename))
                if os.path.exists(asset_path):
                    os.remove(asset_path)
        session_path = api_server_module._get_outline_path(session_id)
        if os.path.exists(session_path):
            os.remove(session_path)
        api_server_module._uploaded_outline_cache.pop(session_id, None)


def test_structured_pptx_embeds_data_url_assets():
    image_data_url = f"data:image/png;base64,{PNG_1X1}"
    tmp_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", ".pytest_tmp"))
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"data-url-asset-{uuid.uuid4().hex}.pptx")

    try:
        build_structured_pptx(
            {
                "pages": [
                    {
                        "title": "素材图测试",
                        "page_type": "opening",
                        "content": "素材图应作为可替换图片对象进入 PPTX。",
                        "fixed_images": [image_data_url],
                    }
                ]
            },
            output_path,
        )

        prs = Presentation(output_path)
        pictures = [shape for shape in prs.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]

        assert pictures
        assert pictures[0].image.blob == base64.b64decode(PNG_1X1)
    finally:
        if os.path.exists(output_path):
            os.remove(output_path)


def test_structured_pptx_embeds_background_data_url_as_bottom_picture():
    image_data_url = f"data:image/png;base64,{PNG_1X1}"
    tmp_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", ".pytest_tmp"))
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"background-data-url-{uuid.uuid4().hex}.pptx")

    try:
        build_structured_pptx(
            {
                "pages": [
                    {
                        "title": "背景图测试",
                        "page_type": "opening",
                        "content": "背景图应作为底层可替换图片对象，文字仍然可编辑。",
                        "background_images": [image_data_url],
                    }
                ]
            },
            output_path,
        )

        prs = Presentation(output_path)
        pictures = [shape for shape in prs.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]

        assert pictures
        assert pictures[0].image.blob == base64.b64decode(PNG_1X1)
        assert any(shape.has_text_frame and "背景图测试" in shape.text for shape in prs.slides[0].shapes)
    finally:
        if os.path.exists(output_path):
            os.remove(output_path)


def test_structured_pptx_uses_campaign_style_catalog():
    page = {
        "title": "小红书生活艺术节",
        "page_type": "creative_claim",
        "content": "围绕生活方式和创作者共创，形成更柔和的种草内容系统。",
    }
    tokens = campaign_style_render_tokens(page)
    assert tokens["id"] == "xhs_lifestyle_grid"

    tmp_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", ".pytest_tmp"))
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"campaign-style-{uuid.uuid4().hex}.pptx")

    try:
        build_structured_pptx({"pages": [page]}, output_path)
        prs = Presentation(output_path)
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "小红书生活方式网格" in notes
        assert "xhs_lifestyle_grid" in notes
    finally:
        if os.path.exists(output_path):
            os.remove(output_path)


def test_structured_pptx_keeps_feibo_page_skill_taxonomy():
    samples = [
        {"page_type": "opening", "layout_skill": "开场定调-观点海报"},
        {"page_type": "chapter_transition", "layout": "章节转场-大标题"},
        {"page_type": "content_bridge", "layout": "内容承接-左右分栏"},
        {"page_type": "creative_claim", "title": "创意主张"},
        {"page_type": "execution_plan", "layout": "执行打法-四步路径"},
        {"page_type": "case_evidence", "layout": "案例证据-截图佐证"},
        {"page_type": "data_result", "layout": "数据结果-指标卡"},
        {"page_type": "video_material", "layout": "视频素材-链接陈列"},
    ]

    resolved = [resolve_page_skill(page) for page in samples]

    assert resolved == list(LAYOUT_SPECS.keys())


def test_legacy_outline_labels_infer_feibo_page_roles():
    samples = [
        (
            {"title": "一、需求回顾【一页】", "type": "content", "layout": "满版图片-全屏背景"},
            1,
            "开场定调",
        ),
        (
            {"title": "三、传播规划", "type": "summary", "layout": "满版图片-全屏背景"},
            11,
            "章节转场",
        ),
        (
            {"title": "目的【1页】", "type": "summary", "layout": "满版图片-全屏背景"},
            14,
            "执行打法",
        ),
        (
            {"title": "step2. 邓峰发布改造视频，妹妹之家爆改美美之家【1页】", "type": "content", "layout": "模块化-步骤流程"},
            20,
            "执行打法",
        ),
        (
            {"title": ". 平台热度正在快速上升：", "type": "chart", "layout": "模块化-表格", "content": "+723% 新发布笔记数"},
            8,
            "数据结果",
        ),
        (
            {
                "title": "类型一：腰部科技达人",
                "type": "content",
                "layout": "模块化-卡片",
                "content": "内容参考：https://www.xiaohongshu.com/explore/xxx",
            },
            31,
            "案例证据",
        ),
        (
            {"title": "核心主题【2页】", "type": "content", "layout": "满版图片-左图右文", "content": "core idea：#尽情享受美的夏天"},
            3,
            "创意主张",
        ),
    ]

    resolved = [resolve_page_skill(page, idx) for page, idx, _ in samples]

    assert resolved == [expected for _, _, expected in samples]


def test_render_plan_contract_for_five_page_roles():
    pages = [
        {
            "title": "开场定调",
            "page_type": "opening",
            "content": "用一句明确判断建立提案气质。",
            "fixed_images": ["hero.png"],
            "background_images": ["hero-bg.png"],
            "font_profile": {"title_font": "京东朗正体", "body_font": "微软雅黑", "latin_font": "Arial", "east_asian_font": "微软雅黑"},
        },
        {"title": "章节转场", "page_type": "chapter_transition", "content": "进入传播策略。"},
        {"title": "创意主张", "page_type": "creative_claim", "content": "提出核心创意钩子。"},
        {"title": "执行打法", "page_type": "execution_plan", "content": "Step 1：预热\nStep 2：爆发\nStep 3：沉淀"},
        {"title": "案例证据", "page_type": "case_evidence", "content": "用真实案例和截图建立可信度。"},
    ]

    plans = build_deck_render_plan({"pages": pages})

    assert len(plans) == 5
    assert all(isinstance(plan, SlideRenderPlan) for plan in plans)
    assert all(isinstance(plan.visual_profile, VisualProfile) for plan in plans)
    assert [plan.intent.page_role for plan in plans] == ["开场定调", "章节转场", "创意主张", "执行打法", "案例证据"]
    assert [plan.visual_profile.archetype for plan in plans] == [
        "hero_photo_claim",
        "section_divider",
        "strategy_claim_collage",
        "campaign_timeline",
        "evidence_wall",
    ]
    assert plans[0].layout_spec.label == "开场定调"
    assert plans[0].image_sources == ("hero.png",)
    assert plans[0].background_sources == ("hero-bg.png",)
    assert plans[0].intent.font_profile.title_font == "京东朗正体"
    assert plans[0].intent.font_profile.east_asian_font == "微软雅黑"
    assert plans[0].component_spec["schema"] == "stylemind.component_spec.v1"
    assert plans[0].component_spec["source"] == "feibo_template_first"
    assert plans[0].component_spec["render_hints"]["show_fixed_tag"] is False
    assert plans[0].component_spec["render_hints"]["show_generic_separator"] is False
    assert "image-2" in "\n".join(plans[0].notes)
    assert "视觉原型" in "\n".join(plans[0].notes)


def test_render_plan_serializes_dashiai_seed_as_feibo_component_spec():
    payload = deck_render_plan_to_dict(
        {
            "pages": [
                {
                    "title": "数据结果页",
                    "page_role": "数据结果",
                    "layout_skill": "数据结果",
                    "content": "小红书生活方式传播数据增长，目标达到 95 分审美。",
                    "dashiai_seed_key": "theme05_page009",
                }
            ]
        }
    )

    component_spec = payload["plans"][0]["component_spec"]

    assert component_spec["schema"] == "stylemind.component_spec.v1"
    assert component_spec["source"] == "dashiai_theme_seed"
    assert component_spec["selected_seed_key"] == "theme05_page009"
    assert component_spec["selected_seed"]["key"] == "theme05_page009"
    assert component_spec["feibo_restyle"]["schema"] == "stylemind.feibo_restyle_policy.v1"
    assert "keep_fixed_corner_index_by_default" in component_spec["blocked_defaults"]
    assert component_spec["render_hints"]["show_left_rail"] is False
    assert component_spec["render_hints"]["show_fixed_tag"] is False
    assert component_spec["render_hints"]["show_generic_separator"] is False


def test_python_renderer_consumes_render_plan_and_writes_cjk_fonts():
    tmp_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", ".pytest_tmp"))
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"render-plan-fonts-{uuid.uuid4().hex}.pptx")
    pages = [
        {
            "title": "字体控制页",
            "page_type": "opening",
            "content": "标题和正文需要写入明确的中文字体。",
            "font_profile": {"title_font": "京东朗正体", "body_font": "微软雅黑", "latin_font": "Arial", "east_asian_font": "微软雅黑"},
        },
        {"title": "章节转场", "page_type": "chapter_transition", "content": "第二页保持同一个 RenderPlan 合同。"},
        {"title": "创意主张", "page_type": "creative_claim", "content": "第三页继续验证可编辑对象。"},
        {"title": "执行打法", "page_type": "execution_plan", "content": "第四页验证流程页。"},
        {"title": "案例证据", "page_type": "case_evidence", "content": "第五页验证案例页。"},
    ]

    try:
        build_structured_pptx({"pages": pages}, output_path)
        prs = Presentation(output_path)

        assert len(prs.slides) == 5
        notes = prs.slides[0].notes_slide.notes_text_frame.text
        assert "结构化 PPTX 导出：开场定调" in notes

        east_asian_faces = []
        latin_faces = []
        for shape in prs.slides[0].shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    r_pr = run._r.get_or_add_rPr()
                    ea = r_pr.find(qn("a:ea"))
                    latin = r_pr.find(qn("a:latin"))
                    if ea is not None and ea.get("typeface"):
                        east_asian_faces.append(ea.get("typeface"))
                    if latin is not None and latin.get("typeface"):
                        latin_faces.append(latin.get("typeface"))

        assert "京东朗正体" in east_asian_faces
        assert "微软雅黑" in east_asian_faces
        assert "Arial" in latin_faces
    finally:
        if os.path.exists(output_path):
            os.remove(output_path)


def test_python_renderer_can_render_plans_without_structured_wrapper():
    tmp_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", ".pytest_tmp"))
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"direct-render-plan-{uuid.uuid4().hex}.pptx")
    plans = build_deck_render_plan(
        {
            "pages": [
                {"title": "直接 RenderPlan", "page_type": "opening", "content": "python-pptx renderer 应只消费 RenderPlan。"},
                {"title": "执行链路", "page_type": "execution_plan", "content": "Step 1：构建计划\nStep 2：渲染对象"},
            ]
        }
    )

    try:
        render_plans_to_pptx(plans, output_path)
        prs = Presentation(output_path)

        assert len(prs.slides) == 2
        assert any(shape.has_text_frame and "直接 RenderPlan" in shape.text for shape in prs.slides[0].shapes)
        assert any(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for shape in prs.slides[1].shapes)
    finally:
        if os.path.exists(output_path):
            os.remove(output_path)


def test_export_pptx_structured_renders_all_feibo_skill_pages():
    pages = [
        {
            "index": idx,
            "title": label,
            "page_type": spec.key,
            "layout_skill": label,
            "content": f"{label} 内容要点\n第二条说明",
            "fixed_images": [f"{spec.key}.png"],
            "video_links": [f"https://example.com/{spec.key}"],
        }
        for idx, (label, spec) in enumerate(LAYOUT_SPECS.items(), start=1)
    ]

    with app.test_client() as client:
        r = client.post("/api/export-pptx-structured", json={"outline": {"pages": pages}})

    assert r.status_code == 200
    body = r.get_json()
    pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", body["filename"]))
    try:
        prs = Presentation(pptx_path)

        assert len(prs.slides) == len(LAYOUT_SPECS)
        for slide, label in zip(prs.slides, LAYOUT_SPECS.keys()):
            assert any(shape.has_text_frame and label in shape.text for shape in slide.shapes)
    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)


def test_export_pptx_structured_can_use_pptxgenjs_renderer():
    require_node_renderer()
    pages = [
        {
            "title": "PptxGenJS 导出",
            "page_type": "opening",
            "content": "同一个 RenderPlan 可以走 Node renderer，仍然输出可编辑对象。",
            "background_images": [f"data:image/png;base64,{PNG_1X1}"],
        },
        {
            "title": "执行打法",
            "page_type": "execution_plan",
            "content": "Step 1：建立 RenderPlan\nStep 2：PptxGenJS 渲染",
        },
    ]

    with app.test_client() as client:
        r = client.post("/api/export-pptx-structured", json={"outline": {"pages": pages}, "renderer": "pptxgenjs"})

    assert r.status_code == 200
    body = r.get_json()
    assert body["renderer"] == "pptxgenjs"
    assert body["pipeline"] == "outline_to_render_plan_to_native_pptx"
    assert body["uses_full_page_image_model"] is False
    pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", body["filename"]))
    try:
        prs = Presentation(pptx_path)
        pictures = [shape for slide in prs.slides for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        full_slide_pictures = [
            shape for shape in pictures
            if abs(shape.left) < 1000 and abs(shape.top) < 1000 and shape.width >= prs.slide_width * 0.98 and shape.height >= prs.slide_height * 0.98
        ]

        assert len(prs.slides) == 2
        assert any(shape.has_text_frame and "PptxGenJS 导出" in shape.text for shape in prs.slides[0].shapes)
        assert pictures
        assert len(full_slide_pictures) <= 1
    finally:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)


def test_export_pptx_structured_can_use_reference_template_renderer():
    require_reference_templates()
    pages = [
        {
            "title": "一、需求回顾【一页】",
            "page_type": "opening",
            "layout": "满版图片-全屏背景",
            "content": "搜索提升与品牌声量沉淀。",
            "background_images": [f"data:image/png;base64,{PNG_1X1}"],
        },
        {
            "title": "核心主题【2页】",
            "page_type": "creative_claim",
            "layout_skill": "创意主张",
            "content": "core idea：#尽情享受美的夏天",
            "fixed_images": [f"data:image/png;base64,{PNG_1X1}"],
        },
    ]

    with app.test_client() as client:
        r = client.post(
            "/api/export-pptx-structured",
            json={"outline": {"pages": pages}, "renderer": "reference-template"},
        )

    assert r.status_code == 200
    body = r.get_json()
    assert body["renderer"] == "reference-template"
    assert body["quality_strategy"] == "reference_template_first"
    assert body["template_match_status"] == "cleaned_reference_templates_experimental_requires_visual_qa"
    assert body["uses_full_page_image_model"] is False
    assert body["reference_template_selected_pages"] == 2
    assert body["reference_template_report_path"].startswith("/api/generated/render_plans/")
    assert body["known_limitations"]

    pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", body["filename"]))
    render_plan_path = os.path.abspath(
        os.path.join(
            os.path.dirname(__file__),
            "..",
            "static",
            "generated",
            "render_plans",
            os.path.basename(body["render_plan_path"]),
        )
    )
    report_path = os.path.abspath(
        os.path.join(
            os.path.dirname(__file__),
            "..",
            "static",
            "generated",
            "render_plans",
            os.path.basename(body["reference_template_report_path"]),
        )
    )
    try:
        counts = pptx_xml_counts(pptx_path)
        assert counts["presentation_slide_ids"] == 2
        assert counts["text_nodes"] >= 2
        assert counts["picture_tags"] >= 1
        assert base64.b64decode(PNG_1X1) in pptx_media_blobs(pptx_path)
        with open(report_path, encoding="utf-8") as f:
            report = json.load(f)
        assert report["strategy"] == "reference_template_first"
        assert report["selectionMode"] == "all"
        assert report["cleanedLibraryUsed"] is True
        assert report["registryUsed"] is True
        assert report["pictureReplacementSummary"]["slidesWithAssets"] == 2
        assert report["pictureReplacementSummary"]["totalReplacements"] >= 2
        assert all(item["matchSource"] == "cleaned-library" for item in report["selected"])
        assert all(item["templateReadiness"] == "cleaned_named_placeholders" for item in report["selected"])
        assert all(item["pictureReplacement"]["replacements"] >= 1 for item in report["selected"])
    finally:
        for path in (pptx_path, render_plan_path, report_path):
            if os.path.exists(path):
                os.remove(path)


def test_pptxgenjs_renderer_wrapper_writes_render_plan_and_editable_pptx():
    require_node_renderer()
    outline = {
        "pages": [
            {
                "title": "PptxGenJS Renderer Wrapper",
                "page_type": "creative_claim",
                "layout_skill": "创意主张",
                "content": "renderer service 负责写 RenderPlan JSON，再调用 Node 输出可编辑对象。",
                "fixed_images": [f"data:image/png;base64,{PNG_1X1}"],
            },
            {
                "title": "底层背景图也是 PPTX picture",
                "page_type": "case_evidence",
                "layout_skill": "案例证据",
                "content": "背景图和素材图都必须进入 RenderPlan，并作为可替换图片对象出现。",
                "background_images": [f"data:image/png;base64,{PNG_1X1}"],
            },
        ]
    }
    repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    tmp_dir = os.path.join(repo_root, ".pytest_tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"pptxgenjs-wrapper-{uuid.uuid4().hex}.pptx")
    render_plan_path = os.path.join(tmp_dir, f"pptxgenjs-wrapper-{uuid.uuid4().hex}.render_plan.json")
    try:
        result = render_outline_to_pptxgenjs(outline, output_path, render_plan_path=render_plan_path)

        assert result.output_path == output_path
        assert result.render_plan_path == render_plan_path
        assert result.page_count == 2
        assert os.path.exists(render_plan_path)
        with open(render_plan_path, encoding="utf-8") as f:
            payload = json.load(f)
        assert payload["schema"] == "stylemind.render_plan.v1"
        assert payload["page_count"] == 2

        prs = Presentation(output_path)
        pictures = [shape for slide in prs.slides for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        full_slide_pictures = [
            shape
            for shape in pictures
            if abs(shape.left) < 1000 and abs(shape.top) < 1000 and shape.width >= prs.slide_width * 0.98 and shape.height >= prs.slide_height * 0.98
        ]

        assert len(prs.slides) == 2
        assert any(shape.has_text_frame and "PptxGenJS Renderer Wrapper" in shape.text for shape in prs.slides[0].shapes)
        assert any(shape.has_text_frame and "背景图" in shape.text for shape in prs.slides[1].shapes)
        assert pictures
        assert len(full_slide_pictures) <= 1
    finally:
        for path in (output_path, render_plan_path):
            if os.path.exists(path):
                os.remove(path)


def test_reference_template_renderer_wrapper_writes_report_and_editable_pptx():
    require_reference_templates()
    outline = {
        "pages": [
            {
                "title": "参考模板优先",
                "page_type": "opening",
                "layout_skill": "开场定调",
                "content": "先继承真实参考页，再替换安全槽位。",
                "fixed_images": [f"data:image/png;base64,{PNG_1X1}"],
            }
        ]
    }
    repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    tmp_dir = os.path.join(repo_root, ".pytest_tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"reference-template-wrapper-{uuid.uuid4().hex}.pptx")
    render_plan_path = os.path.join(tmp_dir, f"reference-template-wrapper-{uuid.uuid4().hex}.render_plan.json")
    report_path = os.path.join(tmp_dir, f"reference-template-wrapper-{uuid.uuid4().hex}.report.json")
    try:
        result = render_outline_to_reference_template(
            outline,
            output_path,
            render_plan_path=render_plan_path,
            report_path=report_path,
        )

        assert result.output_path == output_path
        assert result.render_plan_path == render_plan_path
        assert result.report_path == report_path
        assert result.page_count == 1
        assert result.selected_count == 1
        assert result.strategy == "reference_template_first"
        assert result.known_limitations

        counts = pptx_xml_counts(output_path)
        assert counts["presentation_slide_ids"] == 1
        assert counts["text_nodes"] >= 1
        assert base64.b64decode(PNG_1X1) in pptx_media_blobs(output_path)
        assert os.path.exists(report_path)
        with open(report_path, encoding="utf-8") as f:
            report = json.load(f)
        assert report["strategy"] == "reference_template_first"
        assert report["cleanedLibraryUsed"] is True
        assert report["registryUsed"] is True
        assert report["pictureReplacementSummary"]["slidesWithAssets"] == 1
        assert report["pictureReplacementSummary"]["totalReplacements"] >= 1
        assert report["selected"][0]["matchSource"] == "cleaned-library"
        assert report["selected"][0]["templateReadiness"] == "cleaned_named_placeholders"
        assert report["selected"][0]["templateFitScore"] is not None
        assert isinstance(report["selected"][0]["templateQualityTags"], list)
        assert report["selected"][0]["pictureReplacement"]["replacements"] >= 1
        assert not any("{{" in text or "}}" in text for text in pptx_text_values(output_path))
    finally:
        for path in (output_path, render_plan_path, report_path):
            if os.path.exists(path):
                os.remove(path)


def test_reference_template_renderer_replaces_missing_assets_with_generated_placeholders():
    require_reference_templates()
    outline = {
        "pages": [
            {
                "title": "无当前素材也不能保留旧参考图",
                "page_type": "opening",
                "layout_skill": "开场定调",
                "content": "缺少上传图时，参考模板图片槽应替换为当前页面生成的无文字视觉层。",
            }
        ]
    }
    repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    tmp_dir = os.path.join(repo_root, ".pytest_tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"reference-template-generated-placeholder-{uuid.uuid4().hex}.pptx")
    render_plan_path = os.path.join(tmp_dir, f"reference-template-generated-placeholder-{uuid.uuid4().hex}.render_plan.json")
    report_path = os.path.join(tmp_dir, f"reference-template-generated-placeholder-{uuid.uuid4().hex}.report.json")
    try:
        render_outline_to_reference_template(
            outline,
            output_path,
            render_plan_path=render_plan_path,
            report_path=report_path,
        )
        counts = pptx_xml_counts(output_path)
        assert counts["presentation_slide_ids"] == 1
        assert counts["picture_tags"] >= 1
        with zipfile.ZipFile(output_path) as zf:
            media_names = [name for name in zf.namelist() if name.startswith("ppt/media/")]
        assert any("stylemind_generated_placeholder" in name for name in media_names)
        with open(report_path, encoding="utf-8") as f:
            report = json.load(f)
        summary = report["pictureReplacementSummary"]
        replacement = report["selected"][0]["pictureReplacement"]
        assert summary["slidesWithAssets"] == 0
        assert summary["slidesWithGeneratedFallbacks"] == 1
        assert summary["generatedFallbackAssets"] >= 1
        assert summary["totalReplacements"] >= 1
        assert replacement["currentAssetCount"] == 0
        assert replacement["generatedFallbackAssetCount"] >= 1
        assert replacement["generatedFallbackTone"] in {"light", "dark"}
        assert replacement["replacements"] >= 1
        assert replacement["skipped"] == ""
        assert report["selected"][0]["templateFitScore"] is not None
        assert isinstance(report["selected"][0]["templateQualityTags"], list)
        assert not any("{{" in text or "}}" in text for text in pptx_text_values(output_path))
    finally:
        for path in (output_path, render_plan_path, report_path):
            if os.path.exists(path):
                os.remove(path)


def test_html_dom_renderer_wrapper_transcribes_render_plan_to_editable_pptx():
    require_html_dom_runtime()
    outline = {
        "pages": [
            {
                "title": "HTML DOM 转写",
                "page_role": "数据结果",
                "layout_skill": "数据结果",
                "content": "主指标：95%\n支撑指标：3.8x\nHTML 预览需要转写成可编辑 PPTX。",
                "dashiai_seed_key": "theme05_page009",
            },
            {
                "title": "案例证据层",
                "page_role": "案例证据",
                "layout_skill": "案例证据",
                "content": "复杂视觉可以局部截图，但文字必须保留可编辑对象。",
            },
        ]
    }
    repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
    tmp_dir = os.path.join(repo_root, ".pytest_tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    output_path = os.path.join(tmp_dir, f"html-dom-wrapper-{uuid.uuid4().hex}.pptx")
    render_plan_path = os.path.join(tmp_dir, f"html-dom-wrapper-{uuid.uuid4().hex}.render_plan.json")
    report_path = os.path.join(tmp_dir, f"html-dom-wrapper-{uuid.uuid4().hex}.report.json")
    deck_dir = os.path.join(tmp_dir, f"html-dom-wrapper-deck-{uuid.uuid4().hex}")
    try:
        result = render_outline_to_html_dom(
            outline,
            output_path,
            render_plan_path=render_plan_path,
            report_path=report_path,
            deck_dir=deck_dir,
        )

        assert result.output_path == output_path
        assert result.render_plan_path == render_plan_path
        assert result.report_path == report_path
        assert result.deck_dir == deck_dir
        assert result.page_count == 2
        assert result.strategy == "html_dom_transcription"
        assert result.text_objects >= 8
        assert result.image_objects >= 1

        counts = pptx_xml_counts(output_path)
        assert counts["presentation_slide_ids"] == 2
        assert counts["text_nodes"] >= 8
        assert counts["picture_tags"] >= 1
        assert os.path.exists(os.path.join(deck_dir, "index.html"))
        with open(report_path, encoding="utf-8") as f:
            report = json.load(f)
        assert report["slideCount"] == 2
        assert report["textObjects"] >= 8
    finally:
        for path in (output_path, render_plan_path, report_path):
            if os.path.exists(path):
                os.remove(path)
        if os.path.isdir(deck_dir):
            shutil.rmtree(deck_dir)


def test_export_pptx_structured_can_use_html_dom_renderer():
    require_html_dom_runtime()
    pages = [
        {
            "title": "HTML DOM API",
            "page_role": "数据结果",
            "layout_skill": "数据结果",
            "content": "主指标：95%\n支撑指标：3.8x",
            "dashiai_seed_key": "theme05_page009",
        }
    ]

    with app.test_client() as client:
        r = client.post("/api/export-pptx-structured", json={"outline": {"pages": pages}, "renderer": "html-dom"})

    assert r.status_code == 200
    body = r.get_json()
    assert body["renderer"] == "html-dom"
    assert body["quality_strategy"] == "html_dom_transcription"
    assert body["html_dom_status"] == "experimental_html_preview_to_editable_pptx_requires_visual_qa"
    assert body["uses_full_page_image_model"] is False
    assert body["html_dom_text_objects"] >= 4
    assert body["html_dom_image_objects"] >= 1
    assert body["html_dom_report_path"].startswith("/api/generated/render_plans/")
    assert body["html_dom_deck_path"].startswith("/api/generated/render_plans/")

    pptx_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", body["filename"]))
    render_plan_path = os.path.abspath(
        os.path.join(
            os.path.dirname(__file__),
            "..",
            "static",
            "generated",
            "render_plans",
            os.path.basename(body["render_plan_path"]),
        )
    )
    report_path = os.path.abspath(
        os.path.join(
            os.path.dirname(__file__),
            "..",
            "static",
            "generated",
            "render_plans",
            os.path.basename(body["html_dom_report_path"]),
        )
    )
    deck_dir = os.path.abspath(
        os.path.join(
            os.path.dirname(__file__),
            "..",
            "static",
            "generated",
            "render_plans",
            os.path.basename(os.path.dirname(body["html_dom_deck_path"])),
        )
    )
    try:
        counts = pptx_xml_counts(pptx_path)
        assert counts["presentation_slide_ids"] == 1
        assert counts["text_nodes"] >= 4
        assert counts["picture_tags"] >= 1
        with open(report_path, encoding="utf-8") as f:
            report = json.load(f)
        assert report["slideCount"] == 1
        assert report["textObjects"] >= 4
    finally:
        for path in (pptx_path, render_plan_path, report_path):
            if os.path.exists(path):
                os.remove(path)
        if os.path.isdir(deck_dir):
            shutil.rmtree(deck_dir)


def test_agent_run_exposes_dashiai_seed_choice_and_feibo_restyle_policy():
    require_dashiai_seed_registry()
    outline = {
        "title": "DashiAI seed choice smoke",
        "pages": [
            {
                "title": "数据结果页",
                "page_role": "数据结果",
                "layout_skill": "数据结果",
                "content": "小红书生活方式传播数据增长，目标达到 95 分审美。",
                "dashiai_seed_key": "theme05_page009",
            },
            {
                "title": "纯飞博参考页",
                "page_role": "创意主张",
                "layout_skill": "创意主张",
                "content": "拒绝 DashiAI seed，只用飞博参考模板。",
                "dashiai_seed_rejected": True,
            },
        ],
    }

    with app.test_client() as client:
        r = client.post("/api/agent/run", json={"action": "borrow_dashiai_component_seed", "outline": outline})

    assert r.status_code == 200
    body = r.get_json()
    assert body["status"] == "success"
    assert body["dashiai_seed_summary"]["counts"]["themes"] >= 12

    selected_page, rejected_page = body["pages"]
    assert selected_page["selected_dashiai_seed_key"] == "theme05_page009"
    assert selected_page["dashiai_theme_seeds"]
    assert selected_page["feibo_restyle"]["schema"] == "stylemind.feibo_restyle_policy.v1"
    assert selected_page["component_spec"]["selected_seed_key"] == "theme05_page009"
    assert selected_page["component_spec"]["render_hints"]["show_fixed_tag"] is False
    assert "keep_original_dashiai_palette_as_final" in selected_page["feibo_restyle"]["blocked_defaults"]
    assert selected_page["dashiai_theme_seeds"][0]["feibo_restyle"]["target_style_label"]

    assert rejected_page["dashiai_seed_rejected"] is True
    assert rejected_page["dashiai_theme_seeds"] == []
    assert rejected_page["feibo_restyle"]["status"] == "dashiai_seed_rejected_use_feibo_template_first"
    assert rejected_page["component_spec"]["dashiai_seed_rejected"] is True
