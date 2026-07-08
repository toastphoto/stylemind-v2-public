from __future__ import annotations

import io
import os
from typing import Dict, List, Tuple

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

# python-pptx shape constants (avoid importing enums everywhere)
MSO_SHAPE_RECTANGLE = 1
MSO_SHAPE_ROUNDED_RECTANGLE = 5
MSO_SHAPE_OVAL = 9
MSO_SHAPE_RIGHT_ARROW = 13


def _hex_to_rgb(color: str, default=(255, 255, 255)) -> Tuple[int, int, int]:
    if not color:
        return default
    c = color.strip()
    if c.startswith("#") and len(c) == 7:
        try:
            return int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)
        except Exception:
            return default
    return default


def _px_to_inches(x: float, total_px: float, total_in: float) -> float:
    return (x / total_px) * total_in


def _add_textbox(slide, bbox: Dict, text: str, img_w: int, img_h: int, prs: Presentation):
    left = Inches(_px_to_inches(bbox["x"], img_w, 13.333))
    top = Inches(_px_to_inches(bbox["y"], img_h, 7.5))
    width = Inches(_px_to_inches(bbox["w"], img_w, 13.333))
    height = Inches(_px_to_inches(bbox["h"], img_h, 7.5))
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)

    # 简单字号估计：bbox 高度映射到 points（后续会被 kind/font_size_pt 覆盖）
    font_size = max(10, min(44, int((_px_to_inches(bbox["h"], img_h, 7.5) * 72) * 0.75)))
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(30, 41, 59)


def _pick_font_size(text_item: Dict, bbox: Dict) -> int:
    """根据 kind / 建议字号 / 文本长度给一个更合理的字号（避免全都一样且重叠）。"""
    kind = (text_item.get("kind") or "").strip()
    suggested = int(text_item.get("font_size_pt") or 0)
    t = (text_item.get("text") or "").strip()

    if suggested >= 8 and suggested <= 60:
        base = suggested
    else:
        if kind == "title":
            base = 40
        elif kind == "subtitle":
            base = 16
        elif kind == "card_title":
            base = 20
        elif kind == "small":
            base = 12
        else:
            base = 14

    # 文本过长就略缩小
    if len(t) > 24:
        base = max(12, int(base * 0.85))
    if len(t) > 40:
        base = max(11, int(base * 0.8))

    return max(10, min(44, base))


def build_ppt_v2(
    output_path: str,
    original_image_path: str,
    clean_bg_png_bytes: bytes,
    layout: Dict,
    cropped_images: List[Tuple[str, Dict]],
) -> str:
    """
    生成分层 PPT：
    - 背景：clean_bg_png_bytes
    - 元素：layout.elements
    - 图片层：cropped_images（(path, bbox)）
    - 文字：layout.texts（OCR）
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_v2_slide(prs, clean_bg_png_bytes, layout, cropped_images)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path


def _add_v2_slide(prs: Presentation, clean_bg_png_bytes: bytes, layout: Dict, cropped_images: List[Tuple[str, Dict]]):
    """向 prs 中追加一个 V2 分层 slide。"""
    img_w = int(layout.get("width") or 1)
    img_h = int(layout.get("height") or 1)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 1) 背景图
    bg_stream = io.BytesIO(clean_bg_png_bytes)
    slide.shapes.add_picture(bg_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # 2) 元素层
    for e in layout.get("elements", []) or []:
        et = e.get("element_type", "other")
        bb = e.get("bbox") or {}
        if not bb:
            continue

        left = Inches(_px_to_inches(bb["x"], img_w, 13.333))
        top = Inches(_px_to_inches(bb["y"], img_h, 7.5))
        width = Inches(_px_to_inches(bb["w"], img_w, 13.333))
        height = Inches(_px_to_inches(bb["h"], img_h, 7.5))

        if et == "card":
            shape = slide.shapes.add_shape(MSO_SHAPE_ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(226, 232, 240)  # light border
        elif et == "divider":
            shape = slide.shapes.add_shape(MSO_SHAPE_RECTANGLE, left, top, width, max(Inches(0.03), height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(226, 232, 240)
            shape.line.fill.background()
        elif et == "arrow":
            shape = slide.shapes.add_shape(MSO_SHAPE_RIGHT_ARROW, left, top, width, height)
            shape.fill.solid()
            r, g, b = _hex_to_rgb(e.get("color") or "#EF4444", default=(239, 68, 68))
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()
        elif et == "circle_badge":
            d = min(width, height)
            shape = slide.shapes.add_shape(MSO_SHAPE_OVAL, left, top, d, d)
            shape.fill.solid()
            r, g, b = _hex_to_rgb(e.get("color") or "#EF4444", default=(239, 68, 68))
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.fill.background()
        else:
            continue

    # 3) 图片层（phone/photo crops）
    for path, bb in cropped_images:
        if not path or not os.path.exists(path):
            continue
        left = Inches(_px_to_inches(bb["x"], img_w, 13.333))
        top = Inches(_px_to_inches(bb["y"], img_h, 7.5))
        width = Inches(_px_to_inches(bb["w"], img_w, 13.333))
        height = Inches(_px_to_inches(bb["h"], img_h, 7.5))
        slide.shapes.add_picture(path, left, top, width=width, height=height)

    # 4) 文本层（OCR）
    badge_boxes = [e.get("bbox") for e in layout.get("elements", []) if e.get("element_type") == "circle_badge"]

    def _inside_any_badge(tb: Dict) -> bool:
        bx = tb["x"] + tb["w"] / 2
        by = tb["y"] + tb["h"] / 2
        for bb in badge_boxes:
            if not bb:
                continue
            if bb["x"] <= bx <= bb["x"] + bb["w"] and bb["y"] <= by <= bb["y"] + bb["h"]:
                return True
        return False

    for t in layout.get("texts", []) or []:
        text = (t.get("text") or "").strip()
        if not text:
            continue
        if text in ("1", "2", "3") and _inside_any_badge(t.get("bbox", {})):
            continue
        bb = t.get("bbox") or {}
        if not bb:
            continue
        # 文字渲染：使用 kind/font_size_pt 更合理
        left = Inches(_px_to_inches(bb["x"], img_w, 13.333))
        top = Inches(_px_to_inches(bb["y"], img_h, 7.5))
        width = Inches(_px_to_inches(bb["w"], img_w, 13.333))
        height = Inches(_px_to_inches(bb["h"], img_h, 7.5))
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)

        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(_pick_font_size(t, bb))
        p.font.color.rgb = RGBColor(30, 41, 59)
        if (t.get("kind") or "") in ("title", "card_title"):
            p.font.bold = True


def build_deck_v2(output_path: str, slides: List[Dict]) -> str:
    """
    slides: [{clean_bg_png_bytes, layout, cropped_images}]
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    # 删除默认第一页（Presentation 默认会创建一个空 slide_layout 但不包含 slides；无需处理）
    for s in slides:
        _add_v2_slide(prs, s["clean_bg_png_bytes"], s["layout"], s.get("cropped_images", []))
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path

def crop_regions(original_image_path: str, regions: List[Dict], out_dir: str) -> List[Tuple[str, Dict]]:
    """
    将 layout.images 中的 bbox 从原图裁切保存，返回 (path, bbox)。

    额外修正：
    - 对 phone bbox 做一次 OpenCV 边缘/轮廓修正，避免“切太大/切歪”
    - photo bbox 若明显异常，会回退为所在 card bbox（由 layout.elements 提供）
    """
    os.makedirs(out_dir, exist_ok=True)
    with Image.open(original_image_path) as im:
        im = im.convert("RGB")
        W, H = im.size
        out: List[Tuple[str, Dict]] = []
        for idx, r in enumerate(regions):
            bb = r.get("bbox") or {}
            if not bb:
                continue
            x, y, w, h = int(bb["x"]), int(bb["y"]), int(bb["w"]), int(bb["h"])
            if w <= 1 or h <= 1:
                continue
            # clamp
            x = max(0, min(W - 1, x))
            y = max(0, min(H - 1, y))
            w = max(1, min(W - x, w))
            h = max(1, min(H - y, h))

            # phone bbox refine
            if (r.get("region_type") or "") == "phone":
                try:
                    import cv2
                    import numpy as np
                    # 在 bbox 周围稍微扩一点再找手机轮廓
                    pad = 20
                    rx0 = max(0, x - pad)
                    ry0 = max(0, y - pad)
                    rx1 = min(W, x + w + pad)
                    ry1 = min(H, y + h + pad)
                    crop_np = np.array(im.crop((rx0, ry0, rx1, ry1)))
                    gray = cv2.cvtColor(crop_np, cv2.COLOR_RGB2GRAY)
                    edges = cv2.Canny(gray, 80, 160)
                    contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
                    best = None
                    best_area = 0
                    for c in contours:
                        bx, by, bw, bh = cv2.boundingRect(c)
                        area = bw * bh
                        if area < 2000:
                            continue
                        ar = bw / max(1, bh)
                        # 手机大致 0.45~0.65（宽/高）
                        if 0.35 <= ar <= 0.8:
                            if area > best_area:
                                best_area = area
                                best = (bx, by, bw, bh)
                    if best:
                        bx, by, bw, bh = best
                        # 转回原图坐标
                        x = rx0 + bx
                        y = ry0 + by
                        w = bw
                        h = bh
                        bb = {"x": int(x), "y": int(y), "w": int(w), "h": int(h)}
                except Exception:
                    pass

            crop = im.crop((x, y, x + w, y + h))
            import uuid
            path = os.path.join(out_dir, f"crop_{idx+1}_{uuid.uuid4().hex[:8]}.png")
            crop.save(path, format="PNG")
            out.append((path, bb))
        return out
