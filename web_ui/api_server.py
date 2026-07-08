#!/usr/bin/env python3
"""
StyleMind Web UI Backend API
Flask 后端 API 服务
"""
import os
import json
import re
import tempfile
import time
import uuid

# 临时目录：优先使用 /run（某些环境下更适配），不可写则回退到 /tmp
_preferred_tmp = "/run"
_fallback_tmp = "/tmp"
_tmpdir = _preferred_tmp if os.path.isdir(_preferred_tmp) and os.access(_preferred_tmp, os.W_OK) else _fallback_tmp
os.environ["TMPDIR"] = _tmpdir
os.environ["TEMP"] = _tmpdir
os.environ["TMP"] = _tmpdir
tempfile.tempdir = _tmpdir

from flask import Flask, request, jsonify, send_from_directory, Response, stream_with_context
from flask_cors import CORS
from urllib.parse import urljoin
import socket

# 导入新的提示词
from prompt_v2 import SKELETON_PROMPT_V2

# 导入 SSE 流大纲生成 V3（长章节自动分批）
from api_server_sse_fix import generate_outline_stream_v3 as generate_outline_stream_v2

app = Flask(__name__, static_folder='.')
CORS(app)

# 配置路径
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_FILE = os.path.join(BASE_DIR, '..', 'config.json')
GENERATED_DIR = os.path.join(BASE_DIR, 'static', 'generated')

# 确保路径存在
os.makedirs(os.path.dirname(CONFIG_FILE), exist_ok=True)
os.makedirs(GENERATED_DIR, exist_ok=True)

@app.route('/api/generated/<path:filename>')
def get_generated_file(filename):
    """访问生成的图片文件（用于刷新后仍可预览 & 供 convert-ppt 下载）"""
    return send_from_directory(GENERATED_DIR, filename)


@app.route('/api/server-info')
def server_info():
    """返回可在局域网/本机访问的 URL（用于在 Safari 里打开）"""
    try:
        # 获取一个“对外可达”的本机 IP（常用技巧：UDP connect 不实际发包）
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        try:
            s.connect(("8.8.8.8", 80))
            ip = s.getsockname()[0]
        finally:
            s.close()
        port = request.host.split(":")[-1] if ":" in request.host else "8080"
        return jsonify(
            {
                "status": "success",
                "lan_url": f"http://{ip}:{port}/",
                "local_url": "http://localhost:8080/",
            }
        )
    except Exception as e:
        return jsonify({"status": "error", "message": str(e)}), 500

# 初始化配置
if not os.path.exists(CONFIG_FILE):
    default_config = {
        "api_base_url": "",
        "api_key": "",
        "chat_model": "chatgpt-4o-latest",
        "image_model": "gpt-image-2",
        "embed_model": "text-embedding-3-small"
    }
    with open(CONFIG_FILE, 'w') as f:
        json.dump(default_config, f, indent=2)

# ==================== API 路由 ====================

@app.route('/')
def index():
    """返回主页"""
    return send_from_directory('.', 'index.html')

@app.route('/upload')
def upload_page():
    """返回大纲上传页面"""
    return send_from_directory('.', 'upload.html')

@app.route('/workbench')
def workbench_page():
    """返回 Agent 工作台页面"""
    return send_from_directory('.', 'workbench.html')

@app.route('/api/settings', methods=['GET', 'POST'])
def settings():
    """获取或保存设置（支持三个独立 API）"""
    if request.method == 'GET':
        try:
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)
            return jsonify(config)
        except Exception as e:
            return jsonify({'status': 'error', 'message': str(e)}), 500

    elif request.method == 'POST':
        try:
            new_config = request.json

            # 读取现有配置（兼容旧配置）
            try:
                with open(CONFIG_FILE, 'r') as f:
                    config = json.load(f)
            except:
                config = {}

            # 更新配置
            config.update(new_config)

            # 向后兼容：如果设置了新的独立 API，同时更新旧字段
            if 'chat_api_base_url' in new_config:
                config['api_base_url'] = new_config['chat_api_base_url']
            if 'chat_api_key' in new_config:
                config['api_key'] = new_config['chat_api_key']

            with open(CONFIG_FILE, 'w') as f:
                json.dump(config, f, indent=2)
            return jsonify({'status': 'success', 'message': '设置已保存'})
        except Exception as e:
            return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/test-connection', methods=['POST'])
def test_connection():
    """测试 API 连接"""
    try:
        data = request.json
        api_base_url = data.get('api_base_url')
        api_key = data.get('api_key')

        # 导入并测试 API
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        from core.api_client import APIClient

        api = APIClient(api_base_url, api_key)
        model = data.get('model', 'gpt-4o')
        resp = api.chat(model=model, messages=[{'role': 'user', 'content': 'hi'}], max_tokens=20, timeout=15)

        return jsonify({
            'status': 'success',
            'message': f'连接成功！API 响应: {resp["content"][:50]}'
        })
    except Exception as e:
        return jsonify({
            'status': 'error',
            'message': f'连接失败: {str(e)[:100]}'
        }), 500

@app.route('/api/import', methods=['POST'])
def import_file():
    """导入文件到知识库（SSE 流式进度）"""

    # 手动解析 multipart，完全绕过 Werkzeug 的 tempfile
    content_type = request.headers.get('Content-Type', '')
    raw_data = request.get_data()
    print(f"[DEBUG] 收到上传请求, Content-Type: {content_type[:80]}, 数据大小: {len(raw_data)} bytes")

    # 从 Content-Type 提取 boundary
    boundary = None
    for part in content_type.split(';'):
        part = part.strip()
        if part.startswith('boundary='):
            boundary = part.split('=', 1)[1].strip('"')
            break

    if not boundary:
        return jsonify({'status': 'error', 'message': '无效的请求格式'}), 400

    # 解析 multipart 数据，手动提取文件
    boundary_bytes = boundary.encode('utf-8')
    delimiter = b'--' + boundary_bytes
    parts = raw_data.split(delimiter)

    filename = None
    file_data = None

    for part in parts:
        if b'Content-Disposition' not in part:
            continue

        # 提取 filename
        header_end = part.find(b'\r\n\r\n')
        if header_end == -1:
            continue

        header = part[:header_end].decode('utf-8', errors='ignore')

        if 'filename=' in header:
            # 从 Content-Disposition 行提取 filename
            for line in header.split('\r\n'):
                if 'Content-Disposition' in line:
                    for h in line.split(';'):
                        h = h.strip()
                        if h.startswith('filename='):
                            filename = h.split('=', 1)[1].strip('"')
                            break
                    break

            if not filename:
                continue

            # 提取文件内容（去掉末尾的 \r\n 和 --）
            body = part[header_end + 4:]
            # 移除末尾的换行和边界结束标记
            while body.endswith(b'\r') or body.endswith(b'\n') or body.endswith(b'-'):
                body = body[:-1]
            file_data = body
            break

    if not filename or not file_data:
        print(f"[DEBUG] 文件解析失败: filename={filename}, file_data_len={len(file_data) if file_data else 0}")
        return jsonify({'status': 'error', 'message': '没有找到上传的文件'}), 400

    print(f"[DEBUG] 文件解析成功: {filename}, 大小: {len(file_data)} bytes")

    # 手动写入文件（放在可写的临时目录下）
    temp_dir = os.path.join(_tmpdir, 'stylemind_upload')
    os.makedirs(temp_dir, exist_ok=True)
    temp_path = os.path.join(temp_dir, filename)

    with open(temp_path, 'wb') as f:
        f.write(file_data)

    file_size_mb = os.path.getsize(temp_path) / 1024 / 1024

    def generate():
        """SSE 事件生成器"""
        print("[DEBUG] SSE generate() 开始执行")
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        from core.api_client import APIClient
        from core.rag_knowledge import RAGKnowledge
        from core.pdf_splitter import PDFSplitter
        from storage.vector_store import VectorStore
        from storage.database import Database
        from core.pdf_parser import PDFParser

        # 数据库使用绝对路径
        DB_PATH = os.path.join(BASE_DIR, '..', 'stylemind.db')

        def send_event(stage, current, total, message, extra=None):
            """发送 SSE 事件"""
            data = json.dumps({
                'stage': stage,
                'current': current,
                'total': total,
                'percent': int(current / total * 100) if total > 0 else 0,
                'message': message,
                'extra': extra or {}
            }, ensure_ascii=False)
            return f"data: {data}\n\n"

        try:
            # 加载配置
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)

            # 初始化组件（使用对话模型的 API 配置）
            api = None
            chat_api_url = config.get('chat_api_base_url') or config.get('api_base_url')
            chat_api_key = config.get('chat_api_key') or config.get('api_key')
            if chat_api_url and chat_api_key:
                api = APIClient(chat_api_url, chat_api_key)

            # 向量存储也使用绝对路径
            VECTOR_PATH = os.path.join(BASE_DIR, '..', 'vector_store.json')
            vs = VectorStore(VECTOR_PATH)
            db = Database(DB_PATH)
            rag = RAGKnowledge(vs, db, api)
            parser = PDFParser()

            # === 阶段1: 检查文件 ===
            yield send_event('检查文件', 1, 10, f'文件: {filename} ({file_size_mb:.1f}MB)')
            time.sleep(0.1)

            ext = os.path.splitext(temp_path)[1].lower()

            if ext == '.pdf':
                # === 阶段2: 拆分检查 ===
                yield send_event('拆分检查', 2, 10, '检查文件大小和页数...')
                time.sleep(0.1)

                splitter = PDFSplitter(max_size_mb=40, max_pages=50)
                split_files = splitter.split_if_needed(temp_path)

                if len(split_files) > 1:
                    split_info = splitter.get_split_info(temp_path, split_files)
                    yield send_event('拆分完成', 3, 10,
                        f'自动拆分: {file_size_mb:.1f}MB → {len(split_files)} 个文件',
                        {'split_info': split_info})
                else:
                    yield send_event('拆分检查', 3, 10, '文件大小正常，无需拆分')

                time.sleep(0.1)

                # === 阶段3: 解析PDF ===
                all_pages = []
                all_insights = []

                for idx, split_file in enumerate(split_files):
                    yield send_event('解析PDF', 4, 10,
                        f'解析分片 {idx+1}/{len(split_files)}...')

                    pages, metadata = parser.parse(split_file)
                    if pages:
                        all_pages.extend(pages)

                total_pages = len(all_pages)
                yield send_event('解析完成', 5, 10,
                    f'解析完成，共 {total_pages} 页',
                    {'total_pages': total_pages})

                time.sleep(0.1)

                # === 阶段4: LLM分析（逐页） ===
                if api:
                    # 按批次分析（每批5页）
                    batch_size = 5
                    total_batches = (total_pages + batch_size - 1) // batch_size

                    for batch_idx in range(total_batches):
                        start = batch_idx * batch_size
                        end = min(start + batch_size, total_pages)
                        batch_pages = all_pages[start:end]

                        yield send_event('LLM分析', 5 + batch_idx, 5 + total_batches,
                            f'LLM 分析第 {start+1}-{end}/{total_pages} 页...',
                            {'analyzing_pages': list(range(start+1, end+1))})

                        try:
                            insights = rag._analyze_pdf_with_llm(batch_pages,
                                config.get('chat_model', 'chatgpt-4o-latest'))
                            all_insights.extend(insights)

                            # 每页分析成功都发一个事件
                            for i, ins in enumerate(insights):
                                if ins and len(ins) > 10:
                                    yield send_event('LLM分析', 5 + batch_idx, 5 + total_batches,
                                        f'✅ 第 {start+i+1} 页分析成功',
                                        {'page_success': start+i+1})
                        except Exception as e:
                            yield send_event('LLM分析', 5 + batch_idx, 5 + total_batches,
                                f'⚠️ 第 {start+1}-{end} 页分析失败: {str(e)[:50]}',
                                {'page_error': str(e)})
                            all_insights.extend([''] * len(batch_pages))
                else:
                    yield send_event('LLM分析', 6, 10, '⚠️ API 未配置，跳过 LLM 分析')

                # === 阶段5: 存储到知识库 ===
                for i, page in enumerate(all_pages):
                    if (i + 1) % 10 == 0 or i == 0 or i == total_pages - 1:
                        yield send_event('存储', 7, 7 + (total_pages // 10),
                            f'存储第 {i+1}/{total_pages} 页到知识库...')

                    content = f"{page.title}\n{page.content}"
                    layout_data = {}
                    if hasattr(page, 'layout') and page.layout:
                        layout_data = page.layout.to_dict()

                    vs.add(
                        text=content,
                        metadata={
                            "source": os.path.basename(temp_path),
                            "page_index": page.page_index,
                            "layout_type": layout_data.get("layout_type", "content"),
                            "layout_columns": layout_data.get("columns", 1),
                            "layout_alignment": layout_data.get("text_alignment", "left"),
                            "layout_density": layout_data.get("content_density", "normal"),
                            "has_tables": len(page.tables) > 0,
                            "has_images": len(page.images) > 0,
                            "title_hierarchy": layout_data.get("title_hierarchy", []),
                            "file_type": "pdf",
                            "insights": all_insights[i] if i < len(all_insights) else "",
                        }
                    )

                    db.add_knowledge(
                        source_file=os.path.basename(temp_path),
                        page_index=page.page_index,
                        content=content,
                        insights=all_insights[i] if i < len(all_insights) else ""
                    )

                # === 完成 ===
                llm_count = sum(1 for ins in all_insights if ins and len(ins) > 10)
                result_data = {
                    'pages': total_pages,
                    'llm_analyzed': llm_count,
                    'file': filename,
                    'size_mb': round(file_size_mb, 2)
                }
                if len(split_files) > 1:
                    result_data['split_count'] = len(split_files)

                yield send_event('完成', 1, 1,
                    f'✅ 导入完成！共 {total_pages} 页，LLM 分析 {llm_count} 页',
                    {'result': result_data})

            elif ext in ['.pptx', '.ppt']:
                yield send_event('解析PPT', 3, 10, '正在解析 PPT 文件...')
                result = rag.ingest_ppt(temp_path, config.get('chat_model', 'gpt-4o'))

                if result.get('status') == 'success':
                    result_data = {
                        'pages': result.get('pages', 0),
                        'llm_analyzed': result.get('pages', 0),
                        'file': filename,
                        'size_mb': round(file_size_mb, 2)
                    }
                    yield send_event('完成', 1, 1,
                        f'✅ 导入完成！共 {result.get("pages", 0)} 页',
                        {'result': result_data})
                else:
                    yield send_event('错误', 0, 1, f'❌ {result.get("message", "未知错误")}')

            else:
                yield send_event('错误', 0, 1, f'❌ 不支持的文件格式: {ext}')

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield send_event('错误', 0, 1, f'❌ 导入失败: {str(e)[:200]}')

        finally:
            # 清理临时文件
            try:
                os.remove(temp_path)
            except:
                pass

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive',
        }
    )

@app.route('/api/knowledge/status')
def knowledge_status():
    """获取知识库状态"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        from storage.vector_store import VectorStore

        VECTOR_PATH = os.path.join(BASE_DIR, '..', 'vector_store.json')
        vs = VectorStore(VECTOR_PATH)
        all_data = vs.get_all()

        files = {}
        has_insights = 0

        for item in all_data:
            m = item.get('metadata', {})
            source = m.get('source', '未知')

            if source not in files:
                files[source] = {'total': 0, 'has_insights': 0}

            files[source]['total'] += 1

            insights = m.get('insights', '')
            if insights and len(insights) > 10:
                has_insights += 1
                files[source]['has_insights'] += 1

        return jsonify({
            'total_records': len(all_data),
            'total_files': len(files),
            'has_insights': has_insights,
            'files': files
        })

    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/knowledge/stats')
def knowledge_stats():
    """获取知识库统计"""
    return knowledge_status()

@app.route('/api/knowledge/clear', methods=['POST'])
def knowledge_clear():
    """清空知识库"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        from storage.vector_store import VectorStore
        from storage.database import Database

        DB_PATH = os.path.join(BASE_DIR, '..', 'stylemind.db')
        VECTOR_PATH = os.path.join(BASE_DIR, '..', 'vector_store.json')

        vs = VectorStore(VECTOR_PATH)
        db = Database(DB_PATH)

        # 清空向量存储
        vs.clear()

        # 清空数据库
        db.clear_knowledge()

        return jsonify({
            'status': 'success',
            'message': '知识库已清空'
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/knowledge/delete', methods=['POST'])
def knowledge_delete_file():
    """按文件名删除知识库记录"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        from storage.vector_store import VectorStore
        from storage.database import Database

        DB_PATH = os.path.join(BASE_DIR, '..', 'stylemind.db')
        VECTOR_PATH = os.path.join(BASE_DIR, '..', 'vector_store.json')

        data = request.json
        filename = data.get('filename')

        if not filename:
            return jsonify({'status': 'error', 'message': '缺少文件名参数'}), 400

        vs = VectorStore(VECTOR_PATH)
        db = Database(DB_PATH)

        # 删除向量存储中的记录
        vs_deleted = vs.delete_by_source(filename)

        # 删除数据库中的记录
        db_deleted = db.delete_knowledge_by_file(filename)

        return jsonify({
            'status': 'success',
            'message': f'已删除 "{filename}"：{vs_deleted} 条向量记录，{db_deleted} 条数据库记录'
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/chat', methods=['POST'])
def chat():
    """对话接口 - 与AI对话生成大纲"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        from core.api_client import APIClient
        from core.rag_knowledge import RAGKnowledge
        from storage.vector_store import VectorStore
        from storage.database import Database

        data = request.json
        message = data.get('message', '')
        history = data.get('history', [])
        current_outline = data.get('outline')

        # 加载配置
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)

        # 初始化 API（使用对话模型配置）
        api = None
        chat_api_url = config.get('chat_api_base_url') or config.get('api_base_url')
        chat_api_key = config.get('chat_api_key') or config.get('api_key')
        if chat_api_url and chat_api_key:
            api = APIClient(chat_api_url, chat_api_key)

        # === RAG 知识库检索（3轮） ===
        DB_PATH = os.path.join(BASE_DIR, '..', 'stylemind.db')
        VECTOR_PATH = os.path.join(BASE_DIR, '..', 'vector_store.json')
        knowledge_context = ""
        try:
            vs = VectorStore(VECTOR_PATH)
            db = Database(DB_PATH)
            rag = RAGKnowledge(vs, db, api_client=None)

            import re
            all_results = {}

            # 第1轮：用完整消息检索 top 15
            results1 = rag.query(message, k=15)
            for r in results1:
                key = (r.get('metadata', {}).get('source', ''), r.get('metadata', {}).get('page_index', 0))
                if key not in all_results or r.get('score', 0) > all_results[key].get('score', 0):
                    all_results[key] = r

            # 第2轮：提取关键词组合检索
            keywords = re.findall(r'[\u4e00-\u9fff]{2,6}', message)
            if len(keywords) >= 2:
                keyword_query = ' '.join(keywords[:3])
                results2 = rag.query(keyword_query, k=10)
                for r in results2:
                    key = (r.get('metadata', {}).get('source', ''), r.get('metadata', {}).get('page_index', 0))
                    if key not in all_results or r.get('score', 0) > all_results[key].get('score', 0):
                        all_results[key] = r

            # 第3轮：用不同关键词子集检索（扩大覆盖面）
            if len(keywords) >= 4:
                keyword_query2 = ' '.join(keywords[2:5])  # 取不同的关键词子集
                results3 = rag.query(keyword_query2, k=10)
                for r in results3:
                    key = (r.get('metadata', {}).get('source', ''), r.get('metadata', {}).get('page_index', 0))
                    if key not in all_results or r.get('score', 0) > all_results[key].get('score', 0):
                        all_results[key] = r

            # 按相关度排序，取 top 20
            sorted_results = sorted(all_results.values(), key=lambda x: x.get('score', 0), reverse=True)[:20]

            if sorted_results:
                # 收集排版模式统计
                layout_types = {}
                for r in sorted_results:
                    meta = r.get('metadata', {})
                    lt = meta.get('layout_type', 'content')
                    layout_types[lt] = layout_types.get(lt, 0) + 1

                knowledge_context = "\n\n📚 参考知识库排版经验：\n"
                knowledge_context += "─" * 40 + "\n"
                knowledge_context += f"【排版类型分布】\n"
                for lt, count in sorted(layout_types.items(), key=lambda x: -x[1]):
                    knowledge_context += f"  • {lt}: {count}页\n"

                knowledge_context += f"\n【参考页面详情】\n"
                for i, r in enumerate(sorted_results[:10]):  # 只取前10个最相关的
                    text = r.get('text', '')
                    meta = r.get('metadata', {})
                    source = meta.get('source', '未知来源')
                    page = meta.get('page_index', '?')
                    layout_type = meta.get('layout_type', 'content')
                    layout_columns = meta.get('layout_columns', 1)
                    has_images = meta.get('has_images', False)
                    has_tables = meta.get('has_tables', False)
                    insights = meta.get('insights', '')
                    score = r.get('score', 0)

                    content = text[:400] if text else ""
                    knowledge_context += f"\n【参考{i+1}】{source} 第{page+1}页\n"
                    knowledge_context += f"  排版: {layout_type} | 分栏: {layout_columns} | 图:{has_images} | 表:{has_tables}\n"
                    knowledge_context += f"  内容: {content}\n"
                    if insights:
                        knowledge_context += f"  分析: {insights[:200]}\n"

                knowledge_context += "\n" + "─" * 40 + "\n"
                print(f"[INFO] RAG 3轮检索完成，共 {len(sorted_results)} 条，排版类型: {layout_types}")
        except Exception as e:
            print(f"[WARN] RAG 检索失败: {e}")

        # === 判断是否需要生成大纲 ===
        need_outline = '大纲' in message or '结构' in message or 'PPT' in message or 'ppt' in message or not current_outline
        print(f"[DEBUG] need_outline={need_outline}, message={message[:50] if message else 'EMPTY'}..., current_outline={current_outline is not None}, api={api is not None}")

        if not need_outline or not api:
            # 普通对话，不生成大纲
            system_prompt = f"""你是 StyleMind，一个专业的 PPT 设计助手。
请用中文回复，保持专业、友好的语气。
{knowledge_context}"""

            messages = [{'role': 'system', 'content': system_prompt}]
            _hist = history[-20:]
            for h in _hist:
                messages.append({'role': h['role'], 'content': h['content']})
            messages.append({'role': 'user', 'content': message})

            if api:
                response = api.chat(
                    model=config.get('chat_model', 'gpt-4o'),
                    messages=messages,
                    max_tokens=8000,
                    timeout=180
                )
                ai_response = response['content']
            else:
                ai_response = "API 未配置，无法生成回复。请先配置 API 设置。"

            return jsonify({
                'status': 'success',
                'response': ai_response,
                'outline': current_outline
            })

        # === 生成大纲 - V2: 逐章节生成，不续生成 ===
        print(f"[INFO] 使用逐章节生成策略（V2），不续生成...")

        from outline_generator_v2 import generate_outline_by_sections_v2

        brand_name = data.get('brand_name', '')
        brand_colors = data.get('brand_colors', '')
        style_description = data.get('style_description', '')

        skeleton = generate_outline_by_sections_v2(
            api=api,
            config=config,
            brief_text=message,
            brand_name=brand_name,
            brand_colors=brand_colors,
            style_description=style_description
        )

        print(f"[DEBUG] 逐章节生成完成: {len(skeleton.get('pages', []))} 页")

        if not skeleton or not skeleton.get('pages'):
            # 骨架提取失败，直接用完整提示词生成
            print(f"[WARN] 骨架提取失败，回退到单次生成")
            full_prompt = f"""你是 StyleMind，一个专业的 PPT 设计助手。
请根据用户需求和知识库内容生成完整的 PPT 大纲。
{knowledge_context}

【⚠️ 内容边界 — 绝对禁止违反】
- content 字段的所有文字必须来自用户提供的原文或知识库内容
- 禁止自行创作、改写、扩写或添加原文没有的信息
- 禁止编造数据、数字、百分比等
- 如果信息不足，宁可简短也不要编造

请严格按以下 JSON 格式输出：
```json
{{
  "title": "PPT标题",
  "pages": [
    {{"title": "封面页标题", "type": "cover", "content": "来自原文的内容"}},
    {{"title": "目录页标题", "type": "toc", "content": "来自原文的内容"}},
    {{"title": "内容页标题", "type": "content", "content": "来自原文的内容"}},
    {{"title": "总结页标题", "type": "summary", "content": "来自原文的内容"}}
  ]
}}
```
content 字段使用原文中的实际内容，不要编造。"""

            full_response = api.chat(
                model=config.get('chat_model', 'gpt-4o'),
                messages=[{'role': 'user', 'content': full_prompt}],
                max_tokens=16000,
                timeout=300
            )
            outline = extract_outline_from_response(full_response['content'])

            if not outline:
                outline = {'title': 'PPT大纲', 'pages': []}

            return jsonify({
                'status': 'success',
                'response': f"大纲已生成，共 {len(outline.get('pages', []))} 页。",
                'outline': outline
            })

        # 使用逐章节生成的结果，直接返回
        skeleton_pages = skeleton.get('pages', [])
        total_pages = len(skeleton_pages)
        print(f"[INFO] 逐章节生成完成: {total_pages} 页")

        # 组装最终大纲
        final_outline = {
            'title': 'PPT大纲',
            'pages': skeleton_pages
        }

        # 构建对话中显示的大纲文本
        outline_text = f"📋 **PPT大纲**\n\n"
        outline_text += f"共 {total_pages} 页大纲：\n\n"

        for i, page in enumerate(skeleton_pages, 1):
            title = page.get('title', f'第{i}页')
            ptype = page.get('type', 'content')
            layout = page.get('layout', '标准排版')
            brief = page.get('brief', '')
            content = page.get('content', '')[:100]  # 只显示前100字

            outline_text += f"**{i}. {title}** ({ptype})\n"
            if layout:
                outline_text += f"   排版: {layout}\n"
            if brief:
                outline_text += f"   简介: {brief}\n"
            if content:
                outline_text += f"   内容: {content}...\n"
            outline_text += "\n"

        outline_text += "---\n💡 确认大纲后，点击「根据大纲生成图片」开始制作 PPT。"

        return jsonify({
            'status': 'success',
            'response': outline_text,
            'outline': final_outline
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/chat-stream', methods=['POST'])
def chat_stream():
    """对话接口（SSE）- 显示大纲生成进度条"""
    def sse(obj: dict) -> str:
        return f"data: {json.dumps(obj, ensure_ascii=False)}\n\n"

    def generate():
        try:
            import sys
            sys.path.insert(0, os.path.join(BASE_DIR, '..'))
            from core.api_client import APIClient
            from core.rag_knowledge import RAGKnowledge
            from storage.vector_store import VectorStore
            from storage.database import Database

            data = request.get_json(silent=True) or {}
            message = data.get('message', '')
            history = data.get('history', [])
            current_outline = data.get('outline')

            yield sse({"status": "start", "percent": 1, "message": "开始处理...（检索知识库）"})

            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)

            api = None
            chat_api_url = config.get('chat_api_base_url') or config.get('api_base_url')
            chat_api_key = config.get('chat_api_key') or config.get('api_key')
            if chat_api_url and chat_api_key:
                api = APIClient(chat_api_url, chat_api_key)

            # RAG 检索
            DB_PATH = os.path.join(BASE_DIR, '..', 'stylemind.db')
            VECTOR_PATH = os.path.join(BASE_DIR, '..', 'vector_store.json')
            knowledge_context = ""
            try:
                vs = VectorStore(VECTOR_PATH)
                db = Database(DB_PATH)
                rag = RAGKnowledge(vs, db, api_client=None)

                import re
                all_results = {}
                results1 = rag.query(message, k=15)
                for r in results1:
                    key = (r.get('metadata', {}).get('source', ''), r.get('metadata', {}).get('page_index', 0))
                    if key not in all_results or r.get('score', 0) > all_results[key].get('score', 0):
                        all_results[key] = r

                keywords = re.findall(r'[\u4e00-\u9fff]{2,6}', message)
                if len(keywords) >= 2:
                    keyword_query = ' '.join(keywords[:3])
                    results2 = rag.query(keyword_query, k=10)
                    for r in results2:
                        key = (r.get('metadata', {}).get('source', ''), r.get('metadata', {}).get('page_index', 0))
                        if key not in all_results or r.get('score', 0) > all_results[key].get('score', 0):
                            all_results[key] = r
                if len(keywords) >= 4:
                    keyword_query2 = ' '.join(keywords[2:5])
                    results3 = rag.query(keyword_query2, k=10)
                    for r in results3:
                        key = (r.get('metadata', {}).get('source', ''), r.get('metadata', {}).get('page_index', 0))
                        if key not in all_results or r.get('score', 0) > all_results[key].get('score', 0):
                            all_results[key] = r

                sorted_results = sorted(all_results.values(), key=lambda x: x.get('score', 0), reverse=True)[:20]
                if sorted_results:
                    layout_types = {}
                    for r in sorted_results:
                        meta = r.get('metadata', {})
                        lt = meta.get('layout_type', 'content')
                        layout_types[lt] = layout_types.get(lt, 0) + 1
                    knowledge_context = "\n\n📚 参考知识库排版经验：\n"
                    knowledge_context += "─" * 40 + "\n"
                    knowledge_context += f"【排版类型分布】\n"
                    for lt, count in sorted(layout_types.items(), key=lambda x: -x[1]):
                        knowledge_context += f"  • {lt}: {count}页\n"
                    knowledge_context += f"\n【参考页面详情】\n"
                    for i, r in enumerate(sorted_results[:10]):
                        text = r.get('text', '')
                        meta = r.get('metadata', {})
                        source = meta.get('source', '未知来源')
                        page = meta.get('page_index', '?')
                        layout_type = meta.get('layout_type', 'content')
                        layout_columns = meta.get('layout_columns', 1)
                        has_images = meta.get('has_images', False)
                        has_tables = meta.get('has_tables', False)
                        insights = meta.get('insights', '')
                        content = text[:400] if text else ""
                        knowledge_context += f"\n【参考{i+1}】{source} 第{page+1}页\n"
                        knowledge_context += f"  排版: {layout_type} | 分栏: {layout_columns} | 图:{has_images} | 表:{has_tables}\n"
                        knowledge_context += f"  内容: {content}\n"
                        if insights:
                            knowledge_context += f"  分析: {insights[:200]}\n"
                    knowledge_context += "\n" + "─" * 40 + "\n"
            except Exception as e:
                print(f"[WARN] RAG 检索失败: {e}")

            need_outline = '大纲' in message or '结构' in message or 'PPT' in message or 'ppt' in message or not current_outline
            if not need_outline or not api:
                yield sse({"status": "progress", "percent": 60, "message": "生成回复..."})
                system_prompt = f"""你是 StyleMind，一个专业的 PPT 设计助手。
请用中文回复，保持专业、友好的语气。
{knowledge_context}"""
                messages = [{'role': 'system', 'content': system_prompt}]
                _hist = history[-20:]
                for h in _hist:
                    messages.append({'role': h['role'], 'content': h['content']})
                messages.append({'role': 'user', 'content': message})
                resp = api.chat(model=config.get('chat_model', 'gpt-4o'), messages=messages, max_tokens=8000, timeout=180)
                yield sse({"status": "done", "percent": 100, "message": "完成", "response": resp.get("content", ""), "outline": current_outline})
                return

            # ===== 大纲生成 - V2: 逐章节生成，无续生成 =====
            yield sse({"status": "progress", "percent": 15, "message": "开始逐章节生成大纲..."})
            skeleton = generate_outline_stream_v2(api, config, message, data, knowledge_context, lambda obj: sse(obj))

            yield sse({"status": "progress", "percent": 80, "message": "整理大纲内容..."})

            skeleton_pages = (skeleton or {}).get('pages', []) or []
            detailed_pages = []
            for page in skeleton_pages:
                existing_content = page.get('content', '')
                brief = page.get('brief', '')
                if len(existing_content) > 100:
                    detailed_pages.append({'title': page.get('title', ''), 'type': page.get('type', 'content'), 'layout': page.get('layout', ''), 'brief': brief, 'content': existing_content})
                elif len(brief) > 100:
                    detailed_pages.append({'title': page.get('title', ''), 'type': page.get('type', 'content'), 'layout': page.get('layout', ''), 'brief': brief, 'content': brief})
                else:
                    detailed_pages.append(page)

            final_outline = {'title': (skeleton or {}).get('title', 'PPT大纲'), 'pages': detailed_pages}

            outline_text = f"📋 **{final_outline['title']}**\n\n"
            outline_text += f"共 {len(detailed_pages)} 页大纲：\n\n"
            for i, page in enumerate(detailed_pages, 1):
                if not isinstance(page, dict):
                    print(f"[WARN] 页面 {i} 不是字典类型: {type(page)}，跳过")
                    continue
                title = page.get('title', f'第{i}页')
                ptype = page.get('type', 'content')
                layout = page.get('layout', '标准排版')
                brief = page.get('brief', '')
                _content_raw = page.get('content', '')
                content = str(_content_raw)[:100] if _content_raw else ''
                outline_text += f"**{i}. {title}** ({ptype})\n"
                if layout:
                    outline_text += f"   排版: {layout}\n"
                if brief:
                    outline_text += f"   简介: {brief}\n"
                if content:
                    outline_text += f"   内容: {content}...\n"
                outline_text += "\n"
            outline_text += "---\n💡 确认大纲后，点击「根据大纲生成图片」开始制作 PPT。"

            yield sse({"status": "done", "percent": 100, "message": "大纲已生成", "response": outline_text, "outline": final_outline})

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield sse({"status": "error", "percent": 100, "message": str(e)})

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'},
    )

def extract_outline_from_response(text):
    """从 AI 回复中提取大纲（支持 JSON 格式）"""
    import re

    print(f"[DEBUG] extract_outline_from_response 收到文本长度: {len(text) if text else 0}")

    if not text:
        return None

    def _normalize_outline(obj):
        """把 outline/pages 归一化为 {title, pages:[{...}]}，并过滤明显无效项。"""
        import json as _json

        if obj is None:
            return None

        pages = None
        title = "PPT大纲"
        total_pages = None

        if isinstance(obj, dict):
            title = obj.get("title") or title
            pages = obj.get("pages")
            tp = obj.get("total_pages")
            if isinstance(tp, int):
                total_pages = tp
            else:
                # 某些模型会把 total_pages 输出成字符串
                try:
                    total_pages = int(tp) if tp is not None else None
                except Exception:
                    total_pages = None
        elif isinstance(obj, list):
            pages = obj
        else:
            return None

        if not isinstance(pages, list) or len(pages) == 0:
            return None

        cleaned = []
        for p in pages:
            # 允许字符串形式的 JSON（部分模型会把每页 dict 当成 string）
            if isinstance(p, str):
                s = p.strip()
                if s.startswith("{") and s.endswith("}"):
                    try:
                        p = _json.loads(s)
                    except Exception:
                        continue
                else:
                    continue
            if not isinstance(p, dict):
                continue

            t = str(p.get("title") or "").strip()
            # 过滤“示例/占位/注释”类输出
            if not t:
                continue
            if t.startswith("//") or t.startswith("/*"):
                continue
            if '"index"' in t or t.startswith("{") or t.endswith("}"):
                # 很像把整段 JSON 当标题了
                continue

            cleaned.append(p)

        if not cleaned:
            return None

        # 按 index 排序（若存在）
        try:
            cleaned.sort(key=lambda x: int(x.get("index") or 10**9))
        except Exception:
            pass

        out = {"title": title, "pages": cleaned}
        if isinstance(total_pages, int) and total_pages > 0:
            out["total_pages"] = total_pages
        return out

    # 优先尝试提取 ```json ... ``` 块（可能是对象 {} 或数组 []）
    json_match = re.search(r'```json\s*([\s\S]*?)\s*```', text)
    if json_match:
        json_str = json_match.group(1).strip()
        try:
            outline = json.loads(json_str)
            norm = _normalize_outline(outline)
            if norm and norm.get("pages"):
                print(f"[DEBUG] 从 ```json 块提取到 {len(norm['pages'])} 页")
                return norm
        except json.JSONDecodeError as e:
            print(f"[DEBUG] ```json 解析失败: {e}")

    # 尝试提取 ``` ... ``` 块（无 json 标记）
    json_match = re.search(r'```\s*([\s\S]*?)\s*```', text)
    if json_match:
        try:
            outline = json.loads(json_match.group(1).strip())
            norm = _normalize_outline(outline)
            if norm and norm.get("pages"):
                print(f"[DEBUG] 从 ``` 块提取到 {len(norm['pages'])} 页")
                return norm
        except json.JSONDecodeError as e:
            print(f"[DEBUG] ``` 解析失败: {e}")

    # 尝试提取整个 JSON 对象（更宽松的匹配）
    # 找到第一个 { 和最后一个 }
    start = text.find('{')
    end = text.rfind('}')
    if start != -1 and end != -1 and end > start:
        json_str = text[start:end+1]
        try:
            outline = json.loads(json_str)
            norm = _normalize_outline(outline)
            if norm and norm.get("pages"):
                print(f"[DEBUG] 从文本中提取到 {len(norm['pages'])} 页")
                return norm
        except json.JSONDecodeError as e:
            print(f"[DEBUG] JSON 解析失败: {e}")

    # 尝试提取 JSON 数组 [...]
    start = text.find('[')
    end = text.rfind(']')
    if start != -1 and end != -1 and end > start:
        json_str = text[start:end+1]
        try:
            pages = json.loads(json_str)
            norm = _normalize_outline(pages)
            if norm and norm.get("pages"):
                print(f"[DEBUG] 从数组提取到 {len(norm['pages'])} 页")
                return norm
        except json.JSONDecodeError as e:
            print(f"[DEBUG] 数组解析失败: {e}")

    # 回退：按行提取
    lines = text.split('\n')
    pages = []

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # 匹配 "第X页" 或 "X." 或 "X、" 开头的行
        if any(keyword in line for keyword in ['第', '页', '封面', '目录', '总结']):
            pages.append({
                'title': line.replace('第', '').replace('页', '').strip(),
                'type': 'content',
                'content': line
            })

    if pages:
        print(f"[DEBUG] 按行提取到 {len(pages)} 页")
        return {'title': 'PPT大纲', 'pages': pages}

    print(f"[DEBUG] 未能提取任何页面")
    return None

@app.route('/api/generate-images', methods=['POST'])
def generate_images():
    """根据大纲生成图片（SSE 流式进度）"""
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    data = request.get_json(silent=True) or {}
    outline = data.get('outline', {})
    style_settings = data.get('style', {})
    start_from = data.get('start_from', 0)
    existing_images = data.get('existing_images', [])

    # 解析风格设置
    style_desc = style_settings.get('description', '')
    color_scheme = style_settings.get('colorScheme', 'auto')
    reference_image = style_settings.get('referenceImage', '')
    product_image = style_settings.get('productImage', '')
    brand_name = style_settings.get('brandName', '')

    def generate():
        try:
            # 加载配置
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)

            # 使用图片生成模型的 API 配置
            image_api_url = config.get('image_api_base_url') or config.get('api_base_url')
            image_api_key = config.get('image_api_key') or config.get('api_key')

            if not image_api_url or not image_api_key:
                yield f"data: {json.dumps({'status': 'error', 'message': '请先配置图片生成模型的 API 设置'}, ensure_ascii=False)}\n\n"
                return

            pages = outline.get('pages', [])
            if not pages:
                yield f"data: {json.dumps({'status': 'error', 'message': '大纲为空，请先生成大纲'}, ensure_ascii=False)}\n\n"
                return

            api = APIClient(image_api_url, image_api_key)
            image_model = config.get('image_model', 'gpt-image-2')
            chat_model = config.get('chat_model', 'gpt-4o')
            total = len(pages)
            images = list(existing_images) if existing_images else []

            if start_from > 0:
                yield f"data: {json.dumps({'status': 'progress', 'current': start_from, 'total': total, 'message': f'⏭ 断点续传：从第 {start_from+1} 页继续生成（已跳过前 {start_from} 张已有图片）...'}, ensure_ascii=False)}\n\n"
                print(f"[INFO] 断点续传: 从第{start_from+1}页开始，已有{len(images)}张图片")

            # 获取品牌名称和色系
            color_scheme_val = style_settings.get('colorScheme', '')

            # 初始化 color_prompt，确保始终有值
            color_prompt = color_scheme_val if color_scheme_val else 'appropriate color scheme matching the content'

            # 如果有品牌名称但没有色系，自动识别品牌色系
            brand_colors = ""
            if brand_name and not color_scheme_val:
                try:
                    yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': f'正在识别 {brand_name} 的品牌色系...'}, ensure_ascii=False)}\n\n"

                    color_prompt = f"""Identify the brand colors for "{brand_name}".
Return ONLY a JSON object with:
- primary_color: main brand color (hex code and name)
- secondary_color: accent color (hex code and name)
- background_suggestion: recommended background style
- text_color: recommended text color

Example for "支付宝":
{{"primary_color": "#1677FF (支付宝蓝)", "secondary_color": "#FF6A00 (支付宝橙)", "background_suggestion": "light or white background", "text_color": "#333333 or white on dark backgrounds"}}"""

                    color_result = api.chat(
                        model=chat_model,
                        messages=[{'role': 'user', 'content': color_prompt}],
                        max_tokens=200,
                        timeout=30
                    )

                    brand_colors = color_result.get('content', '')
                    print(f"[INFO] 品牌 {brand_name} 色系: {brand_colors}")

                except Exception as e:
                    print(f"[WARN] 品牌色系识别失败: {e}")
                    brand_colors = ""

            # 如果用户指定了色系，使用用户的
            if color_scheme_val:
                brand_colors = f"User specified color scheme: {color_scheme_val}"

            # 分析参考图风格（如果有）
            reference_style = ""
            if reference_image and reference_image.startswith('data:image'):
                try:
                    yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': '正在分析参考图风格...'}, ensure_ascii=False)}\n\n"

                    # 使用多模态模型分析参考图
                    analysis_prompt = "Analyze this reference image and describe its visual style in detail. Include: color palette (list specific hex colors), typography style, layout characteristics, visual elements, overall mood/atmosphere. Keep it under 150 words."

                    # 构建多模态消息
                    messages = [
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": analysis_prompt},
                                {"type": "image_url", "image_url": {"url": reference_image}}
                            ]
                        }
                    ]

                    analysis_result = api.chat(
                        model=chat_model,
                        messages=messages,
                        max_tokens=300,
                        timeout=60
                    )

                    reference_style = analysis_result.get('content', '')
                    print(f"[INFO] 参考图风格分析: {reference_style[:100]}...")

                except Exception as e:
                    print(f"[WARN] 参考图分析失败: {e}")
                    reference_style = ""

            # === 关键：先用 AI 确定统一的风格规范 ===
            # 这样所有页面都会使用同一套配色和风格，保证 PPT 整体一致性
            unified_style_guide = ""

            # 检测是否为无印良品/MUJI风格，如果是则使用预设的精确规范
            is_muji_style = style_desc and ('无印良品' in style_desc or 'MUJI' in style_desc.upper() or 'muji' in style_desc.lower() or '日式极简' in style_desc or '日式排版' in style_desc)

            if is_muji_style:
                # 使用精确的无印良品风格规范
                unified_style_guide = """MUJI (无印良品) Japanese Minimalist Style - STRICT RULES:
- Background: Warm off-white #F5F5F0 (like washi paper), NEVER dark, NEVER gradient
- Primary text: Soft charcoal #4A4A4A (NOT pure black), left-aligned ONLY
- Accent color: Natural beige #D4C4B0 (wood/linen feel), use sparingly
- Layout: 60% white space (MA/間), asymmetric balance, grid-based
- Typography: Clean sans-serif, generous line spacing, hierarchical sizing
- Visuals: Natural textures only, single focal point per slide, NO decorative borders/gradients/shadows
- Philosophy: "This is enough" (これでいい) - understated, honest, essential
- ABSOLUTELY NO: Bright colors, gradients, centered text (except cover), flashy effects, pure black backgrounds"""
                yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': '已应用无印良品精确风格规范...'}, ensure_ascii=False)}\n\n"
                print("[INFO] 应用无印良品精确风格规范")
            elif style_desc or color_prompt or reference_style:
                try:
                    yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': '正在确定统一风格规范...'}, ensure_ascii=False)}\n\n"

                    style_guide_prompt = f"""You are a professional brand designer. Based on the following style requirements, define a UNIFIED visual style guide for a PPT presentation. All slides must follow this EXACT same style to maintain consistency, BUT each slide should have varied backgrounds and moods.

Style description: {style_desc or 'Not specified'}
Color preference: {color_prompt or 'Not specified'}
Reference image style: {reference_style or 'No reference'}

Please output a concise style guide in this format (under 250 words):
- Primary colors: (list 2-3 specific colors with hex codes)
- Secondary colors: (list 1-2 accent colors)
- Background VARIETY: (list 5+ different background styles to rotate across slides: e.g. "daylight urban", "soft indoor", "blue hour cityscape", "abstract geometric", "warm lifestyle", "cool minimal", etc.)
- Typography mood: (modern/classic/playful/minimal)
- Visual elements: (shapes, lines, textures to use consistently)
- Overall atmosphere: (one sentence)
- ANTI-PATTERN: explicitly list what to AVOID (e.g. "NOT every slide golden hour sunset", "NOT every slide with person's back silhouette")

IMPORTANT: Be specific about colors. The background variety is CRITICAL - each slide must look visually distinct while maintaining color and typography consistency."""

                    guide_result = api.chat(
                        model=chat_model,
                        messages=[{'role': 'user', 'content': style_guide_prompt}],
                        max_tokens=400,
                        timeout=60
                    )

                    unified_style_guide = guide_result.get('content', '')
                    print(f"[INFO] 统一风格规范: {unified_style_guide[:200]}...")

                except Exception as e:
                    print(f"[WARN] 风格规范生成失败: {e}")
                    unified_style_guide = ""

            print(f"[IMG] ========== 开始生图循环: 共{total}页, 模型={image_model}, API={image_api_url} ==========", flush=True)
            for i, page in enumerate(pages):
                if i < start_from:
                    continue
                page_title = page.get('title', f'第{i+1}页')
                page_content = page.get('content', '')
                page_type = page.get('type', 'content')
                page_layout = page.get('layout', '')
                final_prompt = None
                if '模块化-三列卡片' in page_layout:
                    page_layout = page_layout.replace('模块化-三列卡片', '模块化-卡片')
                if '模块化-三列卡片' in page_layout:
                    page_layout = page_layout.replace('模块化-三列卡片', '模块化-卡片')

                # 进度
                yield f"data: {json.dumps({'status': 'progress', 'current': i+1, 'total': total, 'message': f'正在生成第 {i+1}/{total} 页: {page_title}'}, ensure_ascii=False)}\n\n"

                # 根据大纲中的 layout 字段判断布局变体
                # 两种变体共享同一套设计语言，只是图片占比不同
                is_fullbleed_page = '满版图片' in page_layout or \
                    page_type in ('cover', 'visual')

                # 根据大纲中的 layout 字段精确匹配PDF风格
                if '满版图片-全屏背景' in page_layout or '满版图片-电影感' in page_layout:
                    # 满版图片-全屏背景：全屏背景+白字叠加
                    # 根据页面内容动态调整背景风格，避免同质化
                    content_lower = page_content.lower() if page_content else ''

                    # 根据内容关键词选择不同的视觉风格
                    if any(kw in page_content for kw in ['数据', '分析', '趋势', '增长', '用户', '比例', '调研']):
                        bg_style = """Background: abstract data visualization style, geometric patterns, subtle gradient, modern tech feel."""
                    elif any(kw in page_content for kw in ['TVC', '脚本', '画面', '镜头', '旁白']):
                        bg_style = """Background: video production scene, camera equipment, filming location, or storyboard sketches."""
                    elif any(kw in page_content for kw in ['愿望', '心愿', '梦想', '希望', '坚持', '努力']):
                        bg_style = """Background: DIVERSE lifestyle scenes - CAN USE: cafe window with rain, street food stall, bus window, park bench, rooftop at sunset, night market, bookstore aisle, cozy bed with string lights. AVOID: desk lamp, person at desk, indoor study room."""
                    elif any(kw in page_content for kw in ['品牌', '站位', '定位', '策略', '传播']):
                        bg_style = """Background: clean modern workspace or urban architecture, professional atmosphere, bright daylight or blue hour. AVOID: desk lamp scenes."""
                    elif any(kw in page_content for kw in ['洞察', '发现', '现象', '观察']):
                        bg_style = """Background: street photography, documentary style, authentic urban moments, natural daylight. AVOID: desk lamp, indoor cozy scenes."""
                    else:
                        bg_style = """Background: HIGHLY VARIED scenes - use different settings for each slide: outdoor cafe, public park, metro/train, night market, rooftop, bookstore, street corner, beach, mountain trail, market stall. AVOID: desk lamp, person at desk."""

                    # 场景多样性轮换
                    diverse_scenes = [
                        "street food stall with warm lights",
                        "person reading on a park bench",
                        "cafe window view with rain",
                        "night market with colorful lights",
                        "rooftop with city skyline at dusk",
                        "bookstore aisle with soft lighting",
                        "metro window with passing scenery",
                        "coastal walkway at sunset",
                        "flower market early morning",
                        "outdoor market stall",
                    ]

                    # 产品图约束
                    product_constraint = ""
                    if product_image:
                        product_constraint = """

【重要 - 产品图严格约束】
用户提供了实际的产品参考图。如果当前页面内容涉及以下场景，必须在生成的图片中严格保持产品的真实外观和外形特征，不得重新绘制或想象一个不同的产品：
- 产品展示、产品特写、产品使用场景
- 包含产品名称、品牌名称的页面
- 任何需要出现该实物产品的画面

处理方式：
- 如果是产品展示页：直接使用产品原图作为主要视觉元素，叠加文字信息
- 如果是使用场景页：将产品图自然融入场景中，保持产品外形100%一致
- 不要重新设计产品的外形、颜色、材质
"""

                    style_text = f"""PRECISE LAYOUT MATCH - 满版图片-全屏背景 (Full-bleed background):
- FULL SCREEN high-quality background image
- {bg_style}
- ANTI-PATTERN: STRICTLY FORBIDDEN - DO NOT use these common repeating patterns:
  * Person sitting at desk with desk lamp (this pattern is OVERUSED)
  * Person writing in notebook at study desk
  * Indoor cozy study room with warm lamp light
  * Student at desk with books and lamp
  * Heavy cinematic color grading (orange/teal tones)
  * Dramatic movie-style lighting
- Background scene MUST match the page's actual content topic — see SCENE DESCRIPTION below for details
- White text directly overlaid with subtle dark gradient overlay for readability
- Top-left corner: small brand-colored rectangular tag with white text for section label
- Large expressive title (60-80px) in white, bold
- Clean body text (16-20px) in white or light gray
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- NO modular cards, NO split layout - pure full-bleed photography with text overlay
- IMPORTANT: Each slide MUST have a UNIQUE background scene that matches its specific content
- IMPORTANT: Use NATURAL lighting and colors, avoid heavy color grading or cinematic effects"""

                elif '满版图片-左图右文' in page_layout:
                    # 满版图片-左图右文：左侧图+右侧白底文字
                    # 根据内容判断左侧应该放什么
                    has_app = any(kw in page_content for kw in ['APP', 'app', '界面', '截图', '手机', '功能', '玩法'])
                    has_social = any(kw in page_content for kw in ['抖音', '小红书', '微博', '社交', '话题', '晒图', 'KOL', 'KOC'])
                    has_people = any(kw in page_content for kw in ['人物', '用户', '年轻人', '场景'])

                    if has_app:
                        left_desc = "realistic smartphone mockup showing app interface with detailed UI elements"
                    elif has_social:
                        left_desc = "social media feed screenshot or phone showing social platform interface"
                    elif has_people:
                        left_desc = "authentic lifestyle photography of young Chinese people in daily scenes"
                    else:
                        left_desc = "high-quality thematic photograph or illustration matching the page topic"

                    style_text = f"""PRECISE LAYOUT MATCH - 满版图片-左图右文 (Split layout):
- LEFT side (40-50%): {left_desc}
- RIGHT side (50-60%): Clean white background text area with structured content layout
- Top-left: brand-colored rectangular section tag
- Right area: Large brand-colored title + well-organized body text
- Body text should use bullet points, numbered lists, or small info cards for visual interest
- Add subtle design elements: thin colored dividers, small icons, or accent lines
- Clear visual separation between image and text areas
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional, polished, agency-quality design"""

                elif '模块化-卡片' in page_layout:
                    card_count = _guess_card_count(page_content)
                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-卡片 (Card module with dynamic columns):
- {card_count} equal-width cards arranged HORIZONTALLY in a row
- Number of cards should match content: if content has 4 items, use 4 cards; if 5 items, use 5 cards
- Each card: rounded corners (12px radius), subtle shadow, white/light background
- Card header: icon (line style) + title in brand primary color
- Card body: brief description in dark gray
- Cards separated by consistent spacing (24px gap)
- Top section: page title with brand color accent
- Background: light warm gray or subtle texture
- Visual hierarchy: cards are the focal point
- Each card self-contained with clear visual boundaries
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

                elif '模块化-步骤流程' in page_layout:
                    # 模块化-步骤流程：STEP步骤
                    has_app_content = any(kw in page_content for kw in ['APP', 'app', '界面', '截图', '手机', '功能', '玩法'])
                    step_visual = "Include realistic smartphone mockups showing app screens at each step" if has_app_content else "Use icons and illustrations for each step, NO phone mockups unless content mentions app"

                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-步骤流程 (Step-by-step flow):
- Vertical or horizontal flow with STEP 1 / STEP 2 / STEP 3 markers
- Each step: numbered badge (brand primary color circle with white number) + title + description
- {step_visual}
- Connect steps with subtle arrows or lines
- Each step module has light background with subtle shadow
- Mix of text and visual elements (icons, illustrations, or phone screens where relevant)
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional social media campaign presentation style"""

                elif '模块化-时间轴' in page_layout:
                    # 模块化-时间轴：三阶段时间轴
                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-时间轴 (Timeline):
- Horizontal timeline with three nodes: 预热期 → 爆发期 → 长尾期
- Each node: circular marker (brand primary color fill) + phase name + key actions
- Timeline line connecting all nodes (brand color or gray)
- Below timeline: detailed modules for each phase with images and text
- Show campaign progression visually
- Include small preview images or icons for each phase
- Clean, organized, easy to follow at a glance
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional campaign timeline design"""

                elif '模块化-表格' in page_layout:
                    # 模块化-表格：传播总览表格
                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-表格 (Data table):
- Multi-column table with clear headers (时间/阶段/目标/主题/动作)
- Header row: brand primary color background with white text
- Data rows: alternating light backgrounds for readability
- Clean grid lines, professional alignment
- Include small icons or color coding for different phases
- Readable font sizes, clear hierarchy
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Campaign overview table style"""

                elif '纯视觉页' in page_layout:
                    # 纯视觉页：KV展示
                    style_text = f"""PRECISE LAYOUT MATCH - 纯视觉页 (Visual only):
- NO text or minimal text only
- Full-screen key visual (KV) display
- Clean, impactful imagery
- Could be single KV or multiple KV layouts shown together
- Focus on visual impact and brand aesthetics
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional advertising campaign visual presentation"""

                else:
                    # 默认使用满版图片-全屏背景
                    style_text = f"""DEFAULT LAYOUT - 满版图片-全屏背景:
- Full-bleed photography background with natural lighting
- White text overlaid with dark gradient for readability
- Brand-colored tag in top-left corner
- Clean, professional presentation style
- High-quality lifestyle photography
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- AVOID: cinematic color grading, dramatic lighting, movie-style effects"""

                # 根据布局类型决定 prompt 策略
                is_text_only = _is_text_only_layout(page_layout)

                if is_text_only:
                    _bc = brand_colors if brand_colors else 'elegant deep navy & warm gold accent'
                    _bg = f"background: soft gradient from near-white (#FAFAFA) to very light brand-tinted tone, with subtle geometric pattern (thin diagonal lines or dots at 5% opacity)"
                    text_style_map = {
                        '文字-金句引言': f"""TEXT-ONLY LAYOUT - 金句引言 (Premium Quote Slide):
- BACKGROUND: {_bg}
- QUOTE TEXT:超大字号(72-96pt), 品牌主色, 加粗, 居中或偏左黄金分割位置, 行距1.2-1.4
- ATTRIBUTION: 小字(14-16pt), 灰色(#888), 斜体, 位于引言下方右侧对齐
- DECORATIVE: 左侧或顶部有品牌色粗线条装饰(4-8px宽, 占高度30-40%), 或大型半透明品牌色引号装饰(200pt+, 10% opacity)
- TITLE AREA: 页面顶部有小标签(badge), 圆角矩形, 品牌色底+白字, 写页面主题
- SPACING: 大量留白(40%+页面空白), 营造呼吸感和高级感
- FONT: 标题用粗体无衬线(Sans-serif Bold), 正文用常规无衬线
- NO photographs, NO illustrations, NO icons as main elements — pure typography art
- BRAND COLORS: {_bc}
- VIBE: 杂志封面级排版, Apple Keynote / Notion 风格的极简高级感""",
                        '文字-数字大屏': f"""TEXT-ONLY LAYOUT - 数字大屏 (Data Dashboard - Premium):
- BACKGROUND: {_bg} + 微妙的网格线(grid lines at 3% opacity)增加数据感
- NUMBERS: 超大数字(100-140pt), 品牌主色, 加粗(Bold/Black), 每个数字下方有小标签说明(14pt)
- LAYOUT: 数字横向排列(2-4个)或卡片式网格排列, 每个数字在独立的浅色圆角卡片内(card: white bg, subtle shadow, 12px radius)
- ACCENT: 数字旁边有小型趋势箭头↑↓或百分比变化指示器(品牌色)
- DIVIDER: 卡片之间有细线分隔, 底部可能有简短的数据洞察文字段落
- TITLE: 页面顶部品牌色标题(36-48pt bold), 可能有副标题(18pt gray)
- DECORATIVE: 极简的柱状图/折线图线条装饰(pure CSS/SVG style, not photo)
- NO photographs, NO lifestyle imagery
- BRAND COLORS: {_bc}
- VIBE: Bloomberg Terminal / 企业数据大屏风格, 信息密度高但不拥挤""",
                        '文字-问答式': f"""TEXT-ONLY LAYOUT - 问答式 (Q&A - Premium Editorial):
- BACKGROUND: {_bg}, 左侧可有品牌色竖条装饰(8px宽, 全高)
- Q (问题): 品牌主色, 大号(28-36pt), 加粗(Bold), 前面有装饰性「Q」符号或问号图标(圆形背景, 品牌色)
- A (回答): 深灰(#333), 中号(18-20pt), 常规体, 行距1.6-1.8, 左对齐, 有适当缩进
- LAYOUT: 每组Q&A之间有细线分隔(dashed line, 2px, light gray), 或用浅色交替行背景(zebra striping)
- NUMBERING: 多个Q时左侧有序号圆圈(品牌色底白字, 24px直径)
- TITLE AREA: 顶部品牌色大标题 + 可能的副标题描述
- DECORATIVE: 右下角或边缘有淡色几何图形装饰(圆形/三角形, 5-10% opacity)
- NO photographs unless absolutely necessary
- BRAND COLORS: {_bc}
- VIBE: FAQ页面 / 知乎精选回答风格, 信息清晰且视觉舒适""",
                        '文字-双栏对比': f"""TEXT-ONLY LAYOUT - 双栏对比 (Comparison - Premium Split Design):
- BACKGROUND: {_bg}
- LAYOUT: 50/50左右分栏, 中间有VS分隔区域(圆形VS徽章, 品牌色底白字, 或细竖线+箭头)
- LEFT COLUMN HEADER: 品牌色A, 浅色背景卡片(white or tinted), 圆角, 内边距充足
- RIGHT COLUMN HEADER: 品牌色B(可用对比色), 同样卡片样式
- CONTENT: 每栏内用要点列表(带小圆点或check标记), 字号16-18pt, 行距1.5
- HIGHLIGHT: 关键差异词用品牌色加粗标注
- TITLE: 顶部居中大标题(36-48pt), 品牌色, 可能有副标题说明对比维度
- DECORATIVE: 分隔区域有微妙的渐变或图案填充
- NO photographs
- BRAND COLORS: {_bc}
- VIBE: 产品功能对比 / 优劣势分析表风格, 清晰且有设计感""",
                        '文字-要点列表': f"""TEXT-ONLY LAYOUT - 要点列表 (Key Points - Premium Structured):
- BACKGROUND: {_bg}, 可有极淡的品牌色渐变方向光效(top-left glow at 3% opacity)
- TITLE: 页面顶部, 品牌主色, 36-48pt bold, 下方有品牌色下划线装饰(underline bar, 60-80px wide, 4px thick)
- LIST ITEMS: 每个要点独立成行, 左侧有大号序号圆圈(品牌色底白字, 32-40px diameter, Bold数字)
- ITEM TEXT: 要点标题(20-24pt, Bold, #222) + 描述文字(16-18pt, regular, #555), 行距1.5
- SPACING: 要点之间间距20-30px, 整体左对齐但有节奏感
- HOVER EFFECT STYLE: 每个要点行有浅色悬停感(极淡背景色变化, like #F8F8F8)
- DECORATIVE: 右侧或底部有淡色几何点缀(斜线纹理, 圆形元素)
- ALTERNATE: 可选奇偶行微弱背景色差异(zebra: white / #FAFAFA)
- NO photographs, NO lifestyle imagery
- BRAND COLORS: {_bc}
- VIBE: Apple Feature List / Notion Database 风格, 结构清晰且现代""",
                        '文字-引用来源': f"""TEXT-ONLY LAYOUT - 引用来源 (References - Premium Academic):
- BACKGROUND: {_bg}
- TITLE: 「参考资料」or「引用来源」, 品牌色, 28-32pt bold, 左上角
- LIST: 编号列表格式 [1] [2] [3], 每条引用完整格式化
- EACH ENTRY: 作者/来源名(Bold, 18pt, 品牌色) + 内容摘要(regular, 15-16pt, #444) + 年份/链接(gray, 13pt)
- LAYOUT: 双栏或单栏(根据内容量自动适配), 充足行距(1.6-1.8)
- SEPARATOR: 每条引用间有细点线分隔(dotted line, 1px, #DDD)
- DECORATIVE: 底部有品牌色细线作为结束标记, 左下角可能有页码
- CONTAINER: 所有内容在一个浅色圆角容器内(white bg, subtle shadow, 16px radius, padding 40px)
- NO photographs
- BRAND COLORS: {_bc}
- VIBE: 学术论文参考文献页 / 法律文件附录风格, 专业严谨但美观"""
                    }
                    style_text = text_style_map.get(page_layout, f"""TEXT-ONLY LAYOUT (Premium Text Design):
- BACKGROUND: soft gradient from near-white to light brand-tinted tone, with subtle geometric pattern
- TITLE in brand primary color (36-48pt bold), body text in dark gray (#333, 16-18pt)
- Strong typography hierarchy: clear size and weight contrast between heading/subheading/body
- Generous whitespace (35-45% of slide area should be empty)
- Decorative elements: brand-colored accent bars, thin divider lines, subtle geometric shapes
- Content organized in clean cards or structured sections with rounded corners
- NO photographs, but decorative graphic elements (lines, shapes, color blocks) are encouraged
- BRAND COLORS: {_bc}
- VIBE: Modern editorial / magazine layout quality — NOT a Word document""")

                    final_prompt = f"""Create a professional PPT slide.

Title: '{page_title}'
Content: '{page_content}'
Layout Type: '{page_layout}'

{style_text}

CRITICAL REQUIREMENTS - TEXT ONLY SLIDE:
- 16:9 aspect ratio
- This is a TEXT-ONLY slide — do NOT add any photos, illustrations, people, scenes, or decorative images
- Focus entirely on typography, layout, spacing, and text hierarchy
- Clean, minimal, professional text-focused design{product_constraint}"""

                else:
                    bc_text = brand_colors if brand_colors else 'Use appropriate brand colors for accents and highlights'

                    _is_strong_image = any(k in page_layout for k in ['满版图片', '纯视觉页'])
                    _is_module = any(k in page_layout for k in ['模块化-卡片', '模块化-步骤', '模块化-时间轴', '模块化-表格'])

                    _scene_desc = _build_scene_description(page_title, page_content)

                    if _is_strong_image:
                        visual_elements = (
                            "\nKEY VISUAL ELEMENTS:\n"
                            "1. SCENE (content-driven): " + _scene_desc + "\n"
                            "2. Only include smartphone mockups IF the content mentions: app, mobile interface, or mobile features\n"
                            "3. Only include social media screenshots IF the content mentions: social sharing, posting, viral content\n"
                            "4. Brand colors: " + bc_text + "\n"
                            "5. Typography: Clean sans-serif Chinese fonts, large expressive titles\n"
                            "6. Design details: Professional polish, agency-quality finish\n"
                        )
                        crITICAL = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- High quality, professional design
- Use authentic photography or strong imagery as the main visual focal point
- Clean, modern, agency-quality presentation design{product_constraint}"""

                    elif _is_module:
                        visual_elements = (
                            "\nKEY VISUAL ELEMENTS:\n"
                            "1. Visual style: Clean graphic design with icons, illustrations, or abstract shapes — NOT lifestyle photography of people\n"
                            "2. Card/module design: Each section should be self-contained with icon + title + text structure\n"
                            "3. Color coding: Use brand colors to differentiate sections or highlight key information\n"
                            "4. Decorative elements: Subtle geometric patterns, gradients, or line art ONLY if it enhances readability\n"
                            "5. DO NOT add photos of people, lifestyle scenes, or stock photography — they clutter module layouts\n"
                            "6. Brand colors: " + bc_text + "\n"
                            "7. Typography: Clear hierarchy, readable body text, bold headers\n"
                            "8. Design details: Rounded corners, subtle shadows, clean lines, generous whitespace\n"
                        )
                        crITICAL = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- This is a MODULE/CARD layout — focus on structured information display, NOT on photography
- NO lifestyle photos, NO people scenes, NO stock photography
- Let the card/module structure and icons carry the visual weight
- Clean, organized, professional data/presentation design{product_constraint}"""

                    else:
                        visual_elements = (
                            "\nKEY VISUAL ELEMENTS:\n"
                            "1. Visual style: Balanced mix of graphics and text based on content needs\n"
                            "2. Only add photography IF the content specifically calls for it (products, real scenes, specific references)\n"
                            "3. Otherwise use: icons, illustrations, charts, diagrams, or clean typography\n"
                            "4. Brand colors: " + bc_text + "\n"
                            "5. Typography: Clean sans-serif Chinese fonts, clear hierarchy\n"
                            "6. Design details: Professional polish, appropriate whitespace\n"
                        )
                        crITICAL = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- High quality, professional design
- Match the exact layout style specified above
- Do not force photography into layouts that don't need it
- Clean, modern presentation design{product_constraint}"""

                    final_prompt = f"""Create a professional PPT slide.

Title: '{page_title}'
Content: '{page_content}'
Layout Type: '{page_layout}'

{style_text}

{visual_elements}

{crITICAL}"""

                if final_prompt is None:
                    _ve2 = ""
                    _cr2 = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- High quality, professional design
- Clean, modern presentation design"""
                    if any(k in page_layout for k in ['满版图片', '纯视觉页']):
                        _ve2 = "\nKEY VISUAL ELEMENTS:\n1. Photography as main visual element\n2. Brand colors: " + (brand_colors or 'appropriate accent colors') + "\n"
                    elif any(k in page_layout for k in ['模块化-卡片', '模块化-步骤', '模块化-时间轴', '模块化-表格']):
                        _ve2 = "\nKEY VISUAL ELEMENTS:\n1. Icons, illustrations, geometric shapes — NO lifestyle photos\n2. Brand colors: " + (brand_colors or 'appropriate accent colors') + "\n"
                    final_prompt = f"""Create a professional PPT slide.

Title: '{page_title}'
Content: '{page_content}'
Layout Type: '{page_layout}'

{style_text}

{_ve2}
{_cr2}{product_constraint if product_image else ''}"""

                try:
                    _needs_prod = _page_needs_product_image(page_title, page_content, page_type, page_layout)

                    if _needs_prod and product_image and product_image.startswith('data:image'):
                        print(f"[INFO] 第{i+1}页「{page_title}」需要产品图 → 使用 /v1/images/edits")
                        messages = [
                            {"role": "user", "content": [
                                {"type": "text", "text": final_prompt},
                                {"type": "image_url", "image_url": {"url": product_image}}
                            ]}
                        ]
                        result = api.image_generate(
                            model=image_model,
                            prompt=final_prompt,
                            size='1792x1024',
                            n=1,
                            messages=messages
                        )
                    elif reference_image and reference_image.startswith('data:image'):
                        messages = [
                            {"role": "user", "content": [
                                {"type": "text", "text": final_prompt},
                                {"type": "image_url", "image_url": {"url": reference_image}}
                            ]}
                        ]
                        result = api.image_generate(
                            model=image_model,
                            prompt=final_prompt,
                            size='1792x1024',
                            n=1,
                            messages=messages
                        )
                    else:
                        if product_image and not _needs_prod:
                            print(f"[INFO] 第{i+1}页「{page_title}」不需要产品图 → 标准生成")
                        result = api.image_generate(
                            model=image_model,
                            prompt=final_prompt,
                            size='1792x1024',
                            n=1
                        )

                    img_url = None
                    for img in result.get('images', []):
                        if img.get('url'):
                            img_url = img['url']
                        elif img.get('b64_json'):
                            # 将 base64 图片落盘，返回可复用 URL（避免刷新后丢失 & localStorage 超限）
                            try:
                                import base64
                                raw = base64.b64decode(img['b64_json'])
                                fname = f"{int(time.time())}_{i+1}_{uuid.uuid4().hex[:10]}.png"
                                fpath = os.path.join(GENERATED_DIR, fname)
                                with open(fpath, "wb") as f:
                                    f.write(raw)
                                img_url = f"/api/generated/{fname}"
                            except Exception as e:
                                print(f"[WARN] 第{i+1}页图片落盘失败: {e}，跳过该图片")
                                img_url = None

                    if img_url:
                        images.append(img_url)
                        yield f"data: {json.dumps({'status': 'page_done', 'current': i+1, 'total': total, 'image': img_url, 'message': f'第 {i+1} 页生成成功'}, ensure_ascii=False)}\n\n"
                    else:
                        images.append(f"/api/placeholder/800/450")
                        print(f"[WARN] 第{i+1}页图片无有效URL，使用占位图")
                        yield f"data: {json.dumps({'status': 'page_fail', 'current': i+1, 'total': total, 'message': f'第 {i+1} 页生成失败（无有效图片URL）'}, ensure_ascii=False)}\n\n"

                    print(f"[INFO] 第{i+1}页图片生成成功: {page_title}")
                except Exception as e:
                    images.append(f"/api/placeholder/800/450")
                    print(f"[INFO] 第{i+1}页图片生成失败: {e}")
                    yield f"data: {json.dumps({'status': 'page_fail', 'current': i+1, 'total': total, 'message': f'第 {i+1} 页失败: {str(e)[:80]}'}, ensure_ascii=False)}\n\n"

            # 完成
            yield f"data: {json.dumps({'status': 'done', 'images': images, 'message': f'完成！共生成 {len(images)} 张图片'}, ensure_ascii=False)}\n\n"

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield f"data: {json.dumps({'status': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive',
        }
    )

def parse_color(color_desc):
    """解析颜色描述为 RGBColor"""
    from pptx.dml.color import RGBColor
    import re

    if not color_desc:
        return RGBColor(255, 255, 255)

    # 尝试解析 hex 颜色
    hex_match = re.search(r'#([0-9A-Fa-f]{6})', str(color_desc))
    if hex_match:
        hex_str = hex_match.group(1)
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return RGBColor(r, g, b)

    # 尝试解析 rgb() 格式
    rgb_match = re.search(r'rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', str(color_desc))
    if rgb_match:
        return RGBColor(int(rgb_match.group(1)), int(rgb_match.group(2)), int(rgb_match.group(3)))

    # 文字描述映射
    color_names = {
        'white': RGBColor(255, 255, 255),
        'black': RGBColor(0, 0, 0),
        'red': RGBColor(255, 0, 0),
        'orange': RGBColor(255, 107, 0),
        'blue': RGBColor(0, 102, 204),
        'green': RGBColor(0, 153, 76),
        'gray': RGBColor(128, 128, 128),
        'grey': RGBColor(128, 128, 128),
        'yellow': RGBColor(255, 204, 0),
        'semi-transparent black': RGBColor(0, 0, 0),
        'semi-transparent white': RGBColor(255, 255, 255),
    }

    for name, color in color_names.items():
        if name in str(color_desc).lower():
            return color

    return RGBColor(255, 255, 255)


def parse_crop_hint(crop_hint, img_w, img_h):
    """解析裁剪提示为 PIL crop box (left, top, right, bottom)"""
    if not crop_hint:
        return None

    import re

    # "left 40%" -> 裁剪左40%
    left_match = re.search(r'left\s+(\d+)%', crop_hint)
    if left_match:
        pct = int(left_match.group(1))
        right = int(img_w * pct / 100)
        return (0, 0, right, img_h)

    # "right 40%" -> 裁剪右40%
    right_match = re.search(r'right\s+(\d+)%', crop_hint)
    if right_match:
        pct = int(right_match.group(1))
        left = int(img_w * (100 - pct) / 100)
        return (left, 0, img_w, img_h)

    # "top 30%" -> 裁剪上30%
    top_match = re.search(r'top\s+(\d+)%', crop_hint)
    if top_match:
        pct = int(top_match.group(1))
        bottom = int(img_h * pct / 100)
        return (0, 0, img_w, bottom)

    # "bottom 30%" -> 裁剪下30%
    bottom_match = re.search(r'bottom\s+(\d+)%', crop_hint)
    if bottom_match:
        pct = int(bottom_match.group(1))
        top = int(img_h * (100 - pct) / 100)
        return (0, top, img_w, img_h)

    # "center" -> 裁剪中间50%
    if 'center' in crop_hint.lower():
        margin_x = int(img_w * 0.25)
        margin_y = int(img_h * 0.25)
        return (margin_x, margin_y, img_w - margin_x, img_h - margin_y)

    return None


def analyze_image_with_gpt4v(img_path, page_info):
    """使用 GPT-4V 分析图片，提取所有可编辑元素（文本、设计元素、图片区域）"""
    try:
        from core.api_client import UnifiedAPIClient
        import base64

        # 读取图片并转为 base64
        with open(img_path, 'rb') as f:
            img_base64 = base64.b64encode(f.read()).decode('utf-8')

        api = UnifiedAPIClient()
        config = load_config()

        page_title = page_info.get('title', '')
        page_layout = page_info.get('layout', '')
        if '模块化-三列卡片' in page_layout:
            page_layout = page_layout.replace('模块化-三列卡片', '模块化-卡片')

        # GPT-4V 分析提示词 - 提取所有类型的元素
        analysis_prompt = f"""You are a professional presentation designer. Analyze this PPT slide image and extract ALL visual elements as separate editable layers.

Page context: title="{page_title}", layout="{page_layout}"

Extract elements in these categories:

1. **text_elements** - All text blocks:
   - text: exact content (Chinese if present)
   - type: "section_label" | "title" | "subtitle" | "body" | "caption" | "hashtag" | "quote"
   - position: {{"x": 0-100, "y": 0-100}} (percentage from top-left)
   - font_size: "small"(12-14pt) | "medium"(16-20pt) | "large"(24-32pt) | "xlarge"(36-48pt)
   - color: "white" | "black" | "colored"
   - bold: true/false

2. **design_elements** - All design/decorative shapes:
   - element_type: "rectangle" | "rounded_rect" | "circle" | "line" | "divider" | "icon" | "arrow" | "badge" | "tag" | "card_bg" | "gradient_overlay"
   - position: {{"x": 0-100, "y": 0-100}}
   - size: {{"w": 0-100, "h": 0-100}} (percentage of slide)
   - color: hex color or description (e.g. "#FF6B00" or "semi-transparent black")
   - description: what this element looks like

3. **image_regions** - Distinct photo/illustration areas (NOT the main background):
   - region_type: "photo" | "icon" | "illustration" | "logo" | "screenshot" | "mockup"
   - position: {{"x": 0-100, "y": 0-100}}
   - size: {{"w": 0-100, "h": 0-100}}
   - description: what the image shows
   - crop_hint: which part of the original image to crop (e.g. "left 40%" or "center")

4. **background** - The main background description:
   - type: "photo" | "solid_color" | "gradient"
   - description: describe the background scene/style
   - dominant_colors: list of main colors

Output in this exact JSON format:
{{
  "background": {{
    "type": "photo",
    "description": "urban street scene with young people",
    "dominant_colors": ["#2C3E50", "#3498DB"]
  }},
  "text_elements": [
    {{"text": "封面", "type": "section_label", "position": {{"x": 5, "y": 8}}, "font_size": "small", "color": "white", "bold": true}},
    {{"text": "主标题", "type": "title", "position": {{"x": 10, "y": 25}}, "font_size": "xlarge", "color": "white", "bold": true}}
  ],
  "design_elements": [
    {{"element_type": "rounded_rect", "position": {{"x": 3, "y": 5}}, "size": {{"w": 12, "h": 5}}, "color": "#FF6B00", "description": "orange section label background"}},
    {{"element_type": "gradient_overlay", "position": {{"x": 0, "y": 0}}, "size": {{"w": 100, "h": 100}}, "color": "rgba(0,0,0,0.3)", "description": "dark gradient overlay for text readability"}}
  ],
  "image_regions": [
    {{"region_type": "photo", "position": {{"x": 5, "y": 30}}, "size": {{"w": 35, "h": 50}}, "description": "young woman using phone", "crop_hint": "left 40%"}}
  ]
}}

IMPORTANT RULES:
- Be thorough - extract EVERY visible element, not just text
- Design elements include: colored rectangles behind text, divider lines, decorative shapes, icon badges, card backgrounds, gradient overlays
- Image regions are for DISTINCT embedded images (not the main background photo)
- Positions are percentages (0-100) from top-left corner
- Transcribe Chinese text EXACTLY as shown
- If you see a card/module layout, identify each card as a separate design_element with "card_bg" type"""

        # 调用 GPT-4V
        response = api.chat(
            model=config.get('chat_model', 'gpt-4o'),
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": analysis_prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{img_base64}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=4000,
            timeout=120
        )

        content = response.get('content', '')
        print(f"[DEBUG] GPT-4V 分析结果前300字: {content[:300]}")

        # 提取 JSON
        import re
        json_match = re.search(r'\{[\s\S]*\}', content)
        if json_match:
            analysis = json.loads(json_match.group())

            text_count = len(analysis.get('text_elements', []))
            design_count = len(analysis.get('design_elements', []))
            image_count = len(analysis.get('image_regions', []))
            print(f"[INFO] GPT-4V 分析完成: {text_count}个文本, {design_count}个设计元素, {image_count}个图片区域")
            return analysis

        return {'text_elements': [], 'design_elements': [], 'image_regions': [], 'background': {}}

    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"[WARN] GPT-4V 分析失败: {e}")
        return {'text_elements': [], 'design_elements': [], 'image_regions': [], 'background': {}}


@app.route('/api/convert-ppt', methods=['POST'])
def convert_ppt():
    """将图片转换为 PPT（支持分层编辑）"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.dml import MSO_THEME_COLOR
        from pptx.dml.color import RGBColor
        import requests
        import base64

        # 兼容：有些前端/代理可能没带正确的 Content-Type，导致 request.json 为 None
        data = request.get_json(silent=True) or {}
        images = data.get('images', [])
        outline = data.get('outline', {})
        pages = outline.get('pages', [])
        use_layered = data.get('layered', True)  # 默认启用分层模式

        if not data:
            return jsonify({'status': 'error', 'message': '请求体不是有效的 JSON（请确保 Content-Type: application/json）'}), 400

        if not images:
            return jsonify({'status': 'error', 'message': '没有图片'}), 400

        # 创建PPT
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # 下载或处理图片
        temp_images = []
        for i, img_url in enumerate(images):
            try:
                if img_url.startswith('data:image'):
                    base64_data = img_url.split(',')[1]
                    img_data = base64.b64decode(base64_data)
                    temp_path = os.path.join(_tmpdir, f'ppt_img_{i}.png')
                    with open(temp_path, 'wb') as f:
                        f.write(img_data)
                    temp_images.append(temp_path)
                elif img_url.startswith('http') or img_url.startswith('/'):
                    # 兼容：/api/generated/xxx 这类相对路径
                    fetch_url = img_url if img_url.startswith('http') else urljoin(request.host_url, img_url.lstrip('/'))
                    response = requests.get(fetch_url, timeout=30)
                    temp_path = os.path.join(_tmpdir, f'ppt_img_{i}.png')
                    with open(temp_path, 'wb') as f:
                        f.write(response.content)
                    temp_images.append(temp_path)
                else:
                    temp_images.append(None)
            except Exception as e:
                print(f"[WARN] 图片 {i} 处理失败: {e}")
                temp_images.append(None)

        # 为每张图片创建幻灯片
        for i, img_path in enumerate(temp_images):
            if not img_path or not os.path.exists(img_path):
                continue

            page_info = pages[i] if i < len(pages) else {}

            # 添加空白幻灯片
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            if use_layered:
                # 分层模式：用 GPT-4V 分析图片，提取所有元素
                analysis = analyze_image_with_gpt4v(img_path, page_info)

                # 添加背景图片
                slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                         width=prs.slide_width, height=prs.slide_height)

                # === 添加设计元素（形状、卡片、分隔线等） ===
                for elem in analysis.get('design_elements', []):
                    try:
                        elem_type = elem.get('element_type', 'rectangle')
                        pos = elem.get('position', {'x': 0, 'y': 0})
                        size = elem.get('size', {'w': 10, 'h': 5})
                        color_desc = elem.get('color', '#FFFFFF')

                        left = Inches(13.333 * pos.get('x', 0) / 100)
                        top = Inches(7.5 * pos.get('y', 0) / 100)
                        width = Inches(13.333 * size.get('w', 10) / 100)
                        height = Inches(7.5 * size.get('h', 5) / 100)

                        # 解析颜色
                        fill_color = parse_color(color_desc)

                        if elem_type == 'gradient_overlay':
                            # 渐变遮罩 - 用半透明矩形模拟
                            shape = slide.shapes.add_shape(
                                1,  # MSO_SHAPE.RECTANGLE
                                left, top, width, height
                            )
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = fill_color
                            # 设置透明度
                            from pptx.oxml.ns import qn
                            solidFill = shape.fill._fill
                            srgbClr = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
                            if srgbClr is not None:
                                alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '30000'})  # 30% 不透明度
                                srgbClr.append(alpha)
                            shape.line.fill.background()

                        elif elem_type in ['rectangle', 'card_bg', 'tag']:
                            shape = slide.shapes.add_shape(
                                1 if elem_type == 'rectangle' else 5,  # RECTANGLE or ROUNDED_RECTANGLE
                                left, top, width, height
                            )
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = fill_color
                            shape.line.fill.background()

                        elif elem_type == 'rounded_rect':
                            shape = slide.shapes.add_shape(
                                5,  # MSO_SHAPE.ROUNDED_RECTANGLE
                                left, top, width, height
                            )
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = fill_color
                            shape.line.fill.background()

                        elif elem_type == 'circle' or elem_type == 'badge':
                            shape = slide.shapes.add_shape(
                                9,  # MSO_SHAPE.OVAL
                                left, top, min(width, height), min(width, height)
                            )
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = fill_color
                            shape.line.fill.background()

                        elif elem_type in ['line', 'divider']:
                            shape = slide.shapes.add_shape(
                                1,  # RECTANGLE as thin line
                                left, top, width, Inches(0.03)
                            )
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = fill_color
                            shape.line.fill.background()

                        elif elem_type == 'arrow':
                            shape = slide.shapes.add_shape(
                                13,  # MSO_SHAPE.RIGHT_ARROW
                                left, top, width, height
                            )
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = fill_color
                            shape.line.fill.background()

                    except Exception as e:
                        print(f"[WARN] 添加设计元素失败: {e}")

                # === 添加图片区域（裁剪嵌入的图片） ===
                for img_region in analysis.get('image_regions', []):
                    try:
                        region_type = img_region.get('region_type', 'photo')
                        pos = img_region.get('position', {'x': 0, 'y': 0})
                        size = img_region.get('size', {'w': 30, 'h': 30})
                        crop_hint = img_region.get('crop_hint', '')

                        left = Inches(13.333 * pos.get('x', 0) / 100)
                        top = Inches(7.5 * pos.get('y', 0) / 100)
                        width = Inches(13.333 * size.get('w', 30) / 100)
                        height = Inches(7.5 * size.get('h', 30) / 100)

                        # 从原始图片裁剪对应区域
                        from PIL import Image as PILImage
                        pil_img = PILImage.open(img_path)
                        img_w, img_h = pil_img.size

                        # 解析 crop_hint
                        crop_box = parse_crop_hint(crop_hint, img_w, img_h)
                        if crop_box:
                            cropped = pil_img.crop(crop_box)
                            crop_path = f'/tmp/ppt_crop_{i}_{region_type}.png'
                            cropped.save(crop_path)

                            # 添加圆角效果（如果是 mockup 或 screenshot）
                            if region_type in ['mockup', 'screenshot']:
                                pic = slide.shapes.add_picture(crop_path, left, top, width, height)
                            else:
                                pic = slide.shapes.add_picture(crop_path, left, top, width, height)

                            os.remove(crop_path)

                    except Exception as e:
                        print(f"[WARN] 添加图片区域失败: {e}")

                # === 添加文本元素（可编辑文本框） ===
                for elem in analysis.get('text_elements', []):
                    try:
                        text = elem.get('text', '')
                        elem_type = elem.get('type', 'body')
                        pos = elem.get('position', {'x': 10, 'y': 10})
                        font_size_cat = elem.get('font_size', 'medium')
                        color = elem.get('color', 'white')
                        is_bold = elem.get('bold', False)

                        # 位置转换为英寸
                        left = Inches(13.333 * pos.get('x', 10) / 100)
                        top = Inches(7.5 * pos.get('y', 10) / 100)

                        # 根据类型设置尺寸
                        size_map = {
                            'section_label': (Inches(2.5), Inches(0.5)),
                            'title': (Inches(10), Inches(1.2)),
                            'subtitle': (Inches(8), Inches(0.7)),
                            'body': (Inches(8), Inches(3)),
                            'caption': (Inches(6), Inches(0.5)),
                            'hashtag': (Inches(8), Inches(0.6)),
                            'quote': (Inches(9), Inches(2)),
                        }
                        width, height = size_map.get(elem_type, (Inches(8), Inches(2)))

                        # 字体大小映射
                        font_size_map = {
                            'small': Pt(13),
                            'medium': Pt(18),
                            'large': Pt(28),
                            'xlarge': Pt(44)
                        }
                        font_size = font_size_map.get(font_size_cat, Pt(18))

                        # 颜色映射
                        color_map = {
                            'white': RGBColor(255, 255, 255),
                            'black': RGBColor(0, 0, 0),
                            'colored': RGBColor(255, 107, 0)
                        }
                        font_color = color_map.get(color, RGBColor(255, 255, 255))

                        # 添加文本框
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf = txBox.text_frame
                        tf.word_wrap = True

                        # 处理多行文本
                        lines = text.split('\n')
                        for line_idx, line in enumerate(lines):
                            if line_idx == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            p.text = line
                            p.font.size = font_size
                            p.font.bold = is_bold or (elem_type in ['title', 'section_label'])
                            p.font.color.rgb = font_color

                            if elem_type == 'section_label':
                                p.alignment = 1  # 居中
                            elif elem_type == 'quote':
                                p.alignment = 2  # 右对齐
                            else:
                                p.alignment = 0  # 左对齐

                    except Exception as e:
                        print(f"[WARN] 添加文本元素失败: {e}")
            else:
                # 旧模式：直接插入整张图片
                slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                         width=prs.slide_width, height=prs.slide_height)

        # 保存PPT
        output_path = os.path.join(BASE_DIR, 'static', 'presentation.pptx')
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)

        # 清理临时文件
        for img_path in temp_images:
            if img_path and os.path.exists(img_path):
                try:
                    os.remove(img_path)
                except:
                    pass

        return jsonify({
            'status': 'success',
            'download_url': '/api/download/presentation.pptx',
            'message': f'PPT 已生成（分层模式），共 {len(temp_images)} 页'
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/convert-ppt-v2', methods=['POST'])
def convert_ppt_v2():
    """V2：背景去字 + 元素/图片/文字分层可编辑的 PPT 转换"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        sys.path.insert(0, BASE_DIR)

        import requests
        from core.api_client import APIClient
        from v2.layout_extractor import analyze_layout, build_background_redraw_prompt
        from v2.background_cleanup import cleanup_background, cleanup_background_ai_only, redraw_background
        from v2.ppt_builder import crop_regions, build_deck_v2
        from v2.storage import save_bytes, url_for_generated

        data = request.get_json(silent=True) or {}
        images = data.get('images', [])
        options = data.get('options', {}) or {}

        if not images:
            return jsonify({'status': 'error', 'message': '没有图片'}), 400

        # 测试/探活：避免在自动化测试里触发 OCR/多模态/修复的昂贵流程
        if options.get("dry_run"):
            return jsonify({'status': 'success', 'message': 'dry_run ok'}), 200

        stream = bool(options.get("stream"))

        # 读取配置
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)

        chat_api_url = config.get('chat_api_base_url') or config.get('api_base_url')
        chat_api_key = config.get('chat_api_key') or config.get('api_key')
        image_api_url = config.get('image_api_base_url') or config.get('api_base_url')
        image_api_key = config.get('image_api_key') or config.get('api_key')

        if not chat_api_url or not chat_api_key:
            return jsonify({'status': 'error', 'message': '未配置对话模型 API'}), 400
        if not image_api_url or not image_api_key:
            return jsonify({'status': 'error', 'message': '未配置图片模型 API'}), 400

        chat_model = config.get('chat_model', 'gpt-4o')

        # inpaint 模型列表：可由前端 options 覆盖
        # 背景处理模式：
        # - redraw（默认）：背景重绘（无涂抹痕迹），再叠加卡片/手机/文字等分层
        # - inpaint_ai：走编辑模型 inpaint 去字
        # - inpaint_local：本地 OpenCV inpaint（免费但有涂抹痕迹）
        background_mode = options.get("background_mode", "redraw")  # "redraw" | "inpaint_ai" | "inpaint_local"
        inpaint_models = (
            options.get('inpaint_models')
            or config.get('inpaint_models')
            or ['nano-banana-pro', 'qwen-image-edit-plus', 'mj_fast_inpaint']
        )
        if background_mode == "inpaint_local":
            inpaint_models = []

        # V2 处理：逐页生成 slide
        api_chat = APIClient(chat_api_url, chat_api_key)
        api_image = APIClient(image_api_url, image_api_key)

        slides_payload = []

        # 预过滤得到总数（用于进度显示）
        valid_images = [
            u
            for u in images
            if u and isinstance(u, str) and (not u.startswith("/api/placeholder/"))
        ]
        total = len(valid_images)

        def _sse(obj: dict) -> str:
            return f"data: {json.dumps(obj, ensure_ascii=False)}\n\n"

        def resolve_image_to_local_path(img_url: str, idx: int) -> str:
            """把 dataUrl/http/url/path 转成可读的本地 png 路径"""
            import base64
            # 1) data url
            if img_url.startswith('data:image'):
                b64_data = img_url.split(',', 1)[1]
                raw = base64.b64decode(b64_data)
                p = os.path.join(_tmpdir, f'v2_input_{idx}.png')
                with open(p, 'wb') as f:
                    f.write(raw)
                return p
            # 2) already generated on server
            if img_url.startswith('/api/generated/'):
                fname = img_url.split('/api/generated/', 1)[1]
                p = os.path.join(GENERATED_DIR, fname)
                return p
            # 3) http(s)
            if img_url.startswith('http'):
                resp = requests.get(img_url, timeout=30)
                p = os.path.join(_tmpdir, f'v2_input_{idx}.png')
                with open(p, 'wb') as f:
                    f.write(resp.content)
                return p
            # 4) fallback: treat as path
            return img_url

        def _process_one(img_url: str, idx0: int, idx1: int):
            """
            idx0: 从 0 开始的索引（用于临时文件命名）
            idx1: 从 1 开始的进度索引（用于 UI 展示）
            """
            local_img = resolve_image_to_local_path(img_url, idx0)
            if not os.path.exists(local_img):
                return None, f"图片不存在，已跳过：{img_url}"

            layout = analyze_layout(api_chat, chat_model=chat_model, image_path=local_img)

            # 背景处理
            if background_mode == "redraw":
                image_model = config.get("image_model", "gpt-image-2")
                prompt = build_background_redraw_prompt(api_chat, chat_model, local_img)
                clean_bg_bytes = redraw_background(api_image, image_model=image_model, prompt=prompt, size="1792x1024")
                mask_bytes = b""
            elif background_mode == "inpaint_ai":
                out, mask_bytes = cleanup_background_ai_only(
                    api_client=api_image,
                    inpaint_models=inpaint_models,
                    image_path=local_img,
                    layout=layout,
                    size="1792x1024",
                    per_try_timeout_s=int(options.get("inpaint_timeout_s") or 60),
                    retries=int(options.get("inpaint_retries") or 1),
                )
                if out is None:
                    # AI 去字超时/失败：自动降级为背景重绘，避免整批卡死
                    image_model = config.get("image_model", "gpt-image-2")
                    prompt = build_background_redraw_prompt(api_chat, chat_model, local_img)
                    clean_bg_bytes = redraw_background(api_image, image_model=image_model, prompt=prompt, size="1792x1024")
                    mask_bytes = b""
                else:
                    clean_bg_bytes = out
            else:
                # inpaint_local
                clean_bg_bytes, mask_bytes = cleanup_background(
                    api_client=api_image,
                    inpaint_models=[],
                    image_path=local_img,
                    layout=layout,
                    size="1792x1024",
                )

            # 落盘背景图（调试/复用）
            bg_fname = save_bytes(GENERATED_DIR, clean_bg_bytes, ext="png", prefix="bg_clean")
            bg_url = url_for_generated(bg_fname)

            # 裁切 phone/photo 等图片区域，作为可编辑图片层
            cropped = crop_regions(local_img, layout.get("images", []), out_dir=GENERATED_DIR)

            slide_payload = {
                "clean_bg_png_bytes": clean_bg_bytes,
                "layout": layout,
                "cropped_images": cropped,
                "debug": {"bg_url": bg_url},
            }
            return slide_payload, ""

        if stream:
            def generate_stream():
                yield _sse({"status": "start", "total": total, "message": f"开始转换，共 {total} 张图片"})
                for i, img_url in enumerate(valid_images):
                    idx = i + 1
                    yield _sse({"status": "progress", "current": idx, "total": total, "stage": "layout", "message": f"解析版式 {idx}/{total}"})
                    try:
                        slide_payload, warn = _process_one(img_url, i, idx)
                        if warn:
                            yield _sse({"status": "page_fail", "current": idx, "total": total, "message": warn})
                            continue
                        slides_payload.append(slide_payload)
                        yield _sse({"status": "page_done", "current": idx, "total": total, "message": f"第 {idx}/{total} 张完成"})
                    except Exception as e:
                        yield _sse({"status": "page_fail", "current": idx, "total": total, "message": f"第 {idx}/{total} 张失败：{str(e)[:120]}"})

                if not slides_payload:
                    yield _sse({"status": "error", "message": "没有可转换的有效图片（可能都生成失败或是占位图）"})
                    return

                yield _sse({"status": "progress", "stage": "ppt", "message": f"正在合成 PPT（{len(slides_payload)} 页）"})
                output_filename = f"presentation_v2_{int(time.time())}_{uuid.uuid4().hex[:6]}.pptx"
                output_path = os.path.join(BASE_DIR, "static", output_filename)
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                build_deck_v2(output_path, slides_payload)
                yield _sse({"status": "done", "download_url": f"/api/download/{output_filename}", "pages": len(slides_payload), "message": f"PPT 已生成，共 {len(slides_payload)} 页"})

            return Response(
                stream_with_context(generate_stream()),
                mimetype="text/event-stream",
                headers={
                    "Cache-Control": "no-cache",
                    "X-Accel-Buffering": "no",
                    "Connection": "keep-alive",
                },
            )

        for i, img_url in enumerate(valid_images):
            slide_payload, warn = _process_one(img_url, i, i + 1)
            if warn:
                print(f"[WARN] {warn}")
                continue
            slides_payload.append(slide_payload)

        # 4) 写 PPTX
        if not slides_payload:
            return jsonify({'status': 'error', 'message': '没有可转换的有效图片（可能都生成失败或是占位图）'}), 400

        output_filename = f"presentation_v2_{int(time.time())}_{uuid.uuid4().hex[:6]}.pptx"
        output_path = os.path.join(BASE_DIR, 'static', output_filename)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        build_deck_v2(output_path, slides_payload)

        return jsonify(
            {
                'status': 'success',
                'download_url': f'/api/download/{output_filename}',
                'message': f'PPT 已生成（V2 分层去字），共 {len(slides_payload)} 页',
            }
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/analyze-layout-v2', methods=['POST'])
def analyze_layout_v2():
    """调试接口：返回 V2 layout JSON（便于迭代 bbox/分类）"""
    try:
        import sys
        sys.path.insert(0, os.path.join(BASE_DIR, '..'))
        sys.path.insert(0, BASE_DIR)
        from core.api_client import APIClient
        from v2.layout_extractor import analyze_layout

        data = request.get_json(silent=True) or {}
        img_url = (data.get('image') or '').strip()
        if not img_url:
            return jsonify({'status': 'error', 'message': '缺少 image'}), 400

        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)
        chat_api_url = config.get('chat_api_base_url') or config.get('api_base_url')
        chat_api_key = config.get('chat_api_key') or config.get('api_key')
        chat_model = config.get('chat_model', 'gpt-4o')
        api_chat = APIClient(chat_api_url, chat_api_key)

        # 复用 convert_ppt_v2 的解析逻辑
        if img_url.startswith('/api/generated/'):
            fname = img_url.split('/api/generated/', 1)[1]
            local_img = os.path.join(GENERATED_DIR, fname)
        else:
            local_img = img_url
        layout = analyze_layout(api_chat, chat_model=chat_model, image_path=local_img)
        return jsonify({'status': 'success', 'layout': layout})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500

@app.route('/api/download/<filename>')
def download_file(filename):
    """下载生成的文件"""
    try:
        file_path = os.path.join(BASE_DIR, 'static', filename)
        if os.path.exists(file_path):
            from flask import send_file
            return send_file(file_path, as_attachment=True)
        else:
            return jsonify({'status': 'error', 'message': '文件不存在'}), 404
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500


def _resolve_outline_from_payload(data):
    session_id = data.get('session_id', '')
    outline = data.get('outline')

    if not outline and session_id:
        store = _uploaded_outline_cache.get(session_id) or _load_outline_session(session_id)
        if store:
            outline = store.get('outline')
            _uploaded_outline_cache[session_id] = store

    return outline


def _page_asset_bucket(asset_role):
    raw = (asset_role or 'image').strip().lower()
    return 'background' if raw in {'background', 'bg', 'background_image', 'background_asset'} else 'image'


def _normalize_asset_list(value):
    if not value:
        return []
    if isinstance(value, list):
        return [item for item in value if item]
    return [value]


def _append_page_asset(page, asset_url, asset_role):
    """Append an asset URL into the correct page-level editable PPTX asset slot."""
    if asset_role == 'background':
        assets = page.get('background_images') or page.get('background_image') or []
        if not isinstance(assets, list):
            assets = [assets] if assets else []
        assets.append(asset_url)
        page['background_images'] = assets
        page.pop('background_image', None)
        return assets

    assets = page.get('fixed_images') or []
    if not isinstance(assets, list):
        assets = [assets] if assets else []
    assets.append(asset_url)
    page['fixed_images'] = assets
    return assets


def _write_page_asset_bytes(raw, original_name='generated-page-asset.png'):
    ext = os.path.splitext(original_name or '')[1].lower()
    if ext not in {'.png', '.jpg', '.jpeg', '.webp'}:
        ext = '.png'
    asset_dir = os.path.join(GENERATED_DIR, 'page_assets')
    os.makedirs(asset_dir, exist_ok=True)
    filename = f"ai_asset_{int(time.time())}_{uuid.uuid4().hex[:8]}{ext}"
    output_path = os.path.join(asset_dir, filename)
    with open(output_path, 'wb') as out:
        out.write(raw)
    return f"/api/generated/page_assets/{filename}", filename


def _asset_prompt_for_page(page, asset_role='background', style_settings=None):
    style_settings = style_settings or {}
    page_title = clean_text(page.get('title')) if 'clean_text' in globals() else str(page.get('title') or '')
    page_content = clean_text(page.get('content')) if 'clean_text' in globals() else str(page.get('content') or '')
    page_brief = clean_text(page.get('brief')) if 'clean_text' in globals() else str(page.get('brief') or '')
    page_layout = clean_text(page.get('layout') or page.get('layout_skill')) if 'clean_text' in globals() else str(page.get('layout') or page.get('layout_skill') or '')
    style_desc = style_settings.get('description') or ''
    brand_name = style_settings.get('brandName') or ''
    color_scheme = style_settings.get('colorScheme') or ''
    scene_description = _build_scene_description(page_title, page_content)

    if asset_role == 'background':
        asset_job = """Generate ONLY a background/key visual image for this PPT page.
- No readable text, no titles, no captions, no UI labels.
- Leave visual quiet zones so native PPTX text can be placed on top later.
- It may be photographic, abstract, lifestyle, or atmospheric, but must match the page topic.
- This image will be inserted as a replaceable bottom-layer picture object in an editable PPTX."""
    else:
        asset_job = """Generate ONLY a page material image for a PPT image slot.
- No readable text, no titles, no captions, no full-slide layout.
- Create a visual asset such as product scene, evidence thumbnail, social screenshot-style placeholder, or supporting illustration.
- This image will be inserted as a replaceable picture object inside an editable PPTX layout."""

    return f"""{asset_job}

Page title: {page_title}
Page layout signal: {page_layout}
Page brief: {page_brief}
Page content: {page_content[:1600]}

{scene_description}

Brand/style context:
- Brand name: {brand_name or 'not specified'}
- Color preference: {color_scheme or 'auto'}
- Style description: {style_desc or 'professional campaign pitch deck, polished but not stock-like'}

Hard constraints:
- 16:9 aspect ratio.
- No PowerPoint chrome, no slide title text, no bullet text, no fake data labels.
- Do not compose the whole slide. Produce only the asset layer.
- Avoid generic stock-photo cliches unless the page content explicitly calls for them."""


def _result_to_local_asset(result):
    import base64
    import urllib.request

    images = []
    if isinstance(result, dict):
        images.extend(result.get('images') or [])
        images.extend(result.get('data') or [])
    for img in images:
        if not isinstance(img, dict):
            continue
        if img.get('b64_json'):
            raw = base64.b64decode(img['b64_json'])
            return _write_page_asset_bytes(raw)
        if img.get('url'):
            with urllib.request.urlopen(img['url'], timeout=60) as response:
                raw = response.read()
            return _write_page_asset_bytes(raw)
    raise RuntimeError('图片模型未返回可保存的图片')


def _generate_page_asset_for_page(page, asset_role, style_settings=None, dry_run=False):
    """Generate one page asset and return local asset metadata without mutating the page."""
    import base64
    import sys

    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    style_settings = style_settings or {}
    prompt = _asset_prompt_for_page(page, asset_role, style_settings)

    if dry_run:
        raw = base64.b64decode('iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9WZ0kAAAAASUVORK5CYII=')
        asset_url, filename = _write_page_asset_bytes(raw, 'dry-run-page-asset.png')
        provider = 'dry_run'
    else:
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)
        image_api_url = config.get('image_api_base_url') or config.get('api_base_url')
        image_api_key = config.get('image_api_key') or config.get('api_key')
        if not image_api_url or not image_api_key:
            raise ValueError('请先配置图片生成模型的 API 设置')

        api = APIClient(image_api_url, image_api_key)
        image_model = config.get('image_model', 'gpt-image-2')
        page_title = str(page.get('title') or '')
        page_content = str(page.get('content') or '')
        page_type = str(page.get('type') or page.get('page_type') or '')
        page_layout = str(page.get('layout') or page.get('layout_skill') or '')
        product_image = style_settings.get('productImage') or ''

        messages = None
        if product_image and isinstance(product_image, str) and product_image.startswith('data:image') and _page_needs_product_image(page_title, page_content, page_type, page_layout):
            messages = [{"role": "user", "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {"url": product_image}},
            ]}]

        result = api.image_generate(
            model=image_model,
            prompt=prompt,
            size='1792x1024',
            n=1,
            messages=messages,
            timeout=120,
        )
        asset_url, filename = _result_to_local_asset(result)
        provider = image_model

    return {
        'asset_url': asset_url,
        'filename': filename,
        'provider': provider,
        'prompt': prompt,
    }


def _build_agent_deck_plan(outline):
    """Build the lightweight workbench plan used by the HTML Agent UI."""
    import sys
    from collections import Counter

    sys.path.insert(0, BASE_DIR)
    from services.dashiai_theme_seed import dashiai_seed_summary, match_restyled_dashiai_theme_seeds
    from services.slide_render_plan import build_deck_render_plan

    plans = build_deck_render_plan(outline)
    seed_summary = dashiai_seed_summary()
    pages_payload = []
    role_counter = Counter()
    archetype_counter = Counter()
    pages_missing_background = 0
    pages_missing_material = 0

    for plan in plans:
        source_page = plan.intent.source_page if isinstance(plan.intent.source_page, dict) else {}
        backgrounds = _normalize_asset_list(source_page.get('background_images') or source_page.get('background_image'))
        materials = _normalize_asset_list(source_page.get('fixed_images'))
        role_counter[plan.intent.page_role] += 1
        archetype_counter[plan.visual_profile.archetype] += 1

        needs = []
        if plan.intent.page_role in {'开场定调', '创意主张', '案例证据', '视频素材'} and not materials:
            needs.append('material_image')
            pages_missing_material += 1
        if plan.visual_profile.image_treatment in {'hero', 'cutout', 'evidence', 'thumbnail'} and not backgrounds:
            needs.append('background_image')
            pages_missing_background += 1
        if plan.intent.page_role == '数据结果' and len(plan.card_texts) < 2:
            needs.append('metric_structure_review')

        if plan.visual_profile.archetype in {
            'hero_photo_claim',
            'strategy_claim_collage',
            'evidence_wall',
            'section_divider',
            'metric_dashboard',
        }:
            recommended_renderer = 'reference-template'
        else:
            recommended_renderer = 'pptxgenjs'

        media_slots = [
            {
                'key': 'background_images',
                'label': '背景图',
                'kind': 'image',
                'count': len(backgrounds),
                'source': 'user_upload_or_image_2',
            },
            {
                'key': 'fixed_images',
                'label': '素材图',
                'kind': 'image',
                'count': len(materials),
                'source': 'user_upload_or_image_2',
            },
            {
                'key': 'video_links',
                'label': '视频链接',
                'kind': 'video_link',
                'count': len(plan.video_links),
                'source': 'outline_or_manual',
            },
        ]
        safe_controls = [
            {'key': 'page_role', 'label': '页面 Skill', 'type': 'select', 'scope': 'page_intent'},
            {'key': 'title', 'label': '标题', 'type': 'text', 'scope': 'copy'},
            {'key': 'brief', 'label': 'Brief', 'type': 'textarea', 'scope': 'copy'},
            {'key': 'content', 'label': '正文', 'type': 'textarea', 'scope': 'copy'},
            {'key': 'background_images', 'label': '背景图槽', 'type': 'image_slot', 'scope': 'asset'},
            {'key': 'fixed_images', 'label': '素材图槽', 'type': 'image_slot', 'scope': 'asset'},
        ]
        agent_actions = [
            'match_reference_template',
            'borrow_dashiai_component_seed',
            'generate_or_upload_page_asset',
            'render_html_preview',
            'transcribe_editable_pptx',
            'qa_native_objects',
        ]
        selected_seed_key = str(source_page.get('dashiai_seed_key') or '').strip()
        dashiai_seed_rejected = bool(source_page.get('dashiai_seed_rejected'))
        dashiai_seeds = [] if dashiai_seed_rejected else match_restyled_dashiai_theme_seeds(
            plan.intent.page_role,
            style_tokens=plan.style_tokens,
            visual_profile=plan.visual_profile,
            limit=3,
        )
        selected_dashiai_seed = None
        if selected_seed_key:
            selected_dashiai_seed = next((seed for seed in dashiai_seeds if seed.get('key') == selected_seed_key), None)
        if not selected_dashiai_seed and selected_seed_key:
            selected_dashiai_seed = {
                'key': selected_seed_key,
                'label': source_page.get('dashiai_seed_label') or selected_seed_key,
                'adaptation_status': 'selected_seed_not_in_current_top_candidates',
            }
        seed_for_restyle = selected_dashiai_seed or (dashiai_seeds[0] if dashiai_seeds else {})
        feibo_restyle_policy = (seed_for_restyle or {}).get('feibo_restyle') or {
            'status': 'no_dashiai_seed_use_feibo_template_first',
            'target_style_id': plan.visual_profile.reference_style_id,
            'target_style_label': plan.visual_profile.reference_style_label,
        }
        component_spec = plan.component_spec if isinstance(getattr(plan, 'component_spec', None), dict) else {}
        feibo_restyle_policy = component_spec.get('feibo_restyle') or feibo_restyle_policy

        pages_payload.append({
            'page_index': plan.index - 1,
            'index': plan.index,
            'title': plan.title,
            'page_role': plan.intent.page_role,
            'layout_label': plan.intent.layout_label,
            'visual_archetype': plan.visual_profile.archetype,
            'visual_density': plan.visual_profile.density,
            'image_treatment': plan.visual_profile.image_treatment,
            'composition': plan.visual_profile.composition,
            'reference_style_id': plan.visual_profile.reference_style_id,
            'reference_style_label': plan.visual_profile.reference_style_label,
            'material_count': len(materials),
            'background_count': len(backgrounds),
            'video_count': len(plan.video_links),
            'content_line_count': len(plan.body_lines),
            'card_count': len(plan.card_texts),
            'recommended_renderer': recommended_renderer,
            'internal_pptx_adapter': recommended_renderer,
            'export_route': 'HTML工作台转写PPTX',
            'agent_status': 'needs_assets' if needs else 'ready_for_preview',
            'needs': needs,
            'selected_dashiai_seed_key': selected_seed_key,
            'dashiai_seed_rejected': dashiai_seed_rejected,
            'component_spec': component_spec,
            'control_manifest': {
                'schema': 'stylemind.page_manifest.v1',
                'copy_fields': ['title', 'brief', 'content'],
                'media_slots': media_slots,
                'safe_controls': safe_controls,
                'agent_actions': agent_actions,
                'dashiai_theme_seeds': dashiai_seeds,
                'selected_dashiai_seed': selected_dashiai_seed,
                'selected_dashiai_seed_key': selected_seed_key,
                'dashiai_seed_rejected': dashiai_seed_rejected,
                'feibo_restyle': feibo_restyle_policy,
                'component_spec': component_spec,
                'policy_flags': [
                    'preserve_outline_parser',
                    'feibo_reference_first',
                    'html_preview_then_editable_pptx',
                    'no_full_slide_png_delivery',
                ],
                'style_source': 'Feibo reference first; DashiAI module workflow as supplement',
            },
        })

    return {
        'schema': 'stylemind.deck_plan.v1',
        'title': outline.get('title') or 'PPT大纲',
        'page_count': len(plans),
        'pages': pages_payload,
        'role_distribution': dict(role_counter.most_common()),
        'archetype_distribution': dict(archetype_counter.most_common()),
        'asset_summary': {
            'pages_missing_background': pages_missing_background,
            'pages_missing_material': pages_missing_material,
            'pages_with_any_asset': sum(1 for p in pages_payload if p['material_count'] or p['background_count']),
        },
        'dashiai_seed_summary': seed_summary,
        'pipeline': 'outline_to_deck_plan_to_agent_workbench',
        'available_actions': [
            'build_deck_plan',
            'qa_deck',
            'match_reference_template',
            'borrow_dashiai_component_seed',
            'render_preview',
            'generate_assets',
            'export_pptx',
        ],
    }


@app.route('/api/agent/run', methods=['POST'])
def agent_run():
    """Small deterministic Agent action layer for the HTML workbench."""
    try:
        data = request.get_json(silent=True) or {}
        action = (data.get('action') or 'build_deck_plan').strip().lower()
        outline = _resolve_outline_from_payload(data)

        if not isinstance(outline, dict):
            return jsonify({'status': 'error', 'message': '请提供 outline 或有效的 session_id'}), 400
        if not outline.get('pages'):
            return jsonify({'status': 'error', 'message': 'outline 中没有 pages 数据'}), 400

        deck_plan = _build_agent_deck_plan(outline)
        if action in {'build_deck_plan', 'classify_pages', 'plan'}:
            return jsonify({
                'status': 'success',
                'action': action,
                'deck_plan': deck_plan,
                'message': 'Agent 工作台计划已生成',
            })
        if action in {'qa_deck', 'qa', 'inspect'}:
            needs_pages = [p for p in deck_plan['pages'] if p['needs']]
            return jsonify({
                'status': 'success',
                'action': action,
                'deck_plan': deck_plan,
                'qa': {
                    'page_count': deck_plan['page_count'],
                    'needs_attention_count': len(needs_pages),
                    'needs_attention_pages': needs_pages[:24],
                    'asset_summary': deck_plan['asset_summary'],
                    'editable_delivery_policy': 'PPTX text/shapes/images must stay native editable objects',
                },
                'message': 'Agent QA 已完成',
            })
        if action in {'borrow_dashiai_component_seed', 'dashiai_seed', 'theme_seed'}:
            seed_pages = [
                {
                    'page_index': page['page_index'],
                    'index': page['index'],
                    'title': page['title'],
                    'page_role': page['page_role'],
                    'selected_dashiai_seed_key': page.get('selected_dashiai_seed_key'),
                    'dashiai_seed_rejected': page.get('dashiai_seed_rejected'),
                    'dashiai_theme_seeds': page.get('control_manifest', {}).get('dashiai_theme_seeds', []),
                    'feibo_restyle': page.get('control_manifest', {}).get('feibo_restyle'),
                    'component_spec': page.get('component_spec') or page.get('control_manifest', {}).get('component_spec'),
                }
                for page in deck_plan['pages']
            ]
            return jsonify({
                'status': 'success',
                'action': action,
                'deck_plan': deck_plan,
                'dashiai_seed_summary': deck_plan.get('dashiai_seed_summary', {}),
                'pages': seed_pages,
                'message': 'DashiAI 组件种子已匹配，仍需按飞博参考做视觉改写',
            })

        return jsonify({
            'status': 'error',
            'message': f'暂不支持的 Agent action: {action}',
            'available_actions': deck_plan['available_actions'],
        }), 400
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/render-plan-preview', methods=['POST'])
def render_plan_preview():
    """从同一个 RenderPlan 生成结构化 PNG 预览，不调用整页生图模型。"""
    try:
        import sys
        from collections import Counter

        sys.path.insert(0, BASE_DIR)
        from services.renderers.plan_preview_renderer import render_plan_preview_data_url
        from services.slide_render_plan import build_deck_render_plan

        data = request.get_json(silent=True) or {}
        outline = _resolve_outline_from_payload(data)

        if not isinstance(outline, dict):
            return jsonify({'status': 'error', 'message': '请提供 outline 或有效的 session_id'}), 400
        if not outline.get('pages'):
            return jsonify({'status': 'error', 'message': 'outline 中没有 pages 数据'}), 400

        plans = build_deck_render_plan(outline)
        previews = []
        for plan in plans:
            previews.append({
                'page_index': plan.index - 1,
                'index': plan.index,
                'title': plan.title,
                'page_role': plan.intent.page_role,
                'layout_label': plan.intent.layout_label,
                'visual_profile': {
                    'archetype': plan.visual_profile.archetype,
                    'density': plan.visual_profile.density,
                    'image_treatment': plan.visual_profile.image_treatment,
                    'composition': plan.visual_profile.composition,
                    'reference_style_id': plan.visual_profile.reference_style_id,
                    'reference_style_label': plan.visual_profile.reference_style_label,
                },
                'renderer': 'render_plan_preview_png',
                'image': render_plan_preview_data_url(plan),
            })

        role_distribution = dict(Counter(plan.intent.page_role for plan in plans).most_common())

        return jsonify({
            'status': 'success',
            'pages': len(plans),
            'previews': previews,
            'role_distribution': role_distribution,
            'pipeline': 'outline_to_slide_layout_spec_to_render_plan',
            'preview_pipeline': 'outline_to_render_plan_to_preview_png',
            'final_pptx_pipeline': 'outline_to_render_plan_to_native_pptx',
            'uses_full_page_image_model': False,
            'message': '结构化预览已生成：与可编辑 PPTX 共用 RenderPlan，不调用整页 image-2'
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/export-pptx-structured', methods=['POST'])
def export_pptx_structured():
    """从结构化大纲直接导出可编辑 PPTX（不把整页 PNG 塞进 PPT）。"""
    try:
        import sys
        sys.path.insert(0, BASE_DIR)
        from services.structured_pptx_renderer import build_structured_pptx

        data = request.get_json(silent=True) or {}
        outline = _resolve_outline_from_payload(data)
        renderer = (data.get('renderer') or 'python-pptx').strip().lower()
        if renderer in {'pptxgen', 'pptxgenjs', 'node'}:
            renderer = 'pptxgenjs'
        elif renderer in {'reference-template', 'reference_template', 'template', 'automizer', 'pptx-automizer', 'pptxautomizer'}:
            renderer = 'reference-template'
        elif renderer in {'html-dom', 'html_dom', 'dom', 'html'}:
            renderer = 'html-dom'
        elif renderer not in {'python-pptx', 'python_pptx', 'python'}:
            return jsonify({'status': 'error', 'message': f'不支持的 PPTX 转写策略: {renderer}'}), 400
        else:
            renderer = 'python-pptx'
        renderer_display = {
            'reference-template': '参考质感优先',
            'html-dom': 'HTML 同源转写',
            'pptxgenjs': '结构化绘制',
            'python-pptx': '稳定基础转写',
        }.get(renderer, 'PPTX 转写')

        if not isinstance(outline, dict):
            return jsonify({'status': 'error', 'message': '请提供 outline 或有效的 session_id'}), 400
        if not outline.get('pages'):
            return jsonify({'status': 'error', 'message': 'outline 中没有 pages 数据'}), 400

        output_filename = f"structured_{renderer.replace('-', '_')}_{int(time.time())}_{uuid.uuid4().hex[:6]}.pptx"
        output_path = os.path.join(BASE_DIR, 'static', output_filename)
        render_plan_path = None
        renderer_meta = {}

        if renderer == 'pptxgenjs':
            from services.renderers.pptxgenjs_renderer import PptxGenJSRendererError, render_outline_to_pptxgenjs

            try:
                result = render_outline_to_pptxgenjs(outline, output_path)
            except PptxGenJSRendererError as e:
                return jsonify({'status': 'error', 'message': f'{renderer_display}失败: {str(e)[:800]}'}), 500
            render_plan_path = result.render_plan_path
        elif renderer == 'reference-template':
            from services.renderers.reference_template_renderer import (
                ReferenceTemplateRendererError,
                render_outline_to_reference_template,
            )

            max_pages_raw = data.get('reference_template_max_pages') or data.get('template_max_pages')
            try:
                max_pages = int(max_pages_raw) if max_pages_raw else None
            except (TypeError, ValueError):
                max_pages = None
            if max_pages is not None and max_pages <= 0:
                max_pages = None

            try:
                result = render_outline_to_reference_template(outline, output_path, max_pages=max_pages)
            except ReferenceTemplateRendererError as e:
                return jsonify({'status': 'error', 'message': f'{renderer_display}失败: {str(e)[:800]}'}), 500
            render_plan_path = result.render_plan_path
            renderer_meta = {
                'quality_strategy': result.strategy,
                'template_match_status': 'cleaned_reference_templates_experimental_requires_visual_qa',
                'reference_template_report_path': f"/api/generated/render_plans/{os.path.basename(result.report_path)}",
                'reference_template_selected_pages': result.selected_count,
                'known_limitations': list(result.known_limitations),
            }
        elif renderer == 'html-dom':
            from services.renderers.html_dom_renderer import (
                HtmlDomRendererError,
                render_outline_to_html_dom,
            )

            try:
                result = render_outline_to_html_dom(outline, output_path)
            except HtmlDomRendererError as e:
                return jsonify({'status': 'error', 'message': f'{renderer_display}失败: {str(e)[:800]}'}), 500
            render_plan_path = result.render_plan_path
            renderer_meta = {
                'quality_strategy': result.strategy,
                'html_dom_status': 'experimental_html_preview_to_editable_pptx_requires_visual_qa',
                'html_dom_report_path': f"/api/generated/render_plans/{os.path.basename(result.report_path)}",
                'html_dom_deck_path': f"/api/generated/render_plans/{os.path.basename(result.deck_dir)}/index.html",
                'html_dom_text_objects': result.text_objects,
                'html_dom_shape_objects': result.shape_objects,
                'html_dom_image_objects': result.image_objects,
                'known_limitations': list(result.known_limitations),
            }
        else:
            build_structured_pptx(outline, output_path)

        response = {
            'status': 'success',
            'download_url': f'/api/download/{output_filename}',
            'filename': output_filename,
            'renderer': renderer,
            'pages': len(outline.get('pages', [])),
            'pipeline': 'outline_to_render_plan_to_native_pptx',
            'preview_pipeline': 'outline_to_render_plan_to_preview_png',
            'fallback_preview_pipeline': 'png_pdf_preview_only',
            'uses_full_page_image_model': False,
            'asset_image_model_role': '素材图生成/补图，不负责整页排版和文字',
            'render_plan_path': f"/api/generated/render_plans/{os.path.basename(render_plan_path)}" if render_plan_path else None,
            'message': f'结构化可编辑 PPTX 已生成（{renderer_display}，不走整页 image-2）'
        }
        response.update(renderer_meta)
        return jsonify(response)
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/upload-page-asset', methods=['POST'])
def upload_page_asset():
    """上传单页素材图，并可选写入该页 fixed_images。"""
    try:
        if 'file' not in request.files:
            return jsonify({'status': 'error', 'message': '没有收到素材文件'}), 400

        f = request.files['file']
        raw = f.read()
        if not raw:
            return jsonify({'status': 'error', 'message': '素材文件为空'}), 400

        original_name = f.filename or 'page_asset.png'
        ext = os.path.splitext(original_name)[1].lower()
        if ext not in {'.png', '.jpg', '.jpeg', '.webp'}:
            content_type = (f.content_type or '').lower()
            if 'jpeg' in content_type or 'jpg' in content_type:
                ext = '.jpg'
            elif 'webp' in content_type:
                ext = '.webp'
            else:
                ext = '.png'

        asset_dir = os.path.join(GENERATED_DIR, 'page_assets')
        os.makedirs(asset_dir, exist_ok=True)
        safe_stem = re.sub(r'[^A-Za-z0-9_.-]+', '_', os.path.splitext(original_name)[0]).strip('._-') or 'page_asset'
        filename = f"{safe_stem}_{int(time.time())}_{uuid.uuid4().hex[:6]}{ext}"
        output_path = os.path.join(asset_dir, filename)

        with open(output_path, 'wb') as out:
            out.write(raw)

        asset_url = f"/api/generated/page_assets/{filename}"

        asset_role_raw = (request.form.get('asset_role') or request.form.get('role') or 'image').strip().lower()
        asset_role = 'background' if asset_role_raw in {'background', 'bg', 'background_image', 'background_asset'} else 'image'
        session_id = request.form.get('session_id', '')
        page_index_raw = request.form.get('page_index', '')
        updated_assets = None
        updated_backgrounds = None
        if session_id and page_index_raw != '':
            try:
                page_index = int(page_index_raw)
            except ValueError:
                page_index = -1

            store = _uploaded_outline_cache.get(session_id) or _load_outline_session(session_id)
            if store and isinstance(store.get('outline'), dict):
                pages = store['outline'].get('pages') or []
                if 0 <= page_index < len(pages) and isinstance(pages[page_index], dict):
                    if asset_role == 'background':
                        backgrounds = pages[page_index].get('background_images') or pages[page_index].get('background_image') or []
                        if not isinstance(backgrounds, list):
                            backgrounds = [backgrounds] if backgrounds else []
                        backgrounds.append(asset_url)
                        pages[page_index]['background_images'] = backgrounds
                        updated_backgrounds = backgrounds
                    else:
                        assets = pages[page_index].get('fixed_images') or []
                        if not isinstance(assets, list):
                            assets = [assets] if assets else []
                        assets.append(asset_url)
                        pages[page_index]['fixed_images'] = assets
                        updated_assets = assets
                    _uploaded_outline_cache[session_id] = store
                    _save_outline_session(session_id, store)

        return jsonify({
            'status': 'success',
            'asset_url': asset_url,
            'asset_role': asset_role,
            'filename': filename,
            'updated_assets': updated_assets,
            'updated_backgrounds': updated_backgrounds,
            'message': '本页背景图已上传' if asset_role == 'background' else '本页素材图已上传'
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/generate-page-asset', methods=['POST'])
def generate_page_asset():
    """Generate one page-level asset image and bind it to fixed_images/background_images."""
    try:
        data = request.get_json(silent=True) or {}
        session_id = data.get('session_id', '')
        page_index_raw = data.get('page_index', 0)
        asset_role = _page_asset_bucket(data.get('asset_role') or data.get('role'))
        style_settings = data.get('style') or {}
        dry_run = bool(data.get('dry_run'))

        try:
            page_index = int(page_index_raw)
        except (TypeError, ValueError):
            return jsonify({'status': 'error', 'message': 'page_index 必须是数字'}), 400

        store = None
        outline = data.get('outline')
        if session_id:
            store = _uploaded_outline_cache.get(session_id) or _load_outline_session(session_id)
            if store:
                outline = store.get('outline')
                _uploaded_outline_cache[session_id] = store

        page = None
        if isinstance(outline, dict):
            pages = outline.get('pages') or []
            if 0 <= page_index < len(pages) and isinstance(pages[page_index], dict):
                page = pages[page_index]
        if page is None and isinstance(data.get('page'), dict):
            page = data['page']

        if not isinstance(page, dict):
            return jsonify({'status': 'error', 'message': '没有找到可生成资产的页面数据'}), 400

        asset_result = _generate_page_asset_for_page(page, asset_role, style_settings, dry_run=dry_run)
        asset_url = asset_result['asset_url']
        filename = asset_result['filename']
        provider = asset_result['provider']
        prompt = asset_result['prompt']

        updated_assets = None
        updated_backgrounds = None
        updated = _append_page_asset(page, asset_url, asset_role)
        if asset_role == 'background':
            updated_backgrounds = updated
        else:
            updated_assets = updated

        if store and isinstance(store.get('outline'), dict):
            _uploaded_outline_cache[session_id] = store
            _save_outline_session(session_id, store)

        return jsonify({
            'status': 'success',
            'asset_role': asset_role,
            'asset_url': asset_url,
            'filename': filename,
            'provider': provider,
            'prompt': prompt,
            'updated_assets': updated_assets,
            'updated_backgrounds': updated_backgrounds,
            'message': '本页背景图已生成并写入资产层' if asset_role == 'background' else '本页素材图已生成并写入资产层',
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/generate-deck-assets', methods=['POST'])
def generate_deck_assets():
    """Generate page-level assets for multiple pages and bind them into editable PPTX asset slots."""
    try:
        data = request.get_json(silent=True) or {}
        session_id = data.get('session_id', '')
        outline = data.get('outline')
        asset_role = _page_asset_bucket(data.get('asset_role') or data.get('role') or 'background')
        style_settings = data.get('style') or {}
        dry_run = bool(data.get('dry_run'))
        missing_only = data.get('missing_only', True)
        max_pages_raw = data.get('max_pages', 12)
        page_indexes = data.get('page_indexes')

        try:
            max_pages = max(1, min(int(max_pages_raw), 80))
        except (TypeError, ValueError):
            max_pages = 12

        store = None
        if session_id:
            store = _uploaded_outline_cache.get(session_id) or _load_outline_session(session_id)
            if store:
                outline = store.get('outline')
                _uploaded_outline_cache[session_id] = store

        if not isinstance(outline, dict) or not isinstance(outline.get('pages'), list):
            return jsonify({'status': 'error', 'message': '请提供 outline 或有效的 session_id'}), 400

        pages = outline.get('pages') or []
        if isinstance(page_indexes, list) and page_indexes:
            candidates = []
            for raw_idx in page_indexes:
                try:
                    idx = int(raw_idx)
                except (TypeError, ValueError):
                    continue
                if 0 <= idx < len(pages):
                    candidates.append(idx)
        else:
            candidates = list(range(len(pages)))

        generated = []
        skipped = []
        for idx in candidates:
            if len(generated) >= max_pages:
                skipped.extend({'page_index': rest, 'reason': 'max_pages_limit'} for rest in candidates[candidates.index(idx):])
                break

            page = pages[idx]
            if not isinstance(page, dict):
                skipped.append({'page_index': idx, 'reason': 'invalid_page'})
                continue

            existing = page.get('background_images') or page.get('background_image') if asset_role == 'background' else page.get('fixed_images')
            if missing_only and _normalize_asset_list(existing):
                skipped.append({'page_index': idx, 'reason': 'already_has_asset'})
                continue

            asset_result = _generate_page_asset_for_page(page, asset_role, style_settings, dry_run=dry_run)
            updated = _append_page_asset(page, asset_result['asset_url'], asset_role)
            generated.append({
                'page_index': idx,
                'index': idx + 1,
                'title': page.get('title') or f'第 {idx + 1} 页',
                'asset_role': asset_role,
                'asset_url': asset_result['asset_url'],
                'filename': asset_result['filename'],
                'provider': asset_result['provider'],
                'prompt': asset_result['prompt'],
                'updated_assets': updated if asset_role == 'image' else None,
                'updated_backgrounds': updated if asset_role == 'background' else None,
            })

        if store and isinstance(store.get('outline'), dict):
            _uploaded_outline_cache[session_id] = store
            _save_outline_session(session_id, store)

        return jsonify({
            'status': 'success',
            'asset_role': asset_role,
            'generated_count': len(generated),
            'skipped_count': len(skipped),
            'generated': generated,
            'skipped': skipped,
            'missing_only': bool(missing_only),
            'dry_run': dry_run,
            'message': f'已生成 {len(generated)} 页{"背景图" if asset_role == "background" else "素材图"}并写入资产层',
        })
    except ValueError as e:
        return jsonify({'status': 'error', 'message': str(e)}), 400
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/download-images-zip', methods=['POST'])
def download_images_zip():
    """将当前生成的图片列表打包为 zip 并返回下载链接"""
    try:
        import re
        import zipfile
        import base64
        import requests

        data = request.get_json(silent=True) or {}
        images = data.get('images', []) or []
        outline = data.get('outline', {}) or {}
        pages = outline.get('pages', []) or []

        if not images:
            return jsonify({'status': 'error', 'message': '没有图片可打包'}), 400

        def safe_name(name: str) -> str:
            name = (name or "").strip()
            # Windows 不允许的字符替换
            name = re.sub(r'[\\\\/:*?"<>|]+', '_', name)
            name = name.replace('\n', ' ').strip()
            return name[:60] if len(name) > 60 else name

        zip_filename = f"images_{int(time.time())}_{uuid.uuid4().hex[:6]}.zip"
        zip_path = os.path.join(BASE_DIR, 'static', zip_filename)
        os.makedirs(os.path.dirname(zip_path), exist_ok=True)

        with zipfile.ZipFile(zip_path, 'w', compression=zipfile.ZIP_DEFLATED) as zf:
            for i, img_url in enumerate(images):
                if not img_url or (isinstance(img_url, str) and img_url.startswith("/api/placeholder/")):
                    continue

                title = ""
                if i < len(pages) and isinstance(pages[i], dict):
                    title = pages[i].get("title") or ""
                base = f"{i+1:02d}"
                if title:
                    base = f"{base}_{safe_name(title)}"
                arcname = f"{base}.png"

                try:
                    # data url
                    if isinstance(img_url, str) and img_url.startswith("data:image"):
                        b64_data = img_url.split(",", 1)[1]
                        raw = base64.b64decode(b64_data)
                        zf.writestr(arcname, raw)
                        continue

                    # /api/generated
                    if isinstance(img_url, str) and img_url.startswith("/api/generated/"):
                        fname = img_url.split("/api/generated/", 1)[1]
                        fpath = os.path.join(GENERATED_DIR, fname)
                        if os.path.exists(fpath):
                            zf.write(fpath, arcname)
                            continue

                    # http(s) or other relative
                    fetch_url = img_url
                    if isinstance(img_url, str) and img_url.startswith("/"):
                        fetch_url = urljoin(request.host_url, img_url.lstrip("/"))
                    if isinstance(fetch_url, str) and fetch_url.startswith("http"):
                        r = requests.get(fetch_url, timeout=60)
                        r.raise_for_status()
                        zf.writestr(arcname, r.content)
                        continue

                    # fallback: local path
                    if isinstance(img_url, str) and os.path.exists(img_url):
                        zf.write(img_url, arcname)
                        continue

                except Exception as e:
                    print(f"[WARN] 打包图片失败（已跳过）: {img_url} err={e}")
                    continue

        return jsonify({'status': 'success', 'download_url': f'/api/download/{zip_filename}', 'message': '图片已打包'})
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500

def _parse_smart_outline(raw_text, skip_filter=False):
    """
    智能解析 Markdown 大纲，支持多种格式：
    1. Sxx｜页面标题 格式（如用户文档中的 S01｜封面）
    2. ## 标题 格式（标准 Markdown）
    3. 格式太乱时调 AI 轻量提取（只做结构化提取，不改写内容）

    返回 pages 列表，每页包含: index, title, type, layout, brief, content
    skip_filter=True 时跳过 _is_likely_page_section 过滤（用于 Word 文档）
    """
    doc_title = ''
    lines = raw_text.split('\n')

    # 提取文档总标题
    for line in lines:
        s = line.strip()
        if s.startswith('# ') and not doc_title:
            doc_title = s[2:].strip()
            break
        elif s.startswith('IMBODY') or s.startswith('【') and len(s) < 60 and not doc_title:
            doc_title = s.strip('【】').strip()
            break

    # 尝试正则匹配多种页面分隔符格式
    page_pattern = re.compile(
        r'^(?:'
        # Sxx｜ 格式
        r'S\d+[｜|:\s]\s*.+'
        r'|'
        # 数字编号格式: 01｜ / 1.1｜
        r'\d+\.?\s*\d*[｜|:\s]\s*.+'
        r'|'
        # ## Markdown标题
        r'##\s+.+'
        r'|'
        # ●一/二/三、中文编号章节（Word文档常见）
        r'[●◆■▪️●]\s*[一二三四五六七八九十]+[、.．]\s*.+'
        r'|'
        # 【X页】标记的行
        r'.*【\d+\s*页】.*'
        r'|'
        # step X / 步骤 X 格式
        r'(?:step|步骤)\s*\d+[\.:\s：].*'
        r'|'
        # 类型一/二/三 格式
        r'类型[一二三四五六七八九十\d]+[：:].*'
        r'|'
        # 独立的中文数字大章节（如 "一、起势期" "四、达人专项策划"）
        r'^[一二三四五六七八九十]+[、.．]\s*[^\n]{4,}'
        r')$',
        re.IGNORECASE | re.MULTILINE
    )
    section_headers = re.compile(r'^\d+\.\d+\s+', re.MULTILINE)

    sections = []
    current_header_line = -1
    current_lines = []

    for i, line in enumerate(lines):
        stripped = line.strip()
        if page_pattern.match(stripped) or (stripped.startswith('## ') and not stripped.startswith('### ')):
            if current_header_line >= 0:
                sections.append((current_header_line, current_lines))
            current_header_line = i
            current_lines = [line]
        else:
            if current_header_line >= 0:
                current_lines.append(line)

    if current_header_line >= 0:
        sections.append((current_header_line, current_lines))

    # 判断是否是真正的 PPT 页面结构（而非说明性章节）
    def _is_likely_page_section(header_text, body_text):
        h = header_text.lower()
        # 明确的页面标识
        page_indicators = ['封面', '目录', '总览', 'recap', '策略', '亮点', '方案', '规划', '执行', '总结', '风险', 'kpi', '目标', '人群', '卖点', '竞品', '学习', '爆款', '公式', '可学', '避坑', '矩阵', '打法', '拍摄清单']
        # 说明性章节标识（非 PPT 页面）
        non_page_indicators = ['设计要求', '视觉风格', '排版规则', '建议页数', '重构目标', '全局', '前言', '附录', '备注']

        for kw in non_page_indicators:
            if kw in h:
                return False

        # 纯数字编号章节（如 "1.4 xxx"、"2.5 xxx"）通常不是独立页面
        if re.match(r'^\d+\.\d+\s+', header_text):
            return False

        # 太短的行（<5字符）且不含页面关键词的，大概率不是页面
        if len(header_text) < 8 and not any(kw in h for kw in page_indicators):
            return False

        # Sxx 编号格式大概率是页面
        if re.match(r'^S\d+[｜|:\s]', header_text):
            return True

        # 有"版式建议"/"视觉元素"字段的大概率是页面
        if '版式建议' in body_text or '视觉元素' in body_text or '视觉关键词' in body_text:
            return True

        # 有"核心结论"或"标题："字段的
        if '核心结论' in body_text[:200] or '页面内容' in body_text[:200]:
            return True

        return any(kw in h for kw in page_indicators)

    pages = []
    first_page_doc_title = {'doc_title': doc_title}

    for idx, (header_idx, sec_lines) in enumerate(sections):
        header = sec_lines[0].strip() if sec_lines else ''
        # 清理标题：去掉 Sxx｜ / ## 前缀
        clean_header = re.sub(r'^(S?\d+[｜|:\s]*|##\s*)', '', header).strip()
        body = '\n'.join(sec_lines[1:]).strip() if len(sec_lines) > 1 else ''

        full_text = '\n'.join(sec_lines).strip()

        if not skip_filter and not _is_likely_page_section(header, full_text):
            continue

        # 从正文提取子字段
        title = clean_header
        core_conclusion = ''
        page_content = ''
        brief = ''
        visual_elements = ''

        # 解析子字段
        field_patterns = {
            '标题': None,
            '副标题': None,
            '核心文案': None,
            '核心结论': None,
            '页面内容': None,
            '版式建议': None,
            '视觉元素': None,
            '视觉关键词': None,
        }

        current_field = None
        current_field_lines = []

        for bline in sec_lines[1:]:
            bst = bline.strip()
            if not bst:
                if current_field and current_field_lines:
                    field_patterns[current_field] = '\n'.join(current_field_lines).strip()
                    current_field_lines = []
                continue

            matched_field = False
            for fk in field_patterns:
                if bst.startswith(fk + '：') or bst.startswith(fk + ':') or bst.startswith(fk + ' '):
                    if current_field and current_field_lines:
                        field_patterns[current_field] = '\n'.join(current_field_lines).strip()
                    current_field = fk
                    val = re.sub(r'^' + fk + r'[：:\s]+', '', bst).strip()
                    current_field_lines = [val] if val else []
                    matched_field = True
                    break

            if not matched_field and current_field:
                current_field_lines.append(bline)

        if current_field and current_field_lines:
            field_patterns[current_field] = '\n'.join(current_field_lines).strip()

        # 用提取到的字段填充
        if field_patterns['标题']:
            title = field_patterns['标题']
        if field_patterns['核心结论']:
            core_conclusion = field_patterns['核心结论']
        if field_patterns['页面内容']:
            page_content = field_patterns['页面内容']
        if field_patterns['版式建议']:
            brief = field_patterns['版式建议']
        if field_patterns['视觉元素'] or field_patterns['视觉关键词']:
            visual_elements = (field_patterns['视觉元素'] or '') + ' | ' + (field_patterns['视觉关键词'] or '')

        # 组装 content：优先用提取到的字段，否则用原始 body
        if page_content or core_conclusion:
            content_parts = []
            if core_conclusion:
                content_parts.append(f"【核心结论】\n{core_conclusion}")
            if page_content:
                content_parts.append(f"【页面内容】\n{page_content}")
            content = '\n\n'.join(content_parts)
        else:
            content = body

        # 推断 layout
        brief_lower = (brief or '').lower()
        content_lower = (content or '').lower()
        if '左图右文' in brief_lower or '左图右文' in content_lower or '左侧' in brief_lower:
            layout = '满版图片-左图右文'
        elif '三列卡片' in brief_lower or '三列' in brief_lower or '矩阵' in brief_lower or '卡片' in brief_lower:
            layout = '模块化-卡片'
        elif '步骤流程' in brief_lower or '步骤' in brief_lower or '流程' in brief_lower:
            layout = '模块化-步骤流程'
        elif '时间轴' in brief_lower or '阶段' in brief_lower or '漏斗' in brief_lower:
            layout = '模块化-时间轴'
        elif '表格' in brief_lower or '对比' in brief_lower or '双栏' in brief_lower:
            layout = '模块化-表格'
        elif '纯视觉' in brief_lower or 'kv' in brief_lower or '大图' in brief_lower:
            layout = '纯视觉页'
        elif '封面' in title.lower() or idx == 0:
            layout = '满版图片-全屏背景'
            brief = brief or f"产品大图 + {doc_title}风格背景"
        else:
            layout = '满版图片-全屏背景'

        if not brief:
            brief = visual_elements or ''

        # 推断 type
        if '封面' in title.lower():
            ptype = 'cover'
        elif '总结' in title.lower() or '总览' in title.lower():
            ptype = 'summary'
        elif '表格' in layout or '矩阵' in layout:
            ptype = 'chart'
        elif '时间轴' in layout or '流程' in layout:
            ptype = 'timeline'
        else:
            ptype = 'content'

        page_obj = {
            'index': len(pages) + 1,
            'title': title,
            'type': ptype,
            'layout': layout,
            'brief': brief,
            'content': content,
        }

        if idx == 0:
            page_obj['doc_title'] = doc_title

        pages.append(page_obj)

    # 如果正则拆分没找到有效页面（格式太乱），调用 AI 轻量提取
    if len(pages) <= 1 and len(raw_text) > 200:
        pages = _ai_extract_pages(raw_text)
        if pages and doc_title:
            pages[0]['doc_title'] = doc_title

    # AI 智能排版分析：根据每页内容确定最优 layout/type/brief
    if pages and len(pages) >= 2:
        pages = _ai_analyze_layouts(pages)

    return pages


def _parse_word_outline(raw_text):
    """
    Word 文档专用拆分策略：按大章节分割 → 内部细粒度拆页

    不使用 _is_likely_page_section 过滤（Word原始文本没有结构化字段）
    分两步走：
    Step1: 按 ●一/二/三/四 或 一/二/三/四 大章节分割
    Step2: 每个章节内部按 【X页】/step/类型/●子项 细拆
    """
    lines = raw_text.split('\n')

    # ===== Step 1: 识别大章节边界 =====
    major_chapter_pattern = re.compile(
        r'^(?:'
        # ●一/二/三/四、格式（最常见）
        r'[●◆■▪️]\s*[一二三四五六七八九十]+[、.．]\s*'
        # 纯 一/二/三/四、 格式（无前缀bullet）
        r'|^[一二三四五六七八九十]+[、.．]\s*'
        r')',
        re.IGNORECASE
    )

    chapters = []  # [(chapter_title, start_line_idx, end_line_idx), ...]
    current_chapter = None
    current_start = 0
    current_lines = []

    for i, line in enumerate(lines):
        stripped = line.strip()
        if major_chapter_pattern.match(stripped) and len(stripped) > 3:
            if current_chapter is not None:
                chapters.append((current_chapter, current_start, i, current_lines))
            # 提取章节标题：去掉 ● 前缀和 一/二/三 前缀
            clean_title = re.sub(r'^[●◆■▪️●\s]*', '', stripped)
            clean_title = re.sub(r'^[一二三四五六七八九十]+[、.．\s]*', '', clean_title).strip()
            current_chapter = clean_title or stripped
            current_start = i
            current_lines = [line]
        else:
            if current_chapter is not None:
                current_lines.append(line)
            elif not current_chapter and stripped:  # 第一个章节之前的内容
                current_lines.append(line)

    if current_chapter is not None:
        chapters.append((current_chapter, current_start, len(lines), current_lines))

    print(f"[INFO] Word文档识别到 {len(chapters)} 个大章节")
    for idx, (title, _, _, _) in enumerate(chapters):
        print(f"[INFO]   章节{idx+1}: {title}")

    # 如果没识别到大章节，把整个文本作为一个章节
    if not chapters:
        chapters.append(('全文', 0, len(lines), lines))

    # ===== Step 2: 每个章节内部细粒度拆页 =====
    all_pages = []

    sub_page_patterns = [
        # 【X页】标记 — 最强信号，必须独立成页
        (re.compile(r'.*【(\d+)\s*页】.*'), 10),
        # step X / 步骤 X
        (re.compile(r'(?:step|步骤)\s*(\d+)[\.:\s：]', re.IGNORECASE), 8),
        # 类型一/二/三
        (re.compile(r'类型([一二三四五六七八九十\d]+)[：:]', re.IGNORECASE), 7),
        # ● 子项（非大章节的 ●）
        (re.compile(r'^[●◆■▪️]\s+(?![一二三四五六七八九十]+[、.．])', re.IGNORECASE), 6),
        # 编号列表 1) 2) 3)
        (re.compile(r'^\d+\)[\s、]', re.IGNORECASE), 5),
    ]

    for ch_idx, (ch_title, ch_start, ch_end, ch_lines) in enumerate(chapters):
        ch_text = '\n'.join(ch_lines).strip()

        # 在章节内查找所有子页面分割点
        split_points = []  # [(priority, line_idx, title_text)]

        for li, line in enumerate(ch_lines):
            st = line.strip()
            for pattern, priority in sub_page_patterns:
                m = pattern.search(st)
                if m:
                    try:
                        page_title = m.group(1) if m.lastindex and m.lastindex >= 1 else st[:50]
                    except:
                        page_title = st[:50]
                    split_points.append((priority, li + ch_start, page_title))
                    break

        # 按行号排序
        split_points.sort(key=lambda x: x[1])

        # 如果没有找到子分割点，整个章节作为一页
        if not split_points:
            all_pages.append({
                'index': len(all_pages) + 1,
                'title': ch_title,
                'type': 'content',
                'layout': _guess_layout(ch_title, ch_text),
                'brief': '',
                'content': ch_text[:3000],
            })
            continue

        # 按分割点切分页面
        for sp_idx, (sp_priority, sp_line_idx, sp_title) in enumerate(split_points):
            # 这页的内容范围：从当前分割点到下一个分割点（或章节结束）
            content_start = sp_line_idx
            content_end = split_points[sp_idx + 1][1] if sp_idx + 1 < len(split_points) else ch_end

            page_lines = [l for l in ch_lines if content_start <= (ch_lines.index(l) if l in ch_lines else 0) < content_end]

            # 更精确地取内容
            page_content_raw = '\n'.join(lines[content_start:content_end]).strip() if content_start < len(lines) else ''

            # 清理标题
            clean_sp_title = re.sub(r'^[●◆■▪️●\s]*(?:step|步骤)?\s*\d*[\.:\s：]?\s*', '', sp_title).strip()
            clean_sp_title = re.sub(r'^类型[一二三四五六七八九十\d]+[：:\s]*', '', clean_sp_title).strip()
            clean_sp_title = re.sub(r'.*【\d+\s*页】\s*', '', clean_sp_title).strip()
            final_title = clean_sp_title or sp_title or f'{ch_title}-{sp_idx+1}'

            all_pages.append({
                'index': len(all_pages) + 1,
                'title': final_title,
                'type': _guess_type(final_title),
                'layout': _guess_layout(final_title, page_content_raw),
                'brief': '',
                'content': page_content_raw[:3000],
            })

        print(f"[INFO]   「{ch_title}」拆出 {len([p for p in split_points])} 个子页面")

    print(f"[INFO] Word文档总计拆分: {len(all_pages)} 页")
    return all_pages


def _guess_card_count(content=''):
    """根据内容智能判断卡片列数"""
    if not content:
        return 3
    c = str(content)

    digit_dot = re.findall(r'^\s*[1234567890]\.\s', c, re.MULTILINE)
    if len(digit_dot) >= 2:
        return min(max(len(digit_dot), 3), 8)

    chinese_nums = re.findall(r'[一二三四五六七八九十]、', c)
    if len(chinese_nums) >= 2:
        return min(len(chinese_nums) + 1, 8)

    digit_nums = re.findall(r'[123456789]、', c)
    if len(digit_nums) >= 2:
        return min(len(digit_nums) + 1, 8)

    circle_nums = re.findall(r'[①②③④⑤⑥⑦⑧⑨⑩]', c)
    if len(circle_nums) >= 2:
        return min(len(circle_nums), 8)

    bullet_nums = re.findall(r'[•●○◉◆■□▼△]\s', c)
    if len(bullet_nums) >= 2:
        return min(len(bullet_nums), 8)

    separators = c.count('、') + c.count('；')
    if separators >= 3:
        return min(max(3, separators), 8)

    lines = [l.strip() for l in re.split(r'[\n\r]+', c) if l.strip()]
    if len(lines) >= 3:
        return min(len(lines), 8)

    return 3

_TEXT_ONLY_LAYOUTS = ['文字-金句引言', '文字-数字大屏', '文字-问答式', '文字-双栏对比', '文字-要点列表', '文字-引用来源']

def _is_text_only_layout(layout=''):
    if not layout:
        return False
    return any(t in layout for t in _TEXT_ONLY_LAYOUTS)

def _guess_layout(title, content=''):
    """根据标题和内容猜测布局（13种）"""
    t = (title + ' ' + content[:300]).lower()

    # ===== 强图布局 =====
    if '左图右文' in t or '左侧' in t or '左右' in t:
        return '满版图片-左图右文'
    elif '全屏' in t or '封面' in t or '总结' in t or '总览' in t:
        return '满版图片-全屏背景'
    elif '纯视觉' in t or ' kv' in t or '大图' in t:
        return '纯视觉页'

    # ===== 少图布局 =====
    elif any(k in t for k in ['三列', '矩阵', '卡片']):
        return '模块化-卡片'
    elif any(k in t for k in ['步骤', '流程', 'step', '执行', '拆解']):
        return '模块化-步骤流程'
    elif any(k in t for k in ['时间轴', '阶段', '节奏', 'timeline']):
        return '模块化-时间轴'
    elif any(k in t for k in ['表格', '对比', '竞品']):
        return '模块化-表格'

    # ===== 纯文字布局（新增加）=====
    # 引言/金句/观点页：大字引言 + 小字出处
    elif any(k in t for k in ['引言', '金句', '核心观点', '洞察', '发现', 'quote', '观点']):
        return '文字-金句引言'

    # 数字/数据展示：大数字 + 说明文字
    elif any(k in t for k in ['数据', '增长', '趋势', '指标', 'kpi', '占比', '数字', '统计', '亿', '万', '%']):
        return '文字-数字大屏'

    # 问答/对话页：问题+回答格式
    elif any(k in t for k in ['qa', '问答', 'faq', '常见问题', '问：', '答：']):
        return '文字-问答式'

    # 双栏文字对比/左右排版
    elif any(k in t for k in ['对比', '双栏', 'a与b', '优劣', '异同']):
        return '文字-双栏对比'

    # 列表/要点：带编号的要点罗列
    elif any(k in t for k in ['要点', '核心要点', '注意事项', '须知', 'checklist', '清单']):
        return '文字-要点列表'

    # 引用/来源/注释页
    elif any(k in t for k in ['来源', '参考', '数据来源', '出处', '附注', '备注']):
        return '文字-引用来源'

    # 默认返回少图布局
    return '模块化-卡片'


def _guess_type(title):
    """根据标题猜测页面类型"""
    tl = title.lower()
    if '封面' in tl:
        return 'cover'
    elif any(k in tl for k in ['总结', '总览', 'recap']):
        return 'summary'
    elif any(k in tl for k in ['时间轴', '阶段', '节奏', '步骤', '流程']):
        return 'timeline'
    elif any(k in tl for k in ['表格', '对比', '数据', '分析', '竞品']):
        return 'chart'
    return 'content'


def _page_needs_product_image(page_title, page_content, page_type, page_layout):
    """
    智能判断当前页面是否需要产品参考图。

    需要产品图的页面：
    - 产品展示/介绍/特点/功能页
    - 使用场景/应用场景页
    - 封面、KV纯视觉页
    - 内容明确提到产品名称或品牌

    不需要产品图的页面：
    - 时间轴、流程步骤、表格等抽象排版页
    - 数据分析/趋势/调研等纯信息页
    - 策略/定位/洞察等概念性页面
    - 传播规划/TVC脚本等文字为主页面
    """
    title = (page_title or '').lower()
    content = (page_content or '').lower()
    layout = (page_layout or '').lower()
    ptype = (page_type or '').lower()

    combined = f"{title} {content}"

    # ===== 肯定需要产品图的场景 =====
    need_keywords = [
        '产品展示', '产品介绍', '产品特点', '产品功能', '产品外观',
        '产品特写', '产品细节', '产品图', '产品渲染', '产品效果图',
        '使用场景', '应用场景', '使用方式', '使用演示',
        '产品优势', '核心卖点', '产品卖点', '技术参数',
        'imbody', 'pm up', 'unitree', '力量魔方',
        '产品外观', '外形', '造型', '尺寸', '重量',
        '开箱', '实物', '真机',
    ]

    if any(kw in combined for kw in need_keywords):
        return True

    # 封面页、纯视觉页、KV页 — 通常需要产品
    if ptype in ('cover', 'visual') or '封面' in title or 'kv' in layout or '纯视觉' in layout:
        return True
    # 但如果封面是纯品牌概念（无具体产品词），可能不需要
    if ptype == 'cover' or '封面' in title:
        cover_need = any(kw in combined for kw in [
            '产品', '发布', '亮相', '展示', 'imbody', 'pm', 'unitree'
        ])
        if not cover_need:
            return False
        return True

    # ===== 肯定不需要产品图的场景 =====
    no_need_keywords = [
        '时间轴', '时间线', '阶段规划', '排期',
        '数据', '分析', '趋势', '增长', '调研', '报告',
        '策略', '定位', '洞察', '发现', '现象',
        'tvc', '脚本', '镜头', '画面', '旁白', '分镜',
        '流程', '步骤', '路径', '漏斗',
        '表格', '总览', '概览', '汇总',
        '预算', '费用', '成本', 'kpi',
        '团队', '分工', '职责',
        '风险', '预案', '备选',
    ]

    no_need_layouts = ['时间轴', '步骤流程', '表格', '模块化-时间轴', '模块化-步骤', '模块化-表格']

    if any(kw in combined for kw in no_need_keywords):
        return False
    if any(nl in layout for nl in no_need_layouts):
        return False

    # 模块化卡片类 — 看内容是否有产品相关
    if '模块化' in layout and '三列' in layout:
        card_needs = any(kw in combined for kw in need_keywords)
        return card_needs

    # 默认：有产品图时保守策略，不强制每页都带
    # 只有内容明显涉及产品才带
    product_hint = any(kw in combined for kw in [
        '产品', '设备', '硬件', '机器', '器械', '装置',
        'imbody', 'pm', 'unitree', '力量魔方',
        '展示', '呈现', '效果', '外观',
    ])
    return product_hint


def _ai_analyze_layouts(pages):
    """
    调用 AI 根据每页内容智能分析排版方式（layout/type/brief）
    不改写内容，只做排版建议
    """
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    try:
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)

        chat_api_url = config.get('chat_api_base_url') or config.get('api_base_url')
        chat_api_key = config.get('chat_api_key') or config.get('api_key')

        if not chat_api_url or not chat_api_key:
            print("[WARN] 无 API 配置，跳过 AI 排版分析")
            return pages

        api = APIClient(chat_api_url, chat_api_key)

        pages_summary = []
        for i, p in enumerate(pages):
            content_preview = (p.get('content', '') or '')[:300]
            brief_preview = (p.get('brief', '') or '')[:100]
            pages_summary.append(f"第{i+1}页 | 标题: {p['title']} | 当前布局: {p.get('layout','')} | 内容摘要: {content_preview}")

        prompt = f"""你是一个专业的 PPT 排版设计师。以下是一份 PPT 的逐页内容信息。

【页面列表】
{chr(10).join(pages_summary)}

【你的任务】为每一页选择最合适的排版方式。根据页面内容的类型（封面、数据表格、对比分析、流程步骤、时间轴、图文混排等）来决定。

【可选的 layout（版式）】
- 满版图片-全屏背景：适合封面、总结页、视觉冲击力强的页面
- 满版图片-左图右文：适合有产品图+文字说明的页面
- 模块化-卡片：适合并列展示多个要点/特性/方向（2-5列，AI自动判断）
- 模块化-步骤流程：适合操作流程、策略路径、执行步骤
- 模块化-时间轴：适合时间线、阶段规划、发展历程
- 模块化-表格：适合数据对比、竞品分析、要做vs不做
- 纯视觉页：适合纯图片展示、KV 主视觉

【可选的 type】cover / content / visual / chart / summary / timeline

【输出格式 — 严格 JSON，不要其他文字】
```json
{{
  "layouts": [
    {{"page_index": 1, "title": "保持原标题", "type": "cover", "layout": "满版图片-全屏背景", "brief": "一句话描述这页该用什么画面风格"}},
    {{"page_index": 2, "title": "保持原标题", "type": "chart", "layout": "模块化-表格", "brief": "..."}}
  ]
}}
```

注意：
1. title 必须保持原样，不能修改
2. layout 要根据实际内容合理分配，不要所有页都用同一种
3. 有表格/对比/矩阵数据的用「模块化-表格」
4. 有流程/步骤的用「模块化-步骤流程」
5. 有阶段/时间的用「模块化-时间轴」
6. 有多个并列要点的用「模块化-卡片」（AI会自动判断列数）
7. 封面和总结用「满版图片-全屏背景」或「纯视觉页」"""

        result = api.chat(
            model=config.get('chat_model', 'gpt-4o'),
            messages=[{'role': 'user', 'content': prompt}],
            max_tokens=8000,
            timeout=180
        )

        content = result.get('content', '')
        json_match = re.search(r'```json\s*([\s\S]*?)\s*```', content)
        json_str = json_match.group(1) if json_match else content
        data = json.loads(json_str)

        layouts = data.get('layouts', [])
        for item in layouts:
            idx = item.get('page_index', -1)
            if 1 <= idx <= len(pages):
                page = pages[idx - 1]
                if item.get('layout'):
                    page['layout'] = item['layout']
                if item.get('type'):
                    page['type'] = item['type']
                if item.get('brief'):
                    page['brief'] = item['brief']

        print(f"[INFO] AI 排版分析完成: {len(layouts)} 页已优化")

    except Exception as e:
        print(f"[WARN] AI 排版分析失败: {e}")

    return pages


def _ai_extract_pages(raw_text):
    """当本地解析无法识别页面时，调 AI 做「细粒度拆分」——把长文档拆成多页PPT"""
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    try:
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)

        chat_api_url = config.get('chat_api_base_url') or config.get('api_base_url')
        chat_api_key = config.get('chat_api_key') or config.get('api_key')

        if not chat_api_url or not chat_api_key:
            return []

        api = APIClient(chat_api_url, chat_api_key)

        # 先估算文档长度，如果太长需要分批处理
        MAX_CHARS_PER_BATCH = 15000
        total_chars = len(raw_text)

        if total_chars > MAX_CHARS_PER_BATCH:
            # 长文档分批处理
            return _ai_extract_pages_batched(api, config, raw_text, MAX_CHARS_PER_BATCH)

        return _ai_extract_pages_single(api, config, raw_text)

    except Exception as e:
        print(f"[WARN] AI 提取页面失败: {e}")

    return []


def _ai_extract_pages_single(api, config, raw_text):
    """单次AI调用拆分页面（适用于较短文档）"""

    # 预扫描【X页】标记，提取强制页数要求
    _page_markers = re.findall(r'【(\d+)\s*页】', raw_text)
    _force_page_count = sum(int(x) for x in _page_markers) if _page_markers else 0
    _has_page_markers = len(_page_markers) > 0

    # 提高截断上限，确保【X页】标记能被读到
    _max_chars = 50000 if _has_page_markers else 30000
    _truncated_text = raw_text[:_max_chars]
    if len(raw_text) > _max_chars:
        _truncated_text += f"\n\n[注意：原文档较长（共{len(raw_text)}字符），以上为前{_max_chars}字符内容]"
        if _has_page_markers:
            _truncated_text += f"\n[重要：原文档中发现 {len(_page_markers)} 个【X页】标记，总要求 {_force_page_count} 页，请务必满足此页数要求！]"

    prompt = f"""你是一个 PPT 拆分专家。请将以下文档内容拆分成**独立的PPT页面**。

【原文档】
{_truncated_text}

【拆分规则 — 必须严格遵守】
1. **细粒度拆分**：每个独立的话题、步骤、案例、分析都要拆成单独页面
   - 例："step1. xxx / step2. yyy / step3. zzz" → 拆成3页
   - 例："类型一：xxx / 类型二：yyy" → 拆成2页
   - 例："1. 搜索 / 2. 品牌声量 / 3. 兴趣人群" → 拆成3页

2. **🔴【X页】标记 — 最高优先级**：
   - 原文档中明确标注了【X页】的章节，**必须严格拆分成 X 个独立页面**
   - 这是用户/策划明确指定的页数要求，不可违反
   - 例如：【5页】高潜场景挖掘 → 必须拆成至少5个页面（概述+养宠+烹饪+睡眠+育儿+清洁）
   - 例如：【3页】竞品调研 → 必须拆成至少3个页面（华为+海尔+小米各一页）

3. **识别页面标记**：
   - 带编号的子项（●、-、1. 2. 3.）每个都是潜在独立页面

4. **内容深度**：每个页面只讲一个具体话题，不要合并多个子项

5. **⚠️ 内容边界 — 绝对禁止违反**：
   - content 字段的所有文字**必须来自上方【原文档】原文**
   - 只能对原文进行**拆分、截取、拼接**，**禁止自行创作、改写、扩写或添加原文没有的内容**
   - 禁止编造数据、数字、案例、引用等原文中没有的信息
   - 禁止用"例如"、"比如"后接自己编造的内容
   - 如果原文信息不足，宁可简短也不要编造

6. **页面类型判断**：
   - cover: 封面/标题页
   - content: 普通内容页
   - visual: 纯视觉/KV页
   - chart: 数据/对比/分析页
   - timeline: 时间轴/流程页
   - summary: 总结/结论页

7. **布局建议**：根据内容选择合适的layout
   - 满版图片-全屏背景: 封面、情感页
   - 满版图片-左图右文: 对比、案例
   - 模块化-卡片: 分类、卖点（AI自动判断列数）
   - 模块化-步骤流程: 步骤、流程
   - 模块化-时间轴: 阶段、规划
   - 模块化-表格: 对比、数据
   - 纯视觉页: KV、大图

【输出格式 — 严格JSON】
```json
{{
  "pages": [
    {{
      "title": "页面标题（具体子话题）",
      "type": "cover/content/visual/chart/summary/timeline",
      "layout": "满版图片-全屏背景/满版图片-左图右文/模块化-卡片/模块化-步骤流程/模块化-时间轴/模块化-表格/纯视觉页",
      "brief": "这页要讲什么（一句话）",
      "content": "页面详细内容（原文相关部分）"
    }}
  ],
  "total_pages": 数字
}}
```

{"🔴🔴🔴 最高优先级：【X页】标记强制要求 " + str(_force_page_count) + " 页，请务必满足！" if _has_page_markers else ""}
【重要】请尽可能拆细，确保每个独立观点/步骤/案例都有独立页面。不要合并多个子项到一页。
只输出JSON，不要其他文字。"""

    result = api.chat(
        model=config.get('chat_model', 'gpt-4o'),
        messages=[{'role': 'user', 'content': prompt}],
        max_tokens=16000,
        timeout=300
    )

    content = result.get('content', '')
    return _parse_ai_pages_response(content)


def _ai_extract_pages_batched(api, config, raw_text, batch_size):
    """分批处理长文档"""
    # 按段落分批
    paragraphs = raw_text.split('\n\n')
    batches = []
    current_batch = []
    current_len = 0

    for para in paragraphs:
        if current_len + len(para) > batch_size and current_batch:
            batches.append('\n\n'.join(current_batch))
            current_batch = [para]
            current_len = len(para)
        else:
            current_batch.append(para)
            current_len += len(para) + 2

    if current_batch:
        batches.append('\n\n'.join(current_batch))

    print(f"[INFO] 长文档分批处理: {len(batches)} 批")

    all_pages = []
    for i, batch in enumerate(batches):
        print(f"[INFO] 处理第 {i+1}/{len(batches)} 批...")
        pages = _ai_extract_pages_single(api, config, batch)
        # 调整页码
        for p in pages:
            p['index'] = len(all_pages) + 1
        all_pages.extend(pages)

    return all_pages


def _parse_ai_pages_response(content):
    """解析AI返回的pages JSON"""
    import json
    import re

    try:
        json_match = re.search(r'```json\s*([\s\S]*?)\s*```', content)
        if json_match:
            json_str = json_match.group(1)
        else:
            json_match = re.search(r'\{[\s\S]*"pages"[\s\S]*\}', content)
            json_str = json_match.group(0) if json_match else content

        data = json.loads(json_str)
        extracted = data.get('pages', [])

        if isinstance(extracted, list):
            out = []
            for i, p in enumerate(extracted):
                out.append({
                    'index': i + 1,
                    'title': p.get('title', f'第{i+1}页'),
                    'type': p.get('type', 'content'),
                    'layout': p.get('layout', '满版图片-全屏背景'),
                    'brief': p.get('brief', ''),
                    'content': p.get('content', ''),
                })
            print(f"[INFO] AI 拆分出 {len(out)} 页")

            out = _post_process_page_markers(out)

            return out
    except Exception as e:
        print(f"[WARN] 解析AI返回失败: {e}")

    return []


def _build_scene_description(page_title='', page_content=''):
    _t = (page_title or '').lower()
    _c = (page_content or '').lower()
    _all = _t + ' ' + _c

    _scene = ""
    _mood = ""
    _elements = []

    if any(k in _all for k in ['家', '回家', '温馨', '舒适', '客厅', '沙发', '卧室', '厨房', '暖光', '温暖', '归属感', '想家', '港湾']):
        _scene = "cozy warm home interior — living room with soft natural light streaming through curtains, comfortable sofa, personal touches like family photos or plants"
        _mood = "warm, intimate, peaceful, emotionally comforting"
        _elements = ["soft warm lighting", "home furnishings", "natural textures", "comfortable atmosphere"]

    elif any(k in _all for k in ['夏天', '夏日', '炎热', '清凉', '空调', '风扇', '冰饮', '西瓜', '蝉鸣', '树荫', '泳池', '海滩', '度假']):
        _scene = "bright summer scene matching the content theme — could be outdoor shade, cool interior with AC breeze, or refreshing summer activity"
        _mood = "fresh, vibrant, relaxed, summery energy"
        _elements = ["natural daylight", "cool color palette", "summer elements"]

    elif any(k in _all for k in ['产品', '功能', '使用', '操作', '界面', 'app', '手机', '体验', '流程', '步骤', '玩法']):
        _scene = "clean product or app showcase scene — minimal background that highlights the subject matter without distraction"
        _mood = "modern, clean, professional, trustworthy"
        _elements = ["minimal composition", "clean background", "product/app focus", "professional lighting"]

    elif any(k in _all for k in ['数据', '报告', '分析', '趋势', '增长', '图表', '统计', '洞察', '调研', '用户量', '转化']):
        _scene = "abstract data visualization aesthetic — subtle charts, graphs, or data-inspired geometric patterns as background texture"
        _mood = "analytical, professional, insightful, data-driven"
        _elements = ["subtle data patterns", "geometric shapes", "clean lines", "professional tone"]

    elif any(k in _all for k in ['品牌', '声量', '传播', '营销', '推广', '投放', '曝光', '话题', '热度', '社交', '小红书', '抖音', '微博']):
        _scene = "dynamic social media / marketing concept scene — digital connectivity, content creation, viral spread visualization"
        _mood = "energetic, connected, trending, modern digital lifestyle"
        _elements = ["digital screens", "social connection", "content creation", "brand presence"]

    elif any(k in _all for k in ['用户', '人群', '画像', '受众', '目标', '客群', '细分', '特征', '需求', '痛点']):
        _scene = "diverse people in authentic everyday moments that reflect the target audience described in the content"
        _mood = "relatable, authentic, diverse, real-life scenarios"
        _elements = ["authentic people", "everyday settings", "natural expressions", "diverse representation"]

    elif any(k in _all for k in ['时间轴', '时间线', '阶段', '历程', '发展', '演进', ' roadmap', '规划', '里程碑']):
        _scene = "abstract timeline or journey concept — path, road, or progression visual metaphor"
        _mood = "forward-moving, structured, progressive, purposeful"
        _elements = ["path/road metaphor", "directional flow", "milestone markers", "clean composition"]

    elif any(k in _all for k in ['对比', '差异', '优劣势', 'vs', '比较', '选择', '决策']):
        _scene = "split-composition or comparison visual — two sides or options shown in balanced visual arrangement"
        _mood = "clear, balanced, analytical, decision-oriented"
        _elements = ["split composition", "visual contrast", "balanced layout", "comparison metaphor"]

    elif any(k in _all for k in ['情感', '情绪', '心情', '感受', '共鸣', '触动', '故事', '回忆', '瞬间']):
        _scene = "emotionally evocative scene that captures the feeling described in the content — focus on atmosphere and mood over specific objects"
        _mood = "emotional, atmospheric, cinematic in feeling (not style), deeply resonant"
        _elements = ["emotional atmosphere", "meaningful details", "cinematic mood (natural, not artificial)", "storytelling composition"]

    else:
        _scene = f"thematic photography scene that visually represents: '{(page_title or '')[:80]}' — create a scene that matches this topic's emotional tone and subject matter"
        _mood = "professional, on-brand, contextually appropriate"
        _elements = ["topic-relevant imagery", "appropriate mood", "professional quality"]

    return f"""SCENE DESCRIPTION (based on actual page content):
- REQUIRED SCENE: {_scene}
- MOOD/TONE: {_mood}
- VISUAL ELEMENTS TO INCLUDE: {', '.join(_elements)}
- CRITICAL: The background image MUST match the page's actual content topic.
  If the content talks about 'home/家', show a home interior.
  If about 'summer/夏天', show a summer scene.
  If about 'data/数据', use abstract data aesthetics.
  NEVER use generic stock scenes that don't match the content."""


def _post_process_page_markers(pages):
    """后处理：检测每页content中的【X页】标记，强制拆分成多个独立页面"""
    import re

    if not pages:
        return pages

    result = []
    page_index_offset = 0

    for page in pages:
        content = page.get('content', '')
        if not content:
            result.append(page)
            continue

        markers = re.findall(r'【(\d+)\s*页】', content)
        if not markers:
            result.append(page)
            continue

        total_requested = sum(int(m) for m in markers)
        if total_requested <= 1:
            result.append(page)
            continue

        print(f"[INFO] 🔍 第{page['index']}页「{page['title']}」发现 {len(markers)} 个【X页】标记，总要求 {total_requested} 页 → 拆分中...")

        split_parts = re.split(r'(?=【\d+\s*页】)', content)

        sub_pages = []
        current_base_idx = len(result) + len(sub_pages) + 1

        for part_idx, part in enumerate(split_parts):
            part = part.strip()
            if not part or len(part) < 10:
                continue

            part_marker = re.match(r'【(\d+)\s*页】\s*', part)
            part_content = re.sub(r'^【\d+\s*页】\s*', '', part).strip()

            if not part_content:
                continue

            first_line = part_content.split('\n')[0].strip()
            first_line = re.sub(r'^[●◆■▪️\-\*\d]+[\.、:\s]*', '', first_line).strip()[:40]

            sub_title = first_line if first_line and len(first_line) > 3 else f"{page['title']}-Part{part_idx+1}"

            sub_layout = page.get('layout', '模块化-卡片')

            guess = _guess_layout(sub_title, part_content[:300])
            if guess != '模块化-三列卡片':
                sub_layout = guess

            sub_pages.append({
                'index': current_base_idx + len(sub_pages),
                'title': sub_title,
                'type': page.get('type', 'content'),
                'layout': sub_layout,
                'brief': f'{page.get("brief", "")} - 拆分{part_idx+1}',
                'content': part_content[:2000] if len(part_content) > 2000 else part_content
            })

        if len(sub_pages) >= 2:
            print(f"[INFO] ✅ 第{page['index']}页已拆分为 {len(sub_pages)} 个独立页面")
            result.extend(sub_pages)
            page_index_offset += len(sub_pages) - 1
        else:
            print(f"[INFO] ⚠️ 第{page['index']}页拆分结果只有{len(sub_pages)}页，保留原页面")
            result.append(page)

    if len(result) != len(pages):
        for i, p in enumerate(result):
            p['index'] = i + 1
        print(f"[INFO] 📊 【X页】标记后处理完成: {len(pages)} 页 → {len(result)} 页")

    return result


import os

_OUTLINE_STORAGE_DIR = os.path.join(BASE_DIR, '..', 'data', 'outline_sessions')
os.makedirs(_OUTLINE_STORAGE_DIR, exist_ok=True)

def _get_outline_path(session_id):
    return os.path.join(_OUTLINE_STORAGE_DIR, f'{session_id}.json')

def _save_outline_session(session_id, data):
    path = _get_outline_path(session_id)
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def _load_outline_session(session_id):
    path = _get_outline_path(session_id)
    if not os.path.exists(path):
        return None
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except:
        return None

def _list_outline_sessions():
    """列出所有历史大纲会话"""
    sessions = []
    if not os.path.exists(_OUTLINE_STORAGE_DIR):
        return sessions
    for fname in os.listdir(_OUTLINE_STORAGE_DIR):
        if fname.endswith('.json'):
            fpath = os.path.join(_OUTLINE_STORAGE_DIR, fname)
            try:
                stat = os.stat(fpath)
                session_id = fname[:-5]  # 去掉 .json
                data = json.load(open(fpath, 'r', encoding='utf-8'))
                outline = data.get('outline', {})
                pages = outline.get('pages', [])
                images_count = len(data.get('images', []))
                sessions.append({
                    'session_id': session_id,
                    'title': outline.get('title', '未命名大纲'),
                    'page_count': len(pages),
                    'images_count': images_count,
                    'created_at': data.get('created_at', stat.st_mtime),
                    'updated_at': stat.st_mtime,
                    'file_size': stat.st_size
                })
            except:
                pass
    # 按更新时间倒序
    sessions.sort(key=lambda x: x['updated_at'], reverse=True)
    return sessions

# 兼容旧代码的内存缓存（启动时从磁盘加载）
_uploaded_outline_cache = {}

@app.route('/api/upload-outline', methods=['POST'])
def upload_outline():
    """上传大纲（支持 JSON / Markdown），Markdown 会调 AI 解析为结构化大纲"""
    try:
        outline = None
        raw_text = None
        is_markdown = False

        content_type = request.headers.get('Content-Type', '')

        if 'multipart/form-data' in content_type:
            raw_data = request.get_data()
            boundary = None
            for part in content_type.split(';'):
                part = part.strip()
                if part.startswith('boundary='):
                    boundary = part.split('=', 1)[1].strip('"')
                    break
            if not boundary:
                return jsonify({'status': 'error', 'message': '无效的请求格式，缺少 boundary'}), 400

            boundary_bytes = boundary.encode('utf-8')
            delimiter = b'--' + boundary_bytes
            parts = raw_data.split(delimiter)

            file_data = None
            filename = ''
            for part in parts:
                if b'Content-Disposition' not in part:
                    continue
                header_end = part.find(b'\r\n\r\n')
                if header_end == -1:
                    continue
                body = part[header_end + 4:]
                while body.endswith(b'\r') or body.endswith(b'\n') or body.endswith(b'-'):
                    body = body[:-1]
                if body:
                    file_data = body
                    header_part = part[:header_end].decode('utf-8', errors='ignore')
                    fn_match = re.search(r'filename=["\']?([^"\';\s]+)["\']?', header_part)
                    if fn_match:
                        filename = fn_match.group(1)
                    break

            if not file_data:
                return jsonify({'status': 'error', 'message': '没有找到上传的文件内容'}), 400

            raw_text = file_data.decode('utf-8', errors='ignore')
            is_markdown = filename.lower().endswith('.md') or filename.lower().endswith('.markdown')
            is_docx = filename.lower().endswith('.docx')

            if is_docx:
                import io
                from docx import Document
                doc = Document(io.BytesIO(file_data))
                doc_text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        doc_text_parts.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_cells = [cell.text.strip() for cell in row.cells]
                        doc_text_parts.append(' | '.join(row_cells))
                raw_text = '\n\n'.join(doc_text_parts)
                is_markdown = True
                print(f"[INFO] Word 文档已提取文本: {len(raw_text)} 字符, {len(doc.paragraphs)} 段落, {len(doc.tables)} 表格")

            if is_markdown:
                pass
            else:
                try:
                    outline = json.loads(file_data)
                except json.JSONDecodeError:
                    pass
        else:
            data = request.get_json(silent=True) or {}
            if not data:
                return jsonify({'status': 'error', 'message': '请求体不是有效的 JSON'}), 400

            raw_text_str = data.get('text', '') or data.get('content', '')
            fmt = (data.get('format', '') or '').lower()

            if raw_text_str and fmt in ('md', 'markdown'):
                raw_text = raw_text_str
                is_markdown = True
            elif 'pages' in data:
                outline = data
            elif raw_text_str:
                raw_text = raw_text_str
                is_markdown = True
            else:
                outline = data

        if is_markdown and raw_text and not outline:
            # Word 文档跳过 _is_likely_page_section 过滤（避免前半段被误过滤）
            # Markdown 文件保留过滤（防止非页面内容混入）
            pages = _parse_smart_outline(raw_text, skip_filter=is_docx)

            # 尾部补漏：检查正则是否处理完了全部文本
            # 不限制页数！只要尾部还有未处理的文本就继续拆
            if pages and len(raw_text) > 5000:
                # 找到最后一个页面的标题在原文中的位置
                last_title = pages[-1].get('title', '')

                # 用最后一个页面标题定位，取之后的所有文本
                # 多种方式尝试定位
                last_pos = -1

                # 方式1：直接搜索最后一个标题
                if last_title and len(last_title) > 3:
                    last_pos = raw_text.rfind(last_title[:30])

                # 方式2：如果方式1失败，用最后一个页面的内容末尾来定位
                if last_pos <= 0 and pages[-1].get('content'):
                    lc = pages[-1]['content']
                    search_str = lc[-150:] if len(lc) > 150 else lc
                    last_pos = raw_text.rfind(search_str)
                    if last_pos > 0:
                        last_pos += len(search_str)

                # 如果找到了位置，且后面还有大量文本，说明有尾巴没处理
                if last_pos > 0 and last_pos < len(raw_text) - 800:
                    remaining_text = raw_text[last_pos:].strip()

                    if len(remaining_text) > 500:
                        print(f"[INFO] ⚠️ 正则拆分了 {len(pages)} 页，但尾部还有 {len(remaining_text)} 字符未处理")
                        print(f"[INFO]    最后页面标题: {last_title}")
                        print(f"[INFO]    最后位置: 第{last_pos}字符 / 总{len(raw_text)}字符")
                        print(f"[INFO]    尾部预览: {remaining_text[:100]}...")

                        # 对尾巴做AI拆分
                        extra_pages = _ai_extract_pages(remaining_text)

                        if extra_pages:
                            print(f"[INFO] ✅ 尾部补漏成功: 额外获得 {len(extra_pages)} 页")
                            for ep in extra_pages:
                                ep['index'] = len(pages) + 1
                                pages.append(ep)
                        else:
                            print(f"[WARN] 尾部AI拆分未返回结果，尝试将剩余内容作为最后一页追加...")
                            # 兜底：把剩余文本作为最后一个补充页
                            pages.append({
                                'index': len(pages) + 1,
                                'title': '（后续内容）',
                                'type': 'content',
                                'layout': '模块化-表格',
                                'brief': '',
                                'content': remaining_text[:3000],
                            })

            if not pages:
                outline = {'title': 'PPT大纲', 'pages': []}
            else:
                pages = _post_process_page_markers(pages)
                outline = {'title': pages[0].get('doc_title', '') or 'PPT大纲', 'pages': pages}
                print(f"[INFO] 大纲拆分完成: 共 {len(pages)} 页")

        pages = outline.get('pages', [])
        if not isinstance(pages, list) or len(pages) == 0:
            return jsonify({'status': 'error', 'message': '大纲格式错误：缺少 pages 数组或 pages 为空。如果是 Markdown 文件，请确保内容有明确的章节/标题结构。'}), 400

        for i, page in enumerate(pages):
            if not isinstance(page, dict):
                page = {'index': i + 1, 'title': f'第{i+1}页', 'type': 'content', 'layout': '满版图片-全屏背景', 'content': str(page)}
                pages[i] = page
            if not page.get('title'):
                page['title'] = f'第{i+1}页'
            if not page.get('type'):
                page['type'] = 'content'
            if not page.get('layout'):
                page['layout'] = '满版图片-全屏背景'
            if 'index' not in page:
                page['index'] = i + 1

        session_id = str(uuid.uuid4())
        session_data = {
            'outline': outline,
            'images': [],
            'created_at': time.time()
        }
        _save_outline_session(session_id, session_data)
        _uploaded_outline_cache[session_id] = session_data

        return jsonify({
            'status': 'success',
            'session_id': session_id,
            'outline': outline,
            'page_count': len(pages),
            'message': f'大纲已上传{"（已从 Markdown 解析为结构化大纲）" if is_markdown else ""}，共 {len(pages)} 页'
        })

    except json.JSONDecodeError as e:
        return jsonify({'status': 'error', 'message': f'JSON 解析失败: {str(e)}。如果上传的是 Markdown 文件，请确保文件扩展名为 .md'}), 400
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/generate-from-outline', methods=['POST'])
def generate_from_outline():
    """根据已上传的大纲生成图片（SSE 流式进度）"""
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    data = request.get_json(silent=True) or {}
    session_id = data.get('session_id', '')
    style_settings = data.get('style', {})
    direct_outline = data.get('outline')

    outline = None

    if session_id and session_id in _uploaded_outline_cache:
        store = _uploaded_outline_cache[session_id]
        outline = store['outline']
    elif session_id:
        # 尝试从磁盘加载
        loaded = _load_outline_session(session_id)
        if loaded:
            store = loaded
            outline = store['outline']
            _uploaded_outline_cache[session_id] = loaded
        else:
            outline = None
    elif direct_outline:
        outline = direct_outline
        if not isinstance(outline, dict):
            def _err_gen():
                yield f"data: {json.dumps({'status': 'error', 'message': 'outline 格式错误，需要对象'}, ensure_ascii=False)}\n\n"
            return Response(
                stream_with_context(_err_gen()),
                mimetype='text/event-stream',
                headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'}
            )
        pages = outline.get('pages', [])
        if not pages:
            def _err_gen():
                yield f"data: {json.dumps({'status': 'error', 'message': 'outline 中没有 pages 数据'}, ensure_ascii=False)}\n\n"
            return Response(
                stream_with_context(_err_gen()),
                mimetype='text/event-stream',
                headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'}
            )
    else:
        def _err_gen():
            yield f"data: {json.dumps({'status': 'error', 'message': '请先上传大纲（提供 session_id 或直接传 outline）'}, ensure_ascii=False)}\n\n"
        return Response(
            stream_with_context(_err_gen()),
            mimetype='text/event-stream',
            headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'}
        )

    style_desc = style_settings.get('description', '')
    color_scheme = style_settings.get('colorScheme', 'auto')
    reference_image = style_settings.get('referenceImage', '')
    product_image = style_settings.get('productImage', '')
    brand_name = style_settings.get('brandName', '')

    def generate():
        try:
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)

            image_api_url = config.get('image_api_base_url') or config.get('api_base_url')
            image_api_key = config.get('image_api_key') or config.get('api_key')

            if not image_api_url or not image_api_key:
                yield f"data: {json.dumps({'status': 'error', 'message': '请先配置图片生成模型的 API 设置'}, ensure_ascii=False)}\n\n"
                return

            pages = outline.get('pages', [])
            if not pages:
                yield f"data: {json.dumps({'status': 'error', 'message': '大纲为空'}, ensure_ascii=False)}\n\n"
                return

            api = APIClient(image_api_url, image_api_key)
            image_model = config.get('image_model', 'gpt-image-2')
            chat_model = config.get('chat_model', 'gpt-4o')
            total = len(pages)
            images = []

            color_scheme_val = style_settings.get('colorScheme', '')
            color_prompt = color_scheme_val if color_scheme_val else 'appropriate color scheme matching the content'

            brand_colors = ""
            if brand_name and not color_scheme_val:
                try:
                    yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': f'正在识别 {brand_name} 的品牌色系...'}, ensure_ascii=False)}\n\n"

                    color_prompt_text = f"""Identify the brand colors for "{brand_name}".
Return ONLY a JSON object with:
- primary_color: main brand color (hex code and name)
- secondary_color: accent color (hex code and name)
- background_suggestion: recommended background style
- text_color: recommended text color

Example for "支付宝":
{{"primary_color": "#1677FF (支付宝蓝)", "secondary_color": "#FF6A00 (支付宝橙)", "background_suggestion": "light or white background", "text_color": "#333333 or white on dark backgrounds"}}"""

                    color_result = api.chat(
                        model=chat_model,
                        messages=[{'role': 'user', 'content': color_prompt_text}],
                        max_tokens=200,
                        timeout=30
                    )
                    brand_colors = color_result.get('content', '')
                    print(f"[INFO] 品牌 {brand_name} 色系: {brand_colors}")
                except Exception as e:
                    print(f"[WARN] 品牌色系识别失败: {e}")
                    brand_colors = ""

            if color_scheme_val:
                brand_colors = f"User specified color scheme: {color_scheme_val}"

            reference_style = ""
            if reference_image and reference_image.startswith('data:image'):
                try:
                    yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': '正在分析参考图风格...'}, ensure_ascii=False)}\n\n"

                    analysis_prompt = "Analyze this reference image and describe its visual style in detail. Include: color palette (list specific hex colors), typography style, layout characteristics, visual elements, overall mood/atmosphere. Keep it under 150 words."

                    messages = [
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": analysis_prompt},
                                {"type": "image_url", "image_url": {"url": reference_image}}
                            ]
                        }
                    ]

                    analysis_result = api.chat(
                        model=chat_model,
                        messages=messages,
                        max_tokens=300,
                        timeout=60
                    )
                    reference_style = analysis_result.get('content', '')
                    print(f"[INFO] 参考图风格分析: {reference_style[:100]}...")
                except Exception as e:
                    print(f"[WARN] 参考图分析失败: {e}")
                    reference_style = ""

            unified_style_guide = ""
            is_muji_style = style_desc and ('无印良品' in style_desc or 'MUJI' in style_desc.upper() or 'muji' in style_desc.lower() or '日式极简' in style_desc or '日式排版' in style_desc)

            if is_muji_style:
                unified_style_guide = """MUJI (无印良品) Japanese Minimalist Style - STRICT RULES:
- Background: Warm off-white #F5F5F0 (like washi paper), NEVER dark, NEVER gradient
- Primary text: Soft charcoal #4A4A4A (NOT pure black), left-aligned ONLY
- Accent color: Natural beige #D4C4B0 (wood/linen feel), use sparingly
- Layout: 60% white space (MA/間), asymmetric balance, grid-based
- Typography: Clean sans-serif, generous line spacing, hierarchical sizing
- Visuals: Natural textures only, single focal point per slide, NO decorative borders/gradients/shadows
- Philosophy: "This is enough" (これでいい) - understated, honest, essential
- ABSOLUTELY NO: Bright colors, gradients, centered text (except cover), flashy effects, pure black backgrounds"""
                yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': '已应用无印良品精确风格规范...'}, ensure_ascii=False)}\n\n"
                print("[INFO] 应用无印良品精确风格规范")
            elif style_desc or color_prompt or reference_style:
                try:
                    yield f"data: {json.dumps({'status': 'progress', 'current': 0, 'total': total, 'message': '正在确定统一风格规范...'}, ensure_ascii=False)}\n\n"

                    style_guide_prompt = f"""You are a professional brand designer. Based on the following style requirements, define a UNIFIED visual style guide for a PPT presentation. All slides must follow this EXACT same style to maintain consistency, BUT each slide should have varied backgrounds and moods.

Style description: {style_desc or 'Not specified'}
Color preference: {color_prompt or 'Not specified'}
Reference image style: {reference_style or 'No reference'}

Please output a concise style guide in this format (under 250 words):
- Primary colors: (list 2-3 specific colors with hex codes)
- Secondary colors: (list 1-2 accent colors)
- Background VARIETY: (list 5+ different background styles to rotate across slides: e.g. "daylight urban", "soft indoor", "blue hour cityscape", "abstract geometric", "warm lifestyle", "cool minimal", etc.)
- Typography mood: (modern/classic/playful/minimal)
- Visual elements: (shapes, lines, textures to use consistently)
- Overall atmosphere: (one sentence)
- ANTI-PATTERN: explicitly list what to AVOID (e.g. "NOT every slide golden hour sunset", "NOT every slide with person's back silhouette")

IMPORTANT: Be specific about colors. The background variety is CRITICAL - each slide must look visually distinct while maintaining color and typography consistency."""

                    guide_result = api.chat(
                        model=chat_model,
                        messages=[{'role': 'user', 'content': style_guide_prompt}],
                        max_tokens=400,
                        timeout=60
                    )
                    unified_style_guide = guide_result.get('content', '')
                    print(f"[INFO] 统一风格规范: {unified_style_guide[:200]}...")
                except Exception as e:
                    print(f"[WARN] 风格规范生成失败: {e}")
                    unified_style_guide = ""

            print(f"[IMG] ========== 上传模式-开始生图循环: 共{total}页, 模型={image_model}, API={image_api_url} ==========", flush=True)
            for i, page in enumerate(pages):
                page_title = page.get('title', f'第{i+1}页')
                page_content = page.get('content', '')
                page_type = page.get('type', 'content')
                page_layout = page.get('layout', '')
                final_prompt = None
                if '模块化-三列卡片' in page_layout:
                    page_layout = page_layout.replace('模块化-三列卡片', '模块化-卡片')

                yield f"data: {json.dumps({'status': 'progress', 'current': i+1, 'total': total, 'message': f'正在生成第 {i+1}/{total} 页: {page_title}'}, ensure_ascii=False)}\n\n"

                is_fullbleed_page = '满版图片' in page_layout or \
                    page_type in ('cover', 'visual')

                if '满版图片-全屏背景' in page_layout or '满版图片-电影感' in page_layout:
                    content_lower = page_content.lower() if page_content else ''
                    if any(kw in page_content for kw in ['数据', '分析', '趋势', '增长', '用户', '比例', '调研']):
                        bg_style = """Background: abstract data visualization style, geometric patterns, subtle gradient, modern tech feel."""
                    elif any(kw in page_content for kw in ['TVC', '脚本', '画面', '镜头', '旁白']):
                        bg_style = """Background: video production scene, camera equipment, filming location, or storyboard sketches."""
                    elif any(kw in page_content for kw in ['愿望', '心愿', '梦想', '希望', '坚持', '努力']):
                        bg_style = """Background: DIVERSE lifestyle scenes - CAN USE: cafe window with rain, street food stall, bus window, park bench, rooftop at sunset, night market, bookstore aisle, cozy bed with string lights. AVOID: desk lamp, person at desk, indoor study room."""
                    elif any(kw in page_content for kw in ['品牌', '站位', '定位', '策略', '传播']):
                        bg_style = """Background: clean modern workspace or urban architecture, professional atmosphere, bright daylight or blue hour. AVOID: desk lamp scenes."""
                    elif any(kw in page_content for kw in ['洞察', '发现', '现象', '观察']):
                        bg_style = """Background: street photography, documentary style, authentic urban moments, natural daylight. AVOID: desk lamp, indoor cozy scenes."""
                    else:
                        bg_style = """Background: HIGHLY VARIED scenes - use different settings for each slide: outdoor cafe, public park, metro/train, night market, rooftop, bookstore, street corner, beach, mountain trail, market stall. AVOID: desk lamp, person at desk."""

                    diverse_scenes = [
                        "street food stall with warm lights",
                        "person reading on a park bench",
                        "cafe window view with rain",
                        "night market with colorful lights",
                        "rooftop with city skyline at dusk",
                        "bookstore aisle with soft lighting",
                        "metro window with passing scenery",
                        "coastal walkway at sunset",
                        "flower market early morning",
                        "outdoor market stall",
                    ]

                    # 产品图约束
                    product_constraint = ""
                    if product_image:
                        product_constraint = """

【重要 - 产品图严格约束】
用户提供了实际的产品参考图。如果当前页面内容涉及以下场景，必须在生成的图片中严格保持产品的真实外观和外形特征，不得重新绘制或想象一个不同的产品：
- 产品展示、产品特写、产品使用场景
- 包含产品名称、品牌名称的页面
- 任何需要出现该实物产品的画面

处理方式：
- 如果是产品展示页：直接使用产品原图作为主要视觉元素，叠加文字信息
- 如果是使用场景页：将产品图自然融入场景中，保持产品外形100%一致
- 不要重新设计产品的外形、颜色、材质
"""

                    style_text = f"""PRECISE LAYOUT MATCH - 满版图片-全屏背景 (Full-bleed background):
- FULL SCREEN high-quality background image
- {bg_style}
- ANTI-PATTERN: STRICTLY FORBIDDEN - DO NOT use these common repeating patterns:
  * Person sitting at desk with desk lamp (this pattern is OVERUSED)
  * Person writing in notebook at study desk
  * Indoor cozy study room with warm lamp light
  * Student at desk with books and lamp
  * Heavy cinematic color grading (orange/teal tones)
  * Dramatic movie-style lighting
- Background scene MUST match the page's actual content topic — see SCENE DESCRIPTION below for details
- White text directly overlaid with subtle dark gradient overlay for readability
- Top-left corner: small brand-colored rectangular tag with white text for section label
- Large expressive title (60-80px) in white, bold
- Clean body text (16-20px) in white or light gray
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- NO modular cards, NO split layout - pure full-bleed photography with text overlay
- IMPORTANT: Each slide MUST have a UNIQUE background scene that matches its specific content
- IMPORTANT: Use NATURAL lighting and colors, avoid heavy color grading or cinematic effects"""

                elif '满版图片-左图右文' in page_layout:
                    has_app = any(kw in page_content for kw in ['APP', 'app', '界面', '截图', '手机', '功能', '玩法'])
                    has_social = any(kw in page_content for kw in ['抖音', '小红书', '微博', '社交', '话题', '晒图', 'KOL', 'KOC'])
                    has_people = any(kw in page_content for kw in ['人物', '用户', '年轻人', '场景'])

                    if has_app:
                        left_desc = "realistic smartphone mockup showing app interface with detailed UI elements"
                    elif has_social:
                        left_desc = "social media feed screenshot or phone showing social platform interface"
                    elif has_people:
                        left_desc = "authentic lifestyle photography of young Chinese people in daily scenes"
                    else:
                        left_desc = "high-quality thematic photograph or illustration matching the page topic"

                    style_text = f"""PRECISE LAYOUT MATCH - 满版图片-左图右文 (Split layout):
- LEFT side (40-50%): {left_desc}
- RIGHT side (50-60%): Clean white background text area with structured content layout
- Top-left: brand-colored rectangular section tag
- Right area: Large brand-colored title + well-organized body text
- Body text should use bullet points, numbered lists, or small info cards for visual interest
- Add subtle design elements: thin colored dividers, small icons, or accent lines
- Clear visual separation between image and text areas
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional, polished, agency-quality design"""

                elif '模块化-卡片' in page_layout:
                    card_count = _guess_card_count(page_content)
                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-卡片 (Card module with dynamic columns):
- {card_count} equal-width cards arranged HORIZONTALLY in a row
- Number of cards should match content: if content has 4 items, use 4 cards; if 5 items, use 5 cards
- Each card: rounded corners (12px radius), subtle shadow, white/light background
- Card header: icon (line style) + title in brand primary color
- Card body: brief description in dark gray
- Cards separated by consistent spacing (24px gap)
- Top section: page title with brand color accent
- Background: light warm gray or subtle texture
- Visual hierarchy: cards are the focal point
- Each card self-contained with clear visual boundaries
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

                elif '模块化-步骤流程' in page_layout:
                    has_app_content = any(kw in page_content for kw in ['APP', 'app', '界面', '截图', '手机', '功能', '玩法'])
                    step_visual = "Include realistic smartphone mockups showing app screens at each step" if has_app_content else "Use icons and illustrations for each step, NO phone mockups unless content mentions app"

                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-步骤流程 (Step-by-step flow):
- Vertical or horizontal flow with STEP 1 / STEP 2 / STEP 3 markers
- Each step: numbered badge (brand primary color circle with white number) + title + description
- {step_visual}
- Connect steps with subtle arrows or lines
- Each step module has light background with subtle shadow
- Mix of text and visual elements (icons, illustrations, or phone screens where relevant)
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional social media campaign presentation style"""

                elif '模块化-时间轴' in page_layout:
                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-时间轴 (Timeline):
- Horizontal timeline with three nodes: 预热期 → 爆发期 → 长尾期
- Each node: circular marker (brand primary color fill) + phase name + key actions
- Timeline line connecting all nodes (brand color or gray)
- Below timeline: detailed modules for each phase with images and text
- Show campaign progression visually
- Include small preview images or icons for each phase
- Clean, organized, easy to follow at a glance
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional campaign timeline design"""

                elif '模块化-表格' in page_layout:
                    style_text = f"""PRECISE LAYOUT MATCH - 模块化-表格 (Data table):
- Multi-column table with clear headers (时间/阶段/目标/主题/动作)
- Header row: brand primary color background with white text
- Data rows: alternating light backgrounds for readability
- Clean grid lines, professional alignment
- Include small icons or color coding for different phases
- Readable font sizes, clear hierarchy
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Campaign overview table style"""

                elif '纯视觉页' in page_layout:
                    style_text = f"""PRECISE LAYOUT MATCH - 纯视觉页 (Visual only):
- NO text or minimal text only
- Full-screen key visual (KV) display
- Clean, impactful imagery
- Could be single KV or multiple KV layouts shown together
- Focus on visual impact and brand aesthetics
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- Professional advertising campaign visual presentation"""

                else:
                    is_text_only = _is_text_only_layout(page_layout)

                    if is_text_only:
                        _bc2 = brand_colors if brand_colors else 'elegant deep navy & warm gold accent'
                        _bg2 = f"background: soft gradient from near-white (#FAFAFA) to very light brand-tinted tone, with subtle geometric pattern"
                        text_style_map = {
                            '文字-金句引言': f"""TEXT-ONLY LAYOUT - 金句引言 (Premium Quote Slide):
- BACKGROUND: {_bg2}
- QUOTE TEXT:超大字号(72-96pt), 品牌主色, 加粗, 居中或偏左黄金分割位置
- DECORATIVE: 左侧或顶部有品牌色粗线条装饰(4-8px宽)或大型半透明引号装饰(200pt+, 10% opacity)
- TITLE AREA: 页面顶部有小标签(badge), 圆角矩形, 品牌色底+白字
- SPACING: 大量留白(40%+页面空白)
- NO photographs — pure typography art
- BRAND COLORS: {_bc2}
- VIBE: 杂志封面级排版, Apple Keynote / Notion 风格极简高级感""",
                            '文字-数字大屏': f"""TEXT-ONLY LAYOUT - 数字大屏 (Data Dashboard - Premium):
- BACKGROUND: {_bg2} + 微妙网格线(grid lines at 3% opacity)
- NUMBERS: 超大数字(100-140pt), 品牌主色, Bold/Black, 每个在独立圆角卡片内(white bg, subtle shadow)
- ACCENT: 趋势箭头↑↓或百分比指示器(品牌色)
- TITLE: 顶部品牌色标题(36-48pt bold)
- NO photographs
- BRAND COLORS: {_bc2}
- VIBE: Bloomberg Terminal / 企业数据大屏风格""",
                            '文字-问答式': f"""TEXT-ONLY LAYOUT - 问答式 (Q&A - Premium Editorial):
- BACKGROUND: {_bg2}, 左侧品牌色竖条装饰(8px宽)
- Q: 品牌主色, 28-36pt Bold, 前面有装饰性「Q」符号
- A: 深灰#333, 18-20pt, 行距1.6-1.8
- LAYOUT: Q&A间细线分隔(dashed line), 序号圆圈(品牌色底白字)
- NO photographs
- BRAND COLORS: {_bc2}
- VIBE: FAQ / 知乎精选风格, 清晰且舒适""",
                            '文字-双栏对比': f"""TEXT-ONLY LAYOUT - 双栏对比 (Premium Split Design):
- BACKGROUND: {_bg2}
- 50/50分栏, 中间VS分隔区域(圆形徽章或竖线+箭头)
- 两栏卡片式(headers in brand color, white bg, rounded corners)
- 关键差异词用品牌色加粗标注
- NO photographs
- BRAND COLORS: {_bc2}
- VIBE: 产品功能对比风格, 清晰且有设计感""",
                            '文字-要点列表': f"""TEXT-ONLY LAYOUT - 要点列表 (Premium Structured):
- BACKGROUND: {_bg2}, 极淡品牌色方向光效
- TITLE: 品牌主色 36-48pt bold, 下方品牌色下划线装饰
- LIST ITEMS: 大号序号圆圈(品牌色底白字, 32-40px diameter)
- ITEM TEXT: 标题20-24pt Bold + 描述16-18pt regular, 行距1.5
- DECORATIVE: 淡色几何点缀, 可选奇偶行微弱背景差异
- NO photographs
- BRAND COLORS: {_bc2}
- VIBE: Apple Feature List / Notion Database 风格""",
                            '文字-引用来源': f"""TEXT-ONLY LAYOUT - 引用来源 (Premium Academic):
- BACKGROUND: {_bg2}
- TITLE: 「参考资料」, 品牌色 28-32pt bold
- EACH ENTRY: 来源名(Bold 18pt 品牌色) + 摘要(15-16pt #444) + 年份(gray 13pt)
- CONTAINER: 浅色圆角容器内(white bg, shadow, 16px radius, padding 40px)
- NO photographs
- BRAND COLORS: {_bc2}
- VIBE: 学术参考文献风格, 专业严谨但美观"""
                        }
                        style_text = text_style_map.get(page_layout, f"""TEXT-ONLY LAYOUT (Premium Text Design):
- BACKGROUND: soft gradient from near-white to light brand-tinted tone, with subtle geometric pattern
- TITLE in brand primary color (36-48pt bold), body text in dark gray (#333, 16-18pt)
- Strong typography hierarchy: clear size and weight contrast between heading/subheading/body
- Generous whitespace (35-45% of slide area should be empty)
- Decorative elements: brand-colored accent bars, thin divider lines, subtle geometric shapes
- Content organized in clean cards or structured sections with rounded corners
- NO photographs, but decorative graphic elements (lines, shapes, color blocks) are encouraged
- BRAND COLORS: {_bc2}
- VIBE: Modern editorial / magazine layout quality — NOT a Word document""")

                        final_prompt = f"""Create a professional PPT slide.

Title: '{page_title}'
Content: '{page_content}'
Layout Type: '{page_layout}'

{style_text}

CRITICAL REQUIREMENTS - TEXT ONLY SLIDE:
- 16:9 aspect ratio
- This is a TEXT-ONLY slide — do NOT add any photos, illustrations, people, scenes, or decorative images
- Focus entirely on typography, layout, spacing, and text hierarchy
- Clean, minimal, professional text-focused design{product_constraint}"""

                    else:
                        bc_text = brand_colors if brand_colors else 'Use appropriate brand colors for accents and highlights'

                        _is_strong_image2 = any(k in page_layout for k in ['满版图片', '纯视觉页'])
                        _is_module2 = any(k in page_layout for k in ['模块化-卡片', '模块化-步骤', '模块化-时间轴', '模块化-表格'])

                        _scene_desc2 = _build_scene_description(page_title, page_content)

                        if _is_strong_image2:
                            visual_elements = (
                                "\nKEY VISUAL ELEMENTS:\n"
                                "1. SCENE (content-driven): " + _scene_desc2 + "\n"
                                "2. Only include smartphone mockups IF the content mentions: app, mobile interface, or mobile features\n"
                                "3. Brand colors: " + bc_text + "\n"
                                "4. Typography: Clean sans-serif Chinese fonts, large expressive titles\n"
                                "5. Design details: Professional polish, agency-quality finish\n"
                            )
                            crITICAL2 = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- Use authentic photography or strong imagery as the main visual focal point
- Clean, modern, agency-quality presentation design{product_constraint}"""

                        elif _is_module2:
                            visual_elements = (
                                "\nKEY VISUAL ELEMENTS:\n"
                                "1. Visual style: Clean graphic design with icons, illustrations, or abstract shapes — NOT lifestyle photography of people\n"
                                "2. Card/module design: Each section should be self-contained with icon + title + text structure\n"
                                "3. DO NOT add photos of people, lifestyle scenes, or stock photography — they clutter module layouts\n"
                                "4. Brand colors: " + bc_text + "\n"
                                "5. Typography: Clear hierarchy, readable body text, bold headers\n"
                                "6. Design details: Rounded corners, subtle shadows, clean lines, generous whitespace\n"
                            )
                            crITICAL2 = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- This is a MODULE/CARD layout — focus on structured information display, NOT on photography
- NO lifestyle photos, NO people scenes, NO stock photography
- Let the card/module structure and icons carry the visual weight{product_constraint}"""

                        else:
                            visual_elements = (
                                "\nKEY VISUAL ELEMENTS:\n"
                                "1. Visual style: Balanced mix of graphics and text based on content needs\n"
                                "2. Only add photography IF the content specifically calls for it\n"
                                "3. Otherwise use: icons, illustrations, charts, diagrams, or clean typography\n"
                                "4. Brand colors: " + bc_text + "\n"
                                "5. Design details: Professional polish, appropriate whitespace\n"
                            )
                            crITICAL2 = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- Do not force photography into layouts that don't need it
- Clean, modern presentation design{product_constraint}"""

                        final_prompt = f"""Create a professional PPT slide.

Title: '{page_title}'
Content: '{page_content}'
Layout Type: '{page_layout}'

{style_text}

{visual_elements}

{crITICAL2}"""

                if final_prompt is None:
                    _ve = ""
                    _cr = """CRITICAL REQUIREMENTS:
- 16:9 aspect ratio
- High quality, professional design
- Clean, modern presentation design"""
                    if any(k in page_layout for k in ['满版图片', '纯视觉页']):
                        _ve = "\nKEY VISUAL ELEMENTS:\n1. Photography as main visual element\n2. Brand colors: " + (brand_colors or 'appropriate accent colors') + "\n"
                    elif any(k in page_layout for k in ['模块化-卡片', '模块化-步骤', '模块化-时间轴', '模块化-表格']):
                        _ve = "\nKEY VISUAL ELEMENTS:\n1. Icons, illustrations, geometric shapes — NO lifestyle photos\n2. Brand colors: " + (brand_colors or 'appropriate accent colors') + "\n"
                    final_prompt = f"""Create a professional PPT slide.

Title: '{page_title}'
Content: '{page_content}'
Layout Type: '{page_layout}'

{style_text}

{_ve}
{_cr}{product_constraint if product_image else ''}"""

                _t0 = time.time()
                print(f"[IMG] PROMPT前200字: {final_prompt[:200]}", flush=True)
                try:
                    _needs_prod2 = _page_needs_product_image(page_title, page_content, page_type, page_layout)

                    if _needs_prod2 and product_image and product_image.startswith('data:image'):
                        print(f"[IMG] 第{i+1}/{total}页「{page_title}」→ edits模式(产品图)", flush=True)
                        messages = [
                            {"role": "user", "content": [
                                {"type": "text", "text": final_prompt},
                                {"type": "image_url", "image_url": {"url": product_image}}
                            ]}
                        ]
                        result = api.image_generate(
                            model=image_model,
                            prompt=final_prompt,
                            size='1792x1024',
                            n=1,
                            messages=messages
                        )
                    else:
                        print(f"[IMG] 第{i+1}/{total}页「{page_title}」→ 标准生成(mode={image_model})", flush=True)
                        result = api.image_generate(
                            model=image_model,
                            prompt=final_prompt,
                            size='1792x1024',
                            n=1
                        )

                    _elapsed = time.time() - _t0
                    print(f"[IMG] API返回耗时={_elapsed:.1f}s, result类型={type(result).__name__}, keys={list(result.keys()) if isinstance(result, dict) else 'N/A'}", flush=True)

                    img_url = None
                    _raw_images = result.get('images', []) if isinstance(result, dict) else []
                    print(f"[IMG] images字段长度={len(_raw_images)}", flush=True)
                    for idx, img in enumerate(_raw_images):
                        _has_url = bool(img.get('url')) if isinstance(img, dict) else False
                        _has_b64 = bool(img.get('b64_json')) if isinstance(img, dict) else False
                        print(f"[IMG]   图片[{idx}]: has_url={_has_url}, has_b64={_has_b64}, b64_len={len(str(img.get('b64_json',''))) if _has_b64 else 0}", flush=True)
                        if img.get('url'):
                            img_url = img['url']
                        elif img.get('b64_json'):
                            try:
                                import base64
                                raw = base64.b64decode(img['b64_json'])
                                fname = f"{int(time.time())}_{i+1}_{uuid.uuid4().hex[:10]}.png"
                                fpath = os.path.join(GENERATED_DIR, fname)
                                with open(fpath, "wb") as f:
                                    f.write(raw)
                                img_url = f"/api/generated/{fname}"
                                print(f"[IMG]   落盘成功: {fname} ({len(raw)}bytes)", flush=True)
                            except Exception as e:
                                print(f"[WARN] 第{i+1}页图片落盘失败: {e}", flush=True)
                                img_url = None

                    if img_url:
                        images.append(img_url)
                        print(f"[IMG] ✅ 第{i+1}页成功: {img_url[:60]}...", flush=True)
                        yield f"data: {json.dumps({'status': 'page_done', 'current': i+1, 'total': total, 'image': img_url, 'message': f'第 {i+1} 页生成成功'}, ensure_ascii=False)}\n\n"
                    else:
                        images.append(f"/api/placeholder/800/450")
                        print(f"[IMG] ❌ 第{i+1}页失败: 无有效图片URL (耗时{_elapsed:.1f}s)", flush=True)
                        yield f"data: {json.dumps({'status': 'page_fail', 'current': i+1, 'total': total, 'message': f'第 {i+1} 页生成失败（无图片返回）'}, ensure_ascii=False)}\n\n"
                except Exception as e:
                    images.append(f"/api/placeholder/800/450")
                    _elapsed = time.time() - _t0
                    print(f"[IMG] 💥 第{i+1}页异常: {type(e).__name__}: {str(e)[:200]} (耗时{_elapsed:.1f}s)", flush=True)
                    import traceback
                    traceback.print_exc()
                    yield f"data: {json.dumps({'status': 'page_fail', 'current': i+1, 'total': total, 'message': f'第 {i+1} 页失败: {str(e)[:80]}'}, ensure_ascii=False)}\n\n"

            store['images'] = images
            _save_outline_session(session_id, store)  # 持久化到磁盘
            yield f"data: {json.dumps({'status': 'done', 'images': images, 'session_id': session_id, 'message': f'完成！共生成 {len(images)} 张图片'}, ensure_ascii=False)}\n\n"

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield f"data: {json.dumps({'status': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive',
        }
    )


@app.route('/api/regenerate-single-page', methods=['POST'])
def regenerate_single_page():
    """重新生成指定单页的图片（SSE 流式进度）"""
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    data = request.get_json(silent=True) or {}
    page_index = data.get('page_index', 0)
    page = data.get('page', {})
    style_settings = data.get('style', {})
    session_id = data.get('session_id', '')

    if not page or not isinstance(page, dict):
        def _err_gen():
            yield f"data: {json.dumps({'status': 'error', 'message': '缺少页面数据'}, ensure_ascii=False)}\n\n"
        return Response(stream_with_context(_err_gen()), mimetype='text/event-stream',
            headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'})

    style_desc = style_settings.get('description', '')
    color_scheme = style_settings.get('colorScheme', 'auto')
    reference_image = style_settings.get('referenceImage', '')
    product_image = style_settings.get('productImage', '')
    brand_name = style_settings.get('brandName', '')

    def generate():
        try:
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)

            image_api_url = config.get('image_api_base_url') or config.get('api_base_url')
            image_api_key = config.get('image_api_key') or config.get('api_key')

            if not image_api_url or not image_api_key:
                yield f"data: {json.dumps({'status': 'error', 'message': '请先配置图片生成模型的 API 设置'}, ensure_ascii=False)}\n\n"
                return

            api = APIClient(image_api_url, image_api_key)
            image_model = config.get('image_model', 'gpt-image-2')
            chat_model = config.get('chat_model', 'gpt-4o')

            color_scheme_val = style_settings.get('colorScheme', '')
            color_prompt = color_scheme_val if color_scheme_val else 'appropriate color scheme matching the content'

            brand_colors = ""
            if brand_name and not color_scheme_val:
                try:
                    yield f"data: {json.dumps({'status': 'progress', 'message': f'正在识别品牌色系...'}, ensure_ascii=False)}\n\n"
                    color_prompt_text = f"""Identify the brand colors for "{brand_name}".
Return ONLY a JSON object with:
- primary_color: main brand color (hex code and name)
- secondary_color: accent color (hex code and name)
- background_suggestion: recommended background style
- text_color: recommended text color"""
                    color_result = api.chat(model=chat_model,
                        messages=[{'role': 'user', 'content': color_prompt_text}], max_tokens=200, timeout=30)
                    brand_colors = color_result.get('content', '')
                except Exception:
                    brand_colors = ""

            if color_scheme_val:
                brand_colors = f"User specified color scheme: {color_scheme_val}"

            reference_style = ""
            if reference_image and reference_image.startswith('data:image'):
                try:
                    analysis_prompt = "Analyze this reference image and describe its visual style in detail. Include: color palette, typography style, layout characteristics, visual elements, overall mood. Keep it under 150 words."
                    messages = [{"role": "user", "content": [
                        {"type": "text", "text": analysis_prompt},
                        {"type": "image_url", "image_url": {"url": reference_image}}
                    ]}]
                    analysis_result = api.chat(model=chat_model, messages=messages, max_tokens=300, timeout=60)
                    reference_style = analysis_result.get('content', '')
                except Exception:
                    reference_style = ""

            unified_style_guide = ""
            is_muji_style = style_desc and ('无印良品' in style_desc or 'MUJI' in style_desc.upper() or 'muji' in style_desc.lower() or '日式极简' in style_desc or '日式排版' in style_desc)
            if is_muji_style:
                unified_style_guide = """MUJI (无印良品) Japanese Minimalist Style - STRICT RULES:
- Background: Warm off-white #F5F5F0, NEVER dark, NEVER gradient
- Primary text: Soft charcoal #4A4A4A, left-aligned ONLY
- Accent color: Natural beige #D4C4B0
- Layout: 60% white space, asymmetric balance, grid-based
- Typography: Clean sans-serif, generous line spacing
- Visuals: Natural textures only, single focal point per slide
- ABSOLUTELY NO: Bright colors, gradients, centered text, flashy effects"""
            elif style_desc or color_prompt or reference_style:
                try:
                    style_guide_prompt = f"""You are a professional brand designer. Based on the following, define a UNIFIED visual style guide for a PPT presentation.
Style description: {style_desc or 'Not specified'}
Color preference: {color_prompt or 'Not specified'}
Reference image style: {reference_style or 'No reference'}

Output concise style guide (under 250 words):
- Primary colors: (list 2-3 specific hex colors)
- Secondary colors: (1-2 accent colors)
- Background VARIETY: (5+ different background styles to rotate)
- Typography mood: (modern/classic/playful/minimal)
- Visual elements: (shapes, lines, textures)
- Overall atmosphere: (one sentence)
- ANTI-PATTERN: what to AVOID"""
                    guide_result = api.chat(model=chat_model,
                        messages=[{'role': 'user', 'content': style_guide_prompt}], max_tokens=400, timeout=60)
                    unified_style_guide = guide_result.get('content', '')
                except Exception:
                    unified_style_guide = ""

            page_title = page.get('title', f'第{page_index+1}页')
            page_content = page.get('content', '')
            page_type = page.get('type', 'content')
            page_layout = page.get('layout', '')
            final_prompt = None

            product_constraint = ""
            if product_image:
                product_constraint = """
【重要 - 产品图严格约束】
用户提供了实际的产品参考图。如果当前页面内容涉及产品展示、产品特写、产品使用场景、包含产品名称/品牌名称的页面，必须在生成的图片中严格保持产品的真实外观和外形特征。
处理方式：产品展示页直接使用产品原图作为主要视觉元素；使用场景页将产品图自然融入场景中，保持产品外形100%一致。
不要重新设计产品的外形、颜色、材质"""

            yield f"data: {json.dumps({'status': 'progress', 'message': f'正在构建第 {page_index+1} 页的生图提示...'}, ensure_ascii=False)}\n\n"

            if '模块化-三列卡片' in page_layout:
                page_layout = page_layout.replace('模块化-三列卡片', '模块化-卡片')

            is_fullbleed_page = '满版图片' in page_layout or page_type in ('cover', 'visual')

            if '满版图片-全屏背景' in page_layout or '满版图片-电影感' in page_layout:
                content_lower = page_content.lower() if page_content else ''
                if any(kw in page_content for kw in ['数据', '分析', '趋势', '增长', '用户', '比例', '调研']):
                    bg_style = """Background: abstract data visualization style, geometric patterns, subtle gradient, modern tech feel."""
                elif any(kw in page_content for kw in ['TVC', '脚本', '画面', '镜头', '旁白']):
                    bg_style = """Background: video production scene, camera equipment, filming location, or storyboard sketches."""
                elif any(kw in page_content for kw in ['愿望', '心愿', '梦想', '希望', '坚持', '努力']):
                    bg_style = """Background: DIVERSE lifestyle scenes - CAN USE: cafe window with rain, street food stall, bus window, park bench, rooftop at sunset, night market, bookstore aisle, cozy bed with string lights."""
                elif any(kw in page_content for kw in ['品牌', '站位', '定位', '策略', '传播']):
                    bg_style = """Background: clean modern workspace or urban architecture, professional atmosphere, bright daylight or blue hour."""
                elif any(kw in page_content for kw in ['洞察', '发现', '现象', '观察']):
                    bg_style = """Background: street photography, documentary style, authentic urban moments, natural daylight."""
                else:
                    bg_style = """Background: HIGHLY VARIED scenes - use different settings: outdoor cafe, public park, metro/train, night market, rooftop, bookstore, street corner."""

                diverse_scenes = ["street food stall with warm lights", "person reading on a park bench",
                    "cafe window view with rain", "night market with colorful lights",
                    "rooftop with city skyline at dusk", "bookstore aisle with soft lighting",
                    "metro window with passing scenery", "coastal walkway at sunset"]

                style_text = f"""PRECISE LAYOUT MATCH - 满版图片-全屏背景 (Full-bleed background):
- FULL SCREEN high-quality background image
- {bg_style}
- ANTI-PATTERN: STRICTLY FORBIDDEN - DO NOT use: person sitting at desk with desk lamp, person writing in notebook at study desk, indoor cozy study room with warm lamp light, heavy cinematic color grading
- Background scene MUST match the page's actual content topic — see SCENE DESCRIPTION below for details
- White text directly overlaid with subtle dark gradient overlay for readability
- Top-left corner: small brand-colored rectangular tag with white text for section label
- Large expressive title (60-80px) in white, bold
- Clean body text (16-20px) in white or light gray
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}
- NO modular cards, NO split layout - pure full-bleed photography with text overlay"""

            elif '满版图片-左图右文' in page_layout:
                has_app = any(kw in page_content for kw in ['APP', 'app', '界面', '截图', '手机', '功能', '玩法'])
                has_social = any(kw in page_content for kw in ['抖音', '小红书', '微博', '社交', '话题', '晒图', 'KOL', 'KOC'])
                has_people = any(kw in page_content for kw in ['人物', '用户', '年轻人', '场景'])
                if has_app:
                    left_desc = "realistic smartphone mockup showing app interface with detailed UI elements"
                elif has_social:
                    left_desc = "social media feed screenshot or phone showing social platform interface"
                elif has_people:
                    left_desc = "authentic lifestyle photography of young Chinese people in daily scenes"
                else:
                    left_desc = "high-quality thematic photograph or illustration matching the page topic"

                style_text = f"""PRECISE LAYOUT MATCH - 满版图片-左图右文 (Split layout):
- LEFT side (40-50%): {left_desc}
- RIGHT side (50-60%): Clean white background text area with structured content layout
- Top-left: brand-colored rectangular section tag
- Right area: Large brand-colored title + well-organized body text
- Body text should use bullet points, numbered lists, or small info cards for visual interest
- Clear visual separation between image and text areas
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

            elif '模块化-卡片' in page_layout:
                card_count = _guess_card_count(page_content)
                style_text = f"""PRECISE LAYOUT MATCH - 模块化-卡片 (Card module with dynamic columns):
- {card_count} equal-width cards arranged HORIZONTALLY in a row
- Each card: rounded corners (12px radius), subtle shadow, white/light background
- Card header: icon (line style) + title in brand primary color
- Card body: brief description in dark gray
- Cards separated by consistent spacing (24px gap)
- Top section: page title with brand color accent
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

            elif '模块化-步骤流程' in page_layout:
                style_text = f"""PRECISE LAYOUT MATCH - 模块化-步骤流程 (Step-by-step flow):
- Vertical or horizontal flow with STEP 1 / STEP 2 / STEP 3 markers
- Each step: numbered badge (brand primary color circle with white number) + title + description
- Connect steps with subtle arrows or lines
- Each step module has light background with subtle shadow
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

            elif '模块化-时间轴' in page_layout:
                style_text = f"""PRECISE LAYOUT MATCH - 模块化-时间轴 (Timeline):
- Horizontal timeline with nodes connected by line
- Each node: circular marker (brand primary color fill) + phase name + key actions
- Below timeline: detailed modules for each phase with images and text
- Clean, organized, easy to follow at a glance
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

            elif '模块化-表格' in page_layout:
                style_text = f"""PRECISE LAYOUT MATCH - 模块化-表格 (Data table):
- Multi-column table with clear headers
- Header row: brand primary color background with white text
- Data rows: alternating light backgrounds for readability
- Clean grid lines, professional alignment
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

            elif '纯视觉页' in page_layout:
                style_text = f"""PRECISE LAYOUT MATCH - 纯视觉页 (Visual only):
- NO text or minimal text only
- Full-screen key visual (KV) display
- Clean, impactful imagery
- Use brand colors from: {brand_colors if brand_colors else 'appropriate accent colors'}"""

            else:
                is_text_only = _is_text_only_layout(page_layout)

                if is_text_only:
                    _bc3 = brand_colors if brand_colors else 'elegant deep navy & warm gold accent'
                    _bg3 = "background: soft gradient from near-white (#FAFAFA) to very light brand-tinted tone, with subtle geometric pattern"
                    text_style_map = {
                        '文字-金句引言': f"""TEXT-ONLY LAYOUT - 金句引言 (Premium Quote Slide):\n- BACKGROUND: {_bg3}\n- QUOTE TEXT:超大字号(72-96pt), 品牌主色, 加粗, 居中或偏左黄金分割位置\n- DECORATIVE: 左侧或顶部有品牌色粗线条装饰(4-8px宽)或大型半透明引号装饰(200pt+, 10% opacity)\n- TITLE AREA: 页面顶部有小标签(badge), 圆角矩形, 品牌色底+白字\n- SPACING: 大量留白(40%+页面空白)\n- NO photographs — pure typography art\n- BRAND COLORS: {_bc3}\n- VIBE: 杂志封面级排版, Apple Keynote / Notion 风格极简高级感""",
                        '文字-数字大屏': f"""TEXT-ONLY LAYOUT - 数字大屏 (Data Dashboard - Premium):\n- BACKGROUND: {_bg3} + 微妙网格线(grid lines at 3% opacity)\n- NUMBERS: 超大数字(100-140pt), 品牌主色, Bold/Black, 每个在独立圆角卡片内(white bg, subtle shadow)\n- ACCENT: 趋势箭头↑↓或百分比指示器(品牌色)\n- NO photographs\n- BRAND COLORS: {_bc3}\n- VIBE: Bloomberg Terminal / 企业数据大屏风格""",
                        '文字-问答式': f"""TEXT-ONLY LAYOUT - 问答式 (Q&A - Premium Editorial):\n- BACKGROUND: {_bg3}, 左侧品牌色竖条装饰(8px宽)\n- Q: 品牌主色, 28-36pt Bold, 前面有装饰性「Q」符号\n- A: 深灰#333, 18-20pt, 行距1.6-1.8\n- NO photographs\n- BRAND COLORS: {_bc3}\n- VIBE: FAQ / 知乎精选风格""",
                        '文字-双栏对比': f"""TEXT-ONLY LAYOUT - 双栏对比 (Premium Split Design):\n- BACKGROUND: {_bg3}\n- 50/50分栏, 中间VS分隔区域(圆形徽章或竖线+箭头)\n- 两栏卡片式(headers in brand color, white bg, rounded corners)\n- NO photographs\n- BRAND COLORS: {_bc3}\n- VIBE: 产品功能对比风格, 清晰且有设计感""",
                        '文字-要点列表': f"""TEXT-ONLY LAYOUT - 要点列表 (Premium Structured):\n- BACKGROUND: {_bg3}, 极淡品牌色方向光效\n- TITLE: 品牌主色 36-48pt bold, 下方品牌色下划线装饰\n- LIST ITEMS: 大号序号圆圈(品牌色底白字, 32-40px diameter)\n- ITEM TEXT: 标题20-24pt Bold + 描述16-18pt regular, 行距1.5\n- NO photographs\n- BRAND COLORS: {_bc3}\n- VIBE: Apple Feature List / Notion Database 风格""",
                        '文字-引用来源': f"""TEXT-ONLY LAYOUT - 引用来源 (Premium Academic):\n- BACKGROUND: {_bg3}\n- TITLE: 「参考资料」, 品牌色 28-32pt bold\n- EACH ENTRY: 来源名(Bold 18pt 品牌色) + 摘要(15-16pt #444) + 年份(gray 13pt)\n- CONTAINER: 浅色圆角容器内(white bg, shadow, 16px radius, padding 40px)\n- NO photographs\n- BRAND COLORS: {_bc3}\n- VIBE: 学术参考文献风格, 专业严谨但美观"""
                    }
                    style_text = text_style_map.get(page_layout, f"TEXT-ONLY LAYOUT (Premium Text Design):\n- BACKGROUND: soft gradient from near-white to light brand-tinted tone, subtle geometric pattern\n- TITLE in brand primary color (36-48pt bold), body text in dark gray (#333, 16-18pt)\n- Strong typography hierarchy: clear size and weight contrast between heading/subheading/body\n- Generous whitespace (35-45% of slide area should be empty)\n- Decorative elements: brand-colored accent bars, thin divider lines, subtle geometric shapes\n- Content organized in clean cards or structured sections with rounded corners\n- NO photographs, but decorative graphic elements are encouraged\n- BRAND COLORS: {_bc3}\n- VIBE: Modern editorial / magazine layout quality — NOT a Word document")
                    bc = brand_colors if brand_colors else 'elegant dark accent color'
                    final_prompt = f"""Create a professional PPT slide.\n\nTitle: '{page_title}'\nContent: '{page_content}'\nLayout Type: '{page_layout}'\n\n{style_text}\n\nCRITICAL REQUIREMENTS - TEXT ONLY SLIDE:\n- 16:9 aspect ratio\n- This is a TEXT-ONLY slide — do NOT add any photos, illustrations, people, scenes, or decorative images\n- Focus entirely on typography, layout, spacing, and text hierarchy{product_constraint}"""
                else:
                    bc_text = brand_colors if brand_colors else 'Use appropriate brand colors'
                    _is_strong_img = any(k in page_layout for k in ['满版图片', '纯视觉页'])
                    _is_module = any(k in page_layout for k in ['模块化-卡片', '模块化-步骤', '模块化-时间轴', '模块化-表格'])

                    _scene_desc3 = _build_scene_description(page_title, page_content)

                    if _is_strong_img:
                        ve = f"\nKEY VISUAL ELEMENTS:\n1. SCENE (content-driven): {_scene_desc3}\n2. Brand colors: {bc_text}\n3. Typography: Clean sans-serif Chinese fonts, large expressive titles"
                        crit = f"CRITICAL REQUIREMENTS:\n- 16:9 aspect ratio\n- Use authentic photography as the main visual focal point{product_constraint}"
                    elif _is_module:
                        ve = f"\nKEY VISUAL ELEMENTS:\n1. Visual style: Clean graphic design with icons, illustrations — NOT lifestyle photography\n2. Card/module design: self-contained with icon + title + text structure\n3. Brand colors: {bc_text}"
                        crit = f"CRITICAL REQUIREMENTS:\n- 16:9 aspect ratio\n- This is a MODULE/CARD layout — focus on structured information display{product_constraint}"
                    else:
                        ve = f"\nKEY VISUAL ELEMENTS:\n1. Balanced mix of graphics and text based on content needs\n2. Brand colors: {bc_text}"
                        crit = f"CRITICAL REQUIREMENTS:\n- 16:9 aspect ratio\n- Clean, modern presentation design{product_constraint}"

                    final_prompt = f"""Create a professional PPT slide.\n\nTitle: '{page_title}'\nContent: '{page_content}'\nLayout Type: '{page_layout}'\n\n{style_text}\n\n{ve}\n\n{crit}"""

            if final_prompt is None:
                _ve = ""
                _cr = """CRITICAL REQUIREMENTS:\n- 16:9 aspect ratio\n- High quality, professional design"""
                if any(k in page_layout for k in ['满版图片', '纯视觉页']):
                    _ve = f"\nKEY VISUAL ELEMENTS:\n1. Photography as main visual element\n2. Brand colors: {(brand_colors or 'appropriate accent colors')}"
                elif any(k in page_layout for k in ['模块化-卡片', '模块化-步骤', '模块化-时间轴', '模块化-表格']):
                    _ve = f"\nKEY VISUAL ELEMENTS:\n1. Icons, illustrations, geometric shapes — NO lifestyle photos\n2. Brand colors: {(brand_colors or 'appropriate accent colors')}"
                final_prompt = f"""Create a professional PPT slide.\n\nTitle: '{page_title}'\nContent: '{page_content}'\nLayout Type: '{page_layout}'\n\n{style_text}\n\n{_ve}\n{_cr}{product_constraint if product_image else ''}"""

            yield f"data: {json.dumps({'status': 'progress', 'message': f'正在生成第 {page_index+1} 页图片...'}, ensure_ascii=False)}\n\n"

            try:
                _needs_prod = _page_needs_product_image(page_title, page_content, page_type, page_layout)

                if _needs_prod and product_image and product_image.startswith('data:image'):
                    messages = [{"role": "user", "content": [
                        {"type": "text", "text": final_prompt},
                        {"type": "image_url", "image_url": {"url": product_image}}
                    ]}]
                    result = api.image_generate(model=image_model, prompt=final_prompt, size='1792x1024', n=1, messages=messages)
                else:
                    result = api.image_generate(model=image_model, prompt=final_prompt, size='1792x1024', n=1)

                img_url = None
                _raw_images = result.get('images', []) if isinstance(result, dict) else []
                for img in _raw_images:
                    if isinstance(img, dict) and img.get('url'):
                        img_url = img['url']
                    elif isinstance(img, dict) and img.get('b64_json'):
                        try:
                            import base64
                            raw = base64.b64decode(img['b64_json'])
                            fname = f"{int(time.time())}_{page_index+1}_{uuid.uuid4().hex[:10]}.png"
                            fpath = os.path.join(GENERATED_DIR, fname)
                            with open(fpath, "wb") as f:
                                f.write(raw)
                            img_url = f"/api/generated/{fname}"
                        except Exception:
                            img_url = None

                if img_url:
                    if session_id and session_id in _uploaded_outline_cache:
                        store = _uploaded_outline_cache[session_id]
                        imgs = store.get('images', [])
                        while len(imgs) <= page_index:
                            imgs.append(f"/api/placeholder/800/450")
                        imgs[page_index] = img_url
                        store['images'] = imgs
                        _save_outline_session(session_id, store)
                    yield f"data: {json.dumps({'status': 'page_done', 'page_index': page_index, 'image': img_url, 'message': f'第 {page_index+1} 页生成成功'}, ensure_ascii=False)}\n\n"
                    yield f"data: {json.dumps({'status': 'done', 'page_index': page_index, 'image': img_url, 'message': f'第 {page_index+1} 页重新生成完成'}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {json.dumps({'status': 'error', 'message': f'第 {page_index+1} 页生成失败（无图片返回）'}, ensure_ascii=False)}\n\n"

            except Exception as e:
                import traceback
                traceback.print_exc()
                yield f"data: {json.dumps({'status': 'error', 'message': f'第 {page_index+1} 页失败: {str(e)[:120]}'}, ensure_ascii=False)}\n\n"

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield f"data: {json.dumps({'status': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'}
    )


@app.route('/api/split-page', methods=['POST'])
def ai_split_page():
    """AI 拆分单页为多页"""
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    data = request.get_json(silent=True) or {}
    page_idx = data.get('page_index', 0)
    title = data.get('title', '')
    content = data.get('content', '')
    layout = data.get('layout', '')
    page_type = data.get('type', 'content')

    if not content or len(content.strip()) < 50:
        return jsonify({'status': 'error', 'message': '内容太短，不足以拆分为多页'}), 400

    try:
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)

        api_url = config.get('api_base_url', config.get('chat_api_base_url'))
        api_key = config.get('api_key', config.get('chat_api_key'))
        chat_model = config.get('chat_model', 'gpt-4o')

        if not api_url or not api_key:
            return jsonify({'status': 'error', 'message': 'API未配置'}), 400

        api = APIClient(api_url, api_key)

        _markers = re.findall(r'【(\d+)\s*页】', content)
        _has_markers = len(_markers) > 0
        _force_count = sum(int(x) for x in _markers) if _has_markers else 0

        split_prompt = f"""你是一个专业的PPT大纲拆分专家。请将以下一页PPT的内容，根据其内在逻辑结构，智能拆分为多个独立的PPT页面。

## 原页面信息
- 标题：{title}
- 布局类型：{layout}
- 页面类型：{page_type}
- 原始内容：
```
{content[:8000]}
```

{'## 注意：原文中包含分页标记（如【2页】【3页】等），请严格按照标记数量拆分。' if _has_markers else ''}

## 拆分规则（必须严格遵守）
1. **绝对禁止**编造、改写、扩写或删减原文中的任何实际内容
2. **只做结构拆分**——把一大段内容按逻辑切分成多小段，每段成为独立一页
3. 每个拆分出的子页面必须有：
   - `title`: 子标题（从原内容中提取第一行/核心主题作为标题）
   - `content`: 该子页面的完整正文内容（直接取自原文，不修改措辞）
   - `type`: "content"（固定）
   - `layout`: 根据内容特征推断合适的布局（可选值：满版图片-全屏背景 / 满版图片-左图右文 / 模块化-卡片 / 模块化-步骤流程 / 模块化-时间轴 / 模块化-表格 / 纯视觉页 / 文字-要点列表 / 文字-双栏对比 等）

4. 如果原文中有【X页】标记，**必须**严格按标记数量拆分
5. 如果内容本身很短（不足200字），则不需要拆分，返回空数组

## 输出格式（严格的JSON数组）
```json
[
  {{
    "title": "子页面1的标题",
    "content": "子页面1的完整正文...",
    "type": "content",
    "layout": "推测的布局类型"
  }},
  {{
    "title": "子页面2的标题",
    "content": "子页面2的完整正文...",
    "type": "content",
    "layout": "推测的布局类型"
  }}
]
```

如果此页不需要拆分（内容太短或已经是单一主题），返回空数组 []。
只输出JSON，不要输出其他任何文字。"""

        result = api.chat(
            model=chat_model,
            messages=[{'role': 'user', 'content': split_prompt}],
            max_tokens=8000,
            timeout=120
        )

        raw_output = result.get('content', '').strip()
        print(f"[SPLIT] AI拆分结果长度: {len(raw_output)}")

        json_match = re.search(r'\[.*\]', raw_output, re.DOTALL)
        if not json_match:
            return jsonify({'status': 'error', 'message': 'AI返回格式异常，无法解析为页面列表'})

        import ast
        split_pages = ast.literal_eval(json_match.group(0))

        if not isinstance(split_pages, list):
            return jsonify({'status': 'error', 'message': 'AI返回的不是有效的页面列表'})

        validated = []
        for i, sp in enumerate(split_pages):
            if not isinstance(sp, dict):
                continue
            p = {
                'title': str(sp.get('title', f'{title} - 第{i+1}部分')),
                'content': str(sp.get('content', '')),
                'type': str(sp.get('type', 'content')),
                'layout': str(sp.get('layout', layout or '默认布局')),
                'brief': f"由「{title}」拆分而来",
            }
            if p['content'].strip():
                validated.append(p)

        print(f"[SPLIT] 原页面「{title}」 → 拆分为 {len(validated)} 页")
        return jsonify({
            'status': 'ok',
            'original_title': title,
            'pages': validated,
            'split_count': len(validated)
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/update-outline-session', methods=['POST'])
def update_outline_session():
    """更新已上传大纲的session数据"""
    data = request.get_json(silent=True) or {}
    session_id = data.get('session_id', '')
    outline = data.get('outline')

    if not session_id:
        return jsonify({'status': 'error', 'message': '缺少session_id'}), 400
    if not outline:
        return jsonify({'status': 'error', 'message': '缺少outline数据'}), 400

    try:
        store = None
        if session_id in _uploaded_outline_cache:
            store = _uploaded_outline_cache[session_id]
        else:
            loaded = _load_outline_session(session_id)
            if loaded:
                store = loaded
                _uploaded_outline_cache[session_id] = loaded

        if not store:
            return jsonify({'status': 'error', 'message': 'session不存在'}), 404

        store['outline'] = outline
        _save_outline_session(session_id, store)
        return jsonify({'status': 'ok', 'pages_count': len(outline.get('pages', []))})

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/page-chat', methods=['POST'])
def page_chat():
    """单页AI对话：针对某一页内容进行AI对话、改写、润色等"""
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    data = request.get_json(silent=True) or {}
    page_idx = data.get('page_index', 0)
    messages = data.get('messages', [])
    session_id = data.get('session_id', '')

    if not messages:
        return jsonify({'status': 'error', 'message': '消息不能为空'}), 400

    # 获取页面数据
    page_data = None
    if session_id and session_id in _uploaded_outline_cache:
        _store = _uploaded_outline_cache[session_id]
        if _store.get('outline') and _store['outline'].get('pages'):
            pages = _store['outline']['pages']
            if 0 <= page_idx < len(pages):
                page_data = pages[page_idx]

    if not page_data:
        return jsonify({'status': 'error', 'message': '找不到该页面数据'}), 404

    try:
        with open(CONFIG_FILE, 'r') as f:
            config = json.load(f)

        api_url = config.get('api_base_url', config.get('chat_api_base_url'))
        api_key = config.get('api_key', config.get('chat_api_key'))
        chat_model = config.get('chat_model', 'gpt-4o')

        if not api_url or not api_key:
            return jsonify({'status': 'error', 'message': 'API未配置'}), 400

        api = APIClient(api_url, api_key)

        _page_title = page_data.get('title', f'第{page_idx+1}页')
        _page_content = page_data.get('content', '')
        _page_brief = page_data.get('brief', '')
        _page_layout = page_data.get('layout', '')

        system_prompt = f"""你是一个专业的PPT内容编辑助手。用户正在编辑第 {page_idx+1} 页 PPT 的内容。

## 当前页面信息
- 标题: {_page_title}
- 布局类型: {_page_layout}
- 当前内容:
```
{_page_content[:3000]}
```
- 排版说明: {_page_brief or '(无)'}

## 你的能力
1. **改写/润色** — 让文案更专业、更有感染力、更简洁
2. **扩写** — 在保持原意的基础上丰富细节
3. **精简** — 删除冗余，提炼核心信息
4. **翻译** — 中英互译或其他语言
5. **风格调整** — 改变语气（正式/活泼/感性/理性）
6. **纠错** — 检查语法、逻辑、用词问题

## 规则
- 用户让你修改内容时，直接输出**修改后的完整文本**，不要只说改了什么
- 如果用户只是聊天询问，正常回答即可
- 输出的内容可以直接复制粘贴到 PPT 内容框中使用
- 保持原文的核心意思，不要擅自改变原意（除非用户明确要求）
- 回复语言与用户一致（用户用中文就用中文）"""

        result = api.chat(
            model=chat_model,
            messages=[
                {'role': 'system', 'content': system_prompt},
                *messages[-20:]
            ],
            max_tokens=4000,
            timeout=60
        )

        reply_text = result.get('content', '').strip()
        print(f"[PAGE_CHAT] 第{page_idx+1}页对话完成, 回复长度: {len(reply_text)}")

        return jsonify({
            'status': 'ok',
            'reply': reply_text,
            'page_index': page_idx
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/download-pdf', methods=['POST'])
def download_pdf():
    """将生成的所有图片合成为一个 PDF 文件"""
    try:
        import base64
        import requests
        from PIL import Image as PILImage

        data = request.get_json(silent=True) or {}
        images = data.get('images', []) or []
        outline = data.get('outline', {}) or {}
        pages = outline.get('pages', []) or []
        session_id = data.get('session_id', '')

        if session_id and session_id in _uploaded_outline_cache:
            store_images = _uploaded_outline_cache[session_id].get('images', [])
            if store_images and not images:
                images = store_images
        elif session_id:
            loaded = _load_outline_session(session_id)
            if loaded:
                store_images = loaded.get('images', [])
                if store_images and not images:
                    images = store_images

        if not images:
            return jsonify({'status': 'error', 'message': '没有图片可导出 PDF'}), 400

        temp_img_paths = []
        for i, img_url in enumerate(images):
            try:
                if not img_url or (isinstance(img_url, str) and img_url.startswith("/api/placeholder/")):
                    continue

                local_path = None

                if isinstance(img_url, str) and img_url.startswith("data:image"):
                    b64_data = img_url.split(",", 1)[1]
                    raw = base64.b64decode(b64_data)
                    local_path = os.path.join(_tmpdir, f'pdf_img_{i}.png')
                    with open(local_path, 'wb') as f:
                        f.write(raw)

                elif isinstance(img_url, str) and img_url.startswith("/api/generated/"):
                    fname = img_url.split("/api/generated/", 1)[1]
                    candidate = os.path.join(GENERATED_DIR, fname)
                    if os.path.exists(candidate):
                        local_path = candidate

                elif isinstance(img_url, str) and img_url.startswith("http"):
                    resp = requests.get(img_url, timeout=60)
                    resp.raise_for_status()
                    local_path = os.path.join(_tmpdir, f'pdf_img_{i}.png')
                    with open(local_path, 'wb') as f:
                        f.write(resp.content)

                elif isinstance(img_url, str) and os.path.exists(img_url):
                    local_path = img_url

                if local_path and os.path.exists(local_path):
                    temp_img_paths.append(local_path)

            except Exception as e:
                print(f"[WARN] PDF 导出时处理图片失败（已跳过）: {img_url} err={e}")
                continue

        if not temp_img_paths:
            return jsonify({'status': 'error', 'message': '没有有效的图片可导出 PDF'}), 400

        pdf_filename = f"pdf_{int(time.time())}_{uuid.uuid4().hex[:6]}.pdf"
        pdf_path = os.path.join(BASE_DIR, 'static', pdf_filename)
        os.makedirs(os.path.dirname(pdf_path), exist_ok=True)

        pil_images = []
        for img_path in temp_img_paths:
            try:
                img = PILImage.open(img_path)
                if img.mode == 'RGBA':
                    rgb_img = PILImage.new('RGB', img.size, (255, 255, 255))
                    rgb_img.paste(img, mask=img.split()[3])
                    pil_images.append(rgb_img)
                else:
                    img = img.convert('RGB')
                    pil_images.append(img)
            except Exception as e:
                print(f"[WARN] 打开图片失败（已跳过）: {img_path} err={e}")
                continue

        if not pil_images:
            return jsonify({'status': 'error', 'message': '没有有效的图片可写入 PDF'}), 400

        if len(pil_images) == 1:
            pil_images[0].save(pdf_path, format='PDF')
        else:
            pil_images[0].save(pdf_path, format='PDF', save_all=True, append_images=pil_images[1:])

        from flask import send_file
        pdf_title = (outline.get('title') or 'StyleMind_PPT')
        safe_title = re.sub(r'[^\w\-.]', '_', pdf_title)
        return send_file(
            pdf_path,
            as_attachment=True,
            download_name=f'{safe_title}.pdf',
            mimetype='application/pdf'
        )

    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/outlines', methods=['GET'])
def list_outlines():
    """列出所有历史上传的大纲"""
    try:
        sessions = _list_outline_sessions()
        return jsonify({
            'status': 'success',
            'count': len(sessions),
            'outlines': sessions
        })
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/outline/<session_id>', methods=['GET'])
def get_outline(session_id):
    """获取某个大纲的完整数据（含已生成的图片）"""
    try:
        # 先查内存缓存
        if session_id in _uploaded_outline_cache:
            data = _uploaded_outline_cache[session_id]
        else:
            data = _load_outline_session(session_id)

        if not data:
            return jsonify({'status': 'error', 'message': '大纲不存在或已过期'}), 404

        return jsonify({
            'status': 'success',
            'session_id': session_id,
            'outline': data.get('outline'),
            'images': data.get('images', []),
            'generated_images': data.get('images', []),  # 兼容前端字段名
            'created_at': data.get('created_at')
        })
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/outline/<session_id>', methods=['DELETE'])
def delete_outline(session_id):
    """删除某个大纲"""
    try:
        path = _get_outline_path(session_id)
        if os.path.exists(path):
            os.remove(path)
        if session_id in _uploaded_outline_cache:
            del _uploaded_outline_cache[session_id]
        return jsonify({'status': 'success', 'message': '已删除'})
    except Exception as e:
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/upload-product-image', methods=['POST'])
def upload_product_image():
    """上传产品参考图（图片或PDF），返回 base64 数据"""
    try:
        import base64
        if 'file' in request.files:
            f = request.files['file']
            img_data = f.read()
            filename = f.filename or ''
        elif request.data:
            img_data = request.data
            filename = ''
        else:
            return jsonify({'status': 'error', 'message': '没有收到文件'}), 400

        if filename.lower().endswith('.pdf') or (len(img_data) > 4 and img_data[:4] == b'%PDF'):
            import fitz
            doc = fitz.open(stream=img_data, filetype="pdf")
            page = doc[0]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_bytes = pix.tobytes("png")
            b64 = base64.b64encode(img_bytes).decode()
            doc.close()
        else:
            b64 = base64.b64encode(img_data).decode()

        product_image_url = f"data:image/png;base64,{b64}"

        return jsonify({
            'status': 'success',
            'product_image': product_image_url,
            'message': '产品图上传成功'
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'status': 'error', 'message': str(e)}), 500


@app.route('/api/edit-image', methods=['POST'])
def edit_single_image():
    """单页图片 AI 修改：用户输入修改指令 → AI 重绘该页图片"""
    import sys
    sys.path.insert(0, os.path.join(BASE_DIR, '..'))
    from core.api_client import APIClient

    data = request.get_json(silent=True) or {}

    session_id = data.get('session_id', '')
    page_index = data.get('page_index')
    edit_instruction = data.get('instruction', '')
    current_image_url = data.get('current_image', '')
    product_image = data.get('productImage', '')

    def generate():
        try:
            with open(CONFIG_FILE, 'r') as f:
                config = json.load(f)

            image_api_url = config.get('image_api_base_url') or config.get('api_base_url')
            image_api_key = config.get('image_api_key') or config.get('api_key')

            if not image_api_url or not image_api_key:
                yield f"data: {json.dumps({'status': 'error', 'message': '请先配置图片生成 API'}, ensure_ascii=False)}\n\n"
                return

            api = APIClient(image_api_url, image_api_key)
            chat_model = config.get('chat_model', 'gpt-4o')
            image_model = config.get('image_model', 'gpt-image-2')

            outline = None
            pages = []
            if session_id:
                if session_id in _uploaded_outline_cache:
                    outline = _uploaded_outline_cache[session_id].get('outline')
                else:
                    loaded = _load_outline_session(session_id)
                    if loaded:
                        outline = loaded.get('outline')
                        _uploaded_outline_cache[session_id] = loaded

            if not outline:
                yield f"data: {json.dumps({'status': 'error', 'message': '找不到大纲数据'}, ensure_ascii=False)}\n\n"
                return

            pages = outline.get('pages', [])
            if page_index < 0 or page_index >= len(pages):
                yield f"data: {json.dumps({'status': 'error', 'message': f'无效的页码: {page_index+1}'}, ensure_ascii=False)}\n\n"
                return

            page = pages[page_index]
            page_title = page.get('title', f'第{page_index+1}页')
            page_content = page.get('content', '')
            page_layout = page.get('layout', '')
            if '模块化-三列卡片' in page_layout:
                page_layout = page_layout.replace('模块化-三列卡片', '模块化-卡片')
            page_type = page.get('type', 'content')

            yield f"data: {json.dumps({'status': 'progress', 'message': f'正在分析第 {page_index+1} 页的修改需求...'}, ensure_ascii=False)}\n\n"

            refine_prompt = f"""你是一个专业的 PPT 视觉设计师。用户要对一张 PPT 页面进行修改。

【原页面信息】
- 标题: {page_title}
- 类型: {page_type}
- 布局: {page_layout}
- 内容: {page_content[:500]}

【用户的修改要求】
{edit_instruction}

请根据用户的修改要求，输出一个新的、完整的英文图片生成 Prompt。要求：
1. 保持原有的布局风格和排版结构
2. 只做用户要求的修改
3. 输出格式：直接输出 prompt 文本，不要其他说明
4. prompt 要详细、具体，包含画面描述、色彩、构图、文字排版等"""

            refine_result = api.chat(
                model=chat_model,
                messages=[{'role': 'user', 'content': refine_prompt}],
                max_tokens=1500,
                timeout=60
            )
            refined_prompt = refine_result.get('content', '').strip()

            if not refined_prompt:
                refined_prompt = edit_instruction

            yield f"data: {json.dumps({'status': 'progress', 'message': 'Prompt 已优化，正在重新生成图片...'}, ensure_ascii=False)}\n\n"

            final_prompt = refined_prompt

            if product_image:
                final_prompt += "\n\n【产品图约束】如果有产品出现，必须严格保持产品外形与参考图一致。"

            messages = [
                {"role": "user", "content": [{"type": "text", "text": final_prompt}]}
            ]

            if product_image and product_image.startswith('data:image'):
                messages[0]["content"].append({
                    "type": "image_url",
                    "image_url": {"url": product_image}
                })

            result = api.image_generate(
                model=image_model,
                prompt=final_prompt,
                messages=messages,
                size="1024x1024",
                quality="high",
                timeout=120
            )

            if result and result.get('data'):
                img = result['data'][0]
                if 'b64_json' in img:
                    img_url = f"data:image/png;base64,{img['b64_json']}"
                    yield f"data: {json.dumps({'status': 'page_done', 'page_index': page_index, 'image': img_url, 'message': f'第 {page_index+1} 页修改完成'}, ensure_ascii=False)}\n\n"

                    if session_id and session_id in _uploaded_outline_cache:
                        store_images = _uploaded_outline_cache[session_id].get('images', [])
                        while len(store_images) <= page_index:
                            store_images.append(None)
                        store_images[page_index] = img_url
                        _uploaded_outline_cache[session_id]['images'] = store_images
                        _save_outline_session(session_id, _uploaded_outline_cache[session_id])

                    yield f"data: {json.dumps({'status': 'done', 'message': '修改完成！'}, ensure_ascii=False)}\n\n"
                    return
                elif 'url' in img:
                    img_url = img['url']
                    yield f"data: {json.dumps({'status': 'page_done', 'page_index': page_index, 'image': img_url, 'message': f'第 {page_index+1} 页修改完成'}, ensure_ascii=False)}\n\n"

                    if session_id and session_id in _uploaded_outline_cache:
                        store_images = _uploaded_outline_cache[session_id].get('images', [])
                        while len(store_images) <= page_index:
                            store_images.append(None)
                        store_images[page_index] = img_url
                        _uploaded_outline_cache[session_id]['images'] = store_images
                        _save_outline_session(session_id, _uploaded_outline_cache[session_id])

                    yield f"data: {json.dumps({'status': 'done', 'message': '修改完成！'}, ensure_ascii=False)}\n\n"
                    return

            yield f"data: {json.dumps({'status': 'error', 'message': '图片生成失败，请重试'}, ensure_ascii=False)}\n\n"

        except Exception as e:
            import traceback
            traceback.print_exc()
            yield f"data: {json.dumps({'status': 'error', 'message': str(e)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype='text/event-stream',
        headers={'Cache-Control': 'no-cache', 'X-Accel-Buffering': 'no', 'Connection': 'keep-alive'}
    )


# ==================== 启动服务 ====================

if __name__ == '__main__':
    print("="*50)
    print("StyleMind Web UI 已启动!")
    print("访问地址: http://localhost:8080")
    print("="*50)
    app.run(host='0.0.0.0', port=8080, debug=False, threaded=True, processes=1)
