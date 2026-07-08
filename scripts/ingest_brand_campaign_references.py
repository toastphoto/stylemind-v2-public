#!/usr/bin/env python3
"""Build a local manifest for brand-campaign PPT/PDF reference decks."""

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import tempfile
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation
from pypdf import PdfReader


DEFAULT_SOURCE = Path(os.environ.get("STYLEMIND_REFERENCE_SOURCE", "reference_samples/brand_campaign_source"))
DEFAULT_OUTPUT = Path("reference_samples/brand_campaign_ingest")


def selected_pages(total: int) -> list[int]:
    """Return 1-based page numbers that cover cover/intro/body/end."""
    if total <= 0:
        return []
    candidates = [1, 2, 5, max(1, total // 2), max(1, total - 1), total]
    return sorted({p for p in candidates if 1 <= p <= total})


def clean_text(text: str, limit: int = 180) -> str:
    return " ".join((text or "").split())[:limit]


def pdf_fonts(reader: PdfReader, max_pages: int = 12) -> list[tuple[str, int]]:
    fonts: Counter[str] = Counter()
    for page in list(reader.pages)[:max_pages]:
        resources = page.get("/Resources") or {}
        page_fonts = resources.get("/Font") or {}
        for font_ref in page_fonts.values():
            try:
                font_obj = font_ref.get_object()
                base = str(font_obj.get("/BaseFont") or "")
                if base:
                    fonts[base.lstrip("/")] += 1
            except Exception:
                continue
    return fonts.most_common(12)


def inspect_pdf(path: Path) -> dict[str, Any]:
    reader = PdfReader(str(path))
    pages = len(reader.pages)
    sample_text = []
    for page_no in selected_pages(pages):
        try:
            text = clean_text(reader.pages[page_no - 1].extract_text() or "")
        except Exception:
            text = ""
        sample_text.append({"page": page_no, "text": text})
    return {
        "file": path.name,
        "kind": "pdf",
        "size_mb": round(path.stat().st_size / 1024 / 1024, 1),
        "pages": pages,
        "selected_pages": selected_pages(pages),
        "source_use": "visual_reference_and_text_extraction",
        "top_fonts": pdf_fonts(reader),
        "sample_text": sample_text,
    }


def inspect_pptx(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    shape_types: Counter[str] = Counter()
    font_names: Counter[str] = Counter()
    layout_names: Counter[str] = Counter()
    text_shapes = 0
    picture_shapes = 0
    group_shapes = 0
    table_shapes = 0
    freeform_shapes = 0
    sample_text = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        try:
            layout_names[slide.slide_layout.name] += 1
        except Exception:
            pass

        slide_text_parts = []
        for shape in slide.shapes:
            shape_types[str(shape.shape_type)] += 1
            shape_type = str(shape.shape_type)
            if "PICTURE" in shape_type:
                picture_shapes += 1
            if "GROUP" in shape_type:
                group_shapes += 1
            if "TABLE" in shape_type:
                table_shapes += 1
            if "FREEFORM" in shape_type:
                freeform_shapes += 1
            if getattr(shape, "has_text_frame", False):
                text_shapes += 1
                text = clean_text(shape.text, 120)
                if text:
                    slide_text_parts.append(text)
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        name = getattr(run.font, "name", None)
                        if name:
                            font_names[name] += 1

        if slide_idx in selected_pages(len(prs.slides)):
            sample_text.append({"page": slide_idx, "text": clean_text(" ".join(slide_text_parts))})

    return {
        "file": path.name,
        "kind": "pptx",
        "size_mb": round(path.stat().st_size / 1024 / 1024, 1),
        "pages": len(prs.slides),
        "selected_pages": selected_pages(len(prs.slides)),
        "source_use": "template_candidate_and_visual_reference",
        "text_shapes": text_shapes,
        "picture_shapes": picture_shapes,
        "group_shapes": group_shapes,
        "table_shapes": table_shapes,
        "freeform_shapes": freeform_shapes,
        "top_fonts": font_names.most_common(12),
        "top_layout_names": layout_names.most_common(12),
        "top_shape_types": shape_types.most_common(12),
        "sample_text": sample_text,
    }


def render_pdf_samples(pdf_path: Path, out_dir: Path, page_numbers: list[int]) -> list[dict[str, Any]]:
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        bundled = Path(os.environ.get("PDFTOPPM_BIN", ""))
        pdftoppm = str(bundled) if bundled.exists() else ""
    if not pdftoppm:
        return []

    out_dir.mkdir(parents=True, exist_ok=True)
    rendered = []
    stem = safe_stem(pdf_path)
    for page_no in page_numbers:
        prefix = out_dir / f"{stem}_p{page_no:03d}"
        cmd = [
            pdftoppm,
            "-jpeg",
            "-r",
            "130",
            "-f",
            str(page_no),
            "-l",
            str(page_no),
            str(pdf_path),
            str(prefix),
        ]
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        candidates = sorted(out_dir.glob(f"{prefix.name}-*.jpg"))
        if candidates:
            final = out_dir / f"{stem}_p{page_no:03d}.jpg"
            candidates[0].replace(final)
            rendered.append({"page": page_no, "image": str(final)})
    return rendered


def render_pptx_samples(pptx_path: Path, out_dir: Path, page_numbers: list[int]) -> list[dict[str, Any]]:
    soffice = shutil.which("soffice") or "/opt/homebrew/bin/soffice"
    if not Path(soffice).exists():
        return []

    with tempfile.TemporaryDirectory(prefix="stylemind_ref_pptx_") as tmp:
        tmp_dir = Path(tmp)
        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(tmp_dir),
                str(pptx_path),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        pdfs = sorted(tmp_dir.glob("*.pdf"))
        if not pdfs:
            return []
        return render_pdf_samples(pdfs[0], out_dir, page_numbers)


def safe_stem(path: Path) -> str:
    keep = []
    for ch in path.stem:
        if ch.isalnum() or ch in ("-", "_"):
            keep.append(ch)
        else:
            keep.append("_")
    return "".join(keep).strip("_")[:80] or "deck"


def build_manifest(source_dir: Path, output_dir: Path, render: bool) -> dict[str, Any]:
    files = sorted([p for p in source_dir.iterdir() if p.suffix.lower() in {".pptx", ".pdf"}])
    entries = []
    sample_dir = output_dir / "samples"

    for path in files:
        if path.suffix.lower() == ".pptx":
            entry = inspect_pptx(path)
            if render:
                entry["rendered_samples"] = render_pptx_samples(path, sample_dir, entry["selected_pages"])
        else:
            entry = inspect_pdf(path)
            if render:
                entry["rendered_samples"] = render_pdf_samples(path, sample_dir, entry["selected_pages"])
        entries.append(entry)

    total_pages = sum(item.get("pages", 0) for item in entries)
    pptx_pages = sum(item.get("pages", 0) for item in entries if item.get("kind") == "pptx")
    pdf_pages = sum(item.get("pages", 0) for item in entries if item.get("kind") == "pdf")
    pptx_fonts = Counter()
    for item in entries:
        if item.get("kind") == "pptx":
            pptx_fonts.update(dict(item.get("top_fonts") or []))

    return {
        "source_dir": str(source_dir),
        "generated_by": "scripts/ingest_brand_campaign_references.py",
        "deck_count": len(entries),
        "pptx_count": sum(1 for item in entries if item.get("kind") == "pptx"),
        "pdf_count": sum(1 for item in entries if item.get("kind") == "pdf"),
        "total_pages": total_pages,
        "pptx_pages": pptx_pages,
        "pdf_pages": pdf_pages,
        "top_pptx_fonts": pptx_fonts.most_common(16),
        "entries": entries,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source-dir", type=Path, default=DEFAULT_SOURCE)
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--render-samples", action="store_true")
    args = parser.parse_args()

    if not args.source_dir.exists():
        raise SystemExit(f"source dir not found: {args.source_dir}")

    args.output_dir.mkdir(parents=True, exist_ok=True)
    manifest = build_manifest(args.source_dir, args.output_dir, args.render_samples)
    manifest_path = args.output_dir / "reference_manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    print(json.dumps({
        "manifest": str(manifest_path),
        "deck_count": manifest["deck_count"],
        "pptx_count": manifest["pptx_count"],
        "pdf_count": manifest["pdf_count"],
        "total_pages": manifest["total_pages"],
        "pptx_pages": manifest["pptx_pages"],
        "pdf_pages": manifest["pdf_pages"],
    }, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
