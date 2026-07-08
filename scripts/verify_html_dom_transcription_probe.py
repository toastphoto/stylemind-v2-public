#!/usr/bin/env python3
"""Verify the StyleMind HTML DOM -> editable PPTX transcription probe."""

from __future__ import annotations

import json
import subprocess
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / ".pytest_tmp" / "html_dom_transcription_probe"
PPTX_OUT = OUT_DIR / "stylemind_html_dom_probe.pptx"
REPORT_OUT = OUT_DIR / "stylemind_html_dom_probe.report.json"
NODE_SCRIPT = ROOT / "scripts" / "run_html_dom_transcription_probe.mjs"


def pptx_xml_counts(path: Path) -> dict:
    with zipfile.ZipFile(path) as zf:
        slide_names = sorted(
            name
            for name in zf.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        root = ET.fromstring(zf.read("ppt/presentation.xml"))
        ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        slide_ids = root.findall(".//p:sldId", ns)
        text_nodes = 0
        picture_tags = 0
        shape_tags = 0
        for slide_name in slide_names:
            slide_root = ET.fromstring(zf.read(slide_name))
            text_nodes += len(slide_root.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t"))
            picture_tags += len(slide_root.findall(".//{http://schemas.openxmlformats.org/presentationml/2006/main}pic"))
            shape_tags += len(slide_root.findall(".//{http://schemas.openxmlformats.org/presentationml/2006/main}sp"))
    return {
        "slide_xml": len(slide_names),
        "presentation_slide_ids": len(slide_ids),
        "text_nodes": text_nodes,
        "picture_tags": picture_tags,
        "shape_tags": shape_tags,
    }


def inspect_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    text_shapes = 0
    picture_shapes = 0
    auto_shapes = 0
    full_slide_pictures = 0
    sample_text: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
                if len(sample_text) < 8:
                    sample_text.append(shape.text.strip()[:100])
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
                if shape.width >= prs.slide_width * 0.9 and shape.height >= prs.slide_height * 0.9:
                    full_slide_pictures += 1
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                auto_shapes += 1
    return {
        "slides": len(prs.slides),
        "text_shapes": text_shapes,
        "picture_shapes": picture_shapes,
        "auto_shapes": auto_shapes,
        "full_slide_pictures": full_slide_pictures,
        "sample_text": sample_text,
        "size_kb": round(path.stat().st_size / 1024, 1),
    }


def main() -> int:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    completed = subprocess.run(
        [
            "node",
            str(NODE_SCRIPT),
            "--out-dir",
            str(OUT_DIR),
            "--output",
            str(PPTX_OUT),
            "--report",
            str(REPORT_OUT),
        ],
        cwd=str(ROOT),
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=120,
    )

    if not PPTX_OUT.exists() or PPTX_OUT.stat().st_size == 0:
        raise SystemExit("HTML DOM probe did not create a PPTX")
    if not REPORT_OUT.exists():
        raise SystemExit("HTML DOM probe did not create an export report")

    stats = inspect_pptx(PPTX_OUT)
    xml_counts = pptx_xml_counts(PPTX_OUT)
    report = json.loads(REPORT_OUT.read_text(encoding="utf-8"))

    if stats["slides"] != 3:
        raise SystemExit(f"expected 3 slides, got {stats['slides']}")
    if stats["text_shapes"] < 12 or xml_counts["text_nodes"] < 12:
        raise SystemExit(f"expected editable text objects, got {stats['text_shapes']} shapes / {xml_counts['text_nodes']} XML text nodes")
    if stats["picture_shapes"] < 1 or xml_counts["picture_tags"] < 1:
        raise SystemExit(f"expected screenshot/image fallback picture objects, got {stats['picture_shapes']}")
    if report.get("slideCount") != 3:
        raise SystemExit(f"expected report slideCount=3, got {report.get('slideCount')}")
    if report.get("textObjects", 0) < 12:
        raise SystemExit(f"expected report textObjects >= 12, got {report.get('textObjects')}")
    if not any("StyleMind" in text or "美的" in text for text in stats["sample_text"]):
        raise SystemExit(f"expected StyleMind/Midea text in PPTX, got samples: {stats['sample_text']}")

    print(
        json.dumps(
            {
                "status": "ok",
                "renderer": "html-dom",
                "output": str(PPTX_OUT),
                "report": str(REPORT_OUT),
                "node_stdout": completed.stdout.strip(),
                **stats,
                **xml_counts,
                "export_report": {
                    "slideCount": report.get("slideCount"),
                    "textObjects": report.get("textObjects"),
                    "shapeObjects": report.get("shapeObjects"),
                    "imageObjects": report.get("imageObjects"),
                    "warnings": len(report.get("warnings") or []),
                },
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
