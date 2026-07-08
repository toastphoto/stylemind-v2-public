#!/usr/bin/env python3
"""Verify a real DOCX outline through StyleMind RenderPlan and both PPTX renderers."""

from __future__ import annotations

import argparse
import json
import os
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
WEB_UI = ROOT / "web_ui"
if str(WEB_UI) not in sys.path:
    sys.path.insert(0, str(WEB_UI))

from api_server import app  # noqa: E402
from services.renderers.plan_preview_renderer import render_plan_preview_png  # noqa: E402
from services.renderers.pptxgenjs_renderer import render_payload_to_pptxgenjs  # noqa: E402
from services.renderers.python_pptx_renderer import render_plans_to_pptx  # noqa: E402
from services.slide_render_plan import build_deck_render_plan, render_plan_to_dict  # noqa: E402


DEFAULT_DOCX = Path(os.environ.get("STYLEMIND_DOCX_FIXTURE", "fixtures/sample_outline.docx"))
TMP = ROOT / ".pytest_tmp"
OUTLINE_CACHE = TMP / "midea_social_outline_upload_response.json"
PLAN_JSON = TMP / "midea_social_render_plan.json"
PYTHON_PPTX = TMP / "midea_social_python_pptx.pptx"
PPTXGENJS_PPTX = TMP / "midea_social_pptxgenjs_spike.pptx"
INTERNAL_ARCHETYPE_IDS = {
    "hero_photo_claim",
    "section_divider",
    "strategy_claim_collage",
    "metric_dashboard",
    "evidence_wall",
    "campaign_timeline",
    "step_flow",
    "editorial_content_bridge",
    "video_material_board",
}
IMAGE_ARCHETYPES = {"hero_photo_claim", "strategy_claim_collage", "evidence_wall", "video_material_board"}
NATIVE_SHAPE_ARCHETYPES = {"metric_dashboard", "campaign_timeline", "step_flow", "section_divider"}


def upload_docx(docx_path: Path) -> dict:
    with app.test_client() as client:
        with docx_path.open("rb") as f:
            response = client.post(
                "/api/upload-outline",
                data={"file": (f, docx_path.name)},
                content_type="multipart/form-data",
            )
    body = response.get_json(silent=True) or {}
    if response.status_code != 200 or body.get("status") != "success":
        raise RuntimeError(f"upload-outline failed: {response.status_code} {body}")
    return body


def load_or_upload_outline(docx_path: Path, refresh: bool) -> dict:
    TMP.mkdir(parents=True, exist_ok=True)
    if OUTLINE_CACHE.exists() and not refresh:
        cached = json.loads(OUTLINE_CACHE.read_text(encoding="utf-8"))
        if cached.get("source_docx") == str(docx_path) and cached.get("source_mtime") == docx_path.stat().st_mtime:
            return cached["upload_response"]

    body = upload_docx(docx_path)
    OUTLINE_CACHE.write_text(
        json.dumps(
            {
                "source_docx": str(docx_path),
                "source_mtime": docx_path.stat().st_mtime,
                "upload_response": body,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return body


def inspect_pptx(path: Path, plan_payload: dict | None = None) -> dict:
    prs = Presentation(str(path))
    text_shapes = 0
    picture_shapes = 0
    auto_shapes = 0
    full_slide_pictures = 0
    sample_text = []
    visible_internal_ids = []
    archetype_stats: dict[str, dict[str, int]] = {}
    plan_items = plan_payload.get("plans", []) if plan_payload else []
    for slide_idx, slide in enumerate(prs.slides):
        archetype = ""
        if slide_idx < len(plan_items):
            archetype = (plan_items[slide_idx].get("visual_profile") or {}).get("archetype") or ""
            if archetype:
                archetype_stats.setdefault(archetype, {"slides": 0, "text_shapes": 0, "picture_shapes": 0, "auto_shapes": 0})
                archetype_stats[archetype]["slides"] += 1
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
                if archetype:
                    archetype_stats[archetype]["text_shapes"] += 1
                if len(sample_text) < 8:
                    sample_text.append(shape.text.strip()[:90])
                for internal_id in INTERNAL_ARCHETYPE_IDS:
                    if internal_id in shape.text:
                        visible_internal_ids.append({"slide": slide_idx + 1, "id": internal_id})
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
                if archetype:
                    archetype_stats[archetype]["picture_shapes"] += 1
                if shape.width >= prs.slide_width * 0.9 and shape.height >= prs.slide_height * 0.9:
                    full_slide_pictures += 1
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                auto_shapes += 1
                if archetype:
                    archetype_stats[archetype]["auto_shapes"] += 1
    return {
        "slides": len(prs.slides),
        "text_shapes": text_shapes,
        "picture_shapes": picture_shapes,
        "auto_shapes": auto_shapes,
        "full_slide_pictures": full_slide_pictures,
        "visible_internal_ids": visible_internal_ids,
        "archetype_stats": archetype_stats,
        "size_kb": round(path.stat().st_size / 1024, 1),
        "sample_text": sample_text,
    }


def run_node_renderer(plan_payload: dict) -> None:
    render_payload_to_pptxgenjs(plan_payload, PPTXGENJS_PPTX, render_plan_path=PLAN_JSON)


def validate_stats(name: str, stats: dict, page_count: int) -> None:
    if stats["slides"] != page_count:
        raise SystemExit(f"{name}: expected {page_count} slides, got {stats['slides']}")
    if stats["text_shapes"] < page_count:
        raise SystemExit(f"{name}: expected at least {page_count} text shapes, got {stats['text_shapes']}")
    if stats["full_slide_pictures"]:
        raise SystemExit(f"{name}: unexpected full-slide pictures: {stats['full_slide_pictures']}")
    if stats.get("visible_internal_ids"):
        raise SystemExit(f"{name}: internal visual profile IDs leaked into visible text: {stats['visible_internal_ids']}")


def validate_archetype_stats(name: str, stats: dict, visual_distribution: dict) -> None:
    archetype_stats = stats.get("archetype_stats") or {}
    missing = [key for key in visual_distribution if key not in archetype_stats]
    if missing:
        raise SystemExit(f"{name}: missing archetype stats for {missing}")

    for archetype, count in visual_distribution.items():
        current = archetype_stats[archetype]
        if current["text_shapes"] < count:
            raise SystemExit(f"{name}: {archetype} expected native text on {count} slides, got {current}")
        if current["auto_shapes"] < count:
            raise SystemExit(f"{name}: {archetype} expected native shapes on {count} slides, got {current}")
        if archetype in IMAGE_ARCHETYPES and current["picture_shapes"] < count:
            raise SystemExit(f"{name}: {archetype} expected replaceable picture slots on {count} slides, got {current}")
        if archetype in NATIVE_SHAPE_ARCHETYPES and current["auto_shapes"] < count * 2:
            raise SystemExit(f"{name}: {archetype} expected richer native visual objects, got {current}")


def role_distribution(plan_payload: dict) -> dict:
    roles = [plan["intent"]["page_role"] for plan in plan_payload.get("plans", [])]
    return dict(Counter(roles).most_common())


def visual_archetype_distribution(plan_payload: dict) -> dict:
    archetypes = [plan["visual_profile"]["archetype"] for plan in plan_payload.get("plans", []) if plan.get("visual_profile")]
    return dict(Counter(archetypes).most_common())


def validate_role_distribution(distribution: dict, page_count: int) -> None:
    if len(distribution) < 5:
        raise SystemExit(f"expected at least 5 page roles, got {distribution}")
    default_count = distribution.get("内容承接", 0)
    if default_count > page_count * 0.7:
        raise SystemExit(f"too many pages fell back to 内容承接: {distribution}")


def inspect_render_plan_previews(plans: list) -> dict:
    total_bytes = 0
    sample_roles = []
    for plan in plans:
        png = render_plan_preview_png(plan)
        if not png.startswith(b"\x89PNG"):
            raise SystemExit(f"render-plan preview for page {plan.index} is not a PNG")
        total_bytes += len(png)
        if len(sample_roles) < 8:
            sample_roles.append(plan.intent.page_role)
    return {
        "pages": len(plans),
        "format": "png",
        "source": "shared_render_plan",
        "uses_full_page_image_model": False,
        "total_size_kb": round(total_bytes / 1024, 1),
        "sample_roles": sample_roles,
    }


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--docx", type=Path, default=DEFAULT_DOCX)
    parser.add_argument("--refresh-outline", action="store_true", help="Re-run upload parsing instead of using the cached outline.")
    args = parser.parse_args()

    if not args.docx.exists():
        raise SystemExit(f"DOCX not found: {args.docx}")

    body = load_or_upload_outline(args.docx, args.refresh_outline)
    outline = body["outline"]
    page_count = int(body["page_count"])
    if page_count < 40:
        raise SystemExit(f"expected a substantial outline from DOCX, got {page_count} pages")

    plans = build_deck_render_plan(outline)
    plan_payload = {
        "schema": "stylemind.render_plan.v1",
        "page_count": len(plans),
        "plans": [render_plan_to_dict(plan) for plan in plans],
    }
    roles = role_distribution(plan_payload)
    visual_archetypes = visual_archetype_distribution(plan_payload)
    validate_role_distribution(roles, page_count)
    PLAN_JSON.write_text(json.dumps(plan_payload, ensure_ascii=False, indent=2), encoding="utf-8")

    preview_stats = inspect_render_plan_previews(plans)
    render_plans_to_pptx(plans, str(PYTHON_PPTX))
    run_node_renderer(plan_payload)

    python_stats = inspect_pptx(PYTHON_PPTX, plan_payload)
    pptxgenjs_stats = inspect_pptx(PPTXGENJS_PPTX, plan_payload)
    validate_stats("python-pptx", python_stats, page_count)
    validate_stats("pptxgenjs", pptxgenjs_stats, page_count)
    validate_archetype_stats("python-pptx", python_stats, visual_archetypes)
    validate_archetype_stats("pptxgenjs", pptxgenjs_stats, visual_archetypes)
    if pptxgenjs_stats["picture_shapes"] < 1:
        raise SystemExit("pptxgenjs: expected at least one replaceable picture object")

    print(
        json.dumps(
            {
                "status": "ok",
                "source_docx": str(args.docx),
                "page_count": page_count,
                "outline_title": outline.get("title"),
                "role_distribution": roles,
                "visual_archetype_distribution": visual_archetypes,
                "render_plan": str(PLAN_JSON),
                "render_plan_preview": preview_stats,
                "python_pptx": {"path": str(PYTHON_PPTX), **python_stats},
                "pptxgenjs": {"path": str(PPTXGENJS_PPTX), **pptxgenjs_stats},
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
