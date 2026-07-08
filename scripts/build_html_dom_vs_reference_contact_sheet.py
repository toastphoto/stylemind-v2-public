#!/usr/bin/env python3
"""Build a visual contact sheet comparing html-dom and reference-template PPTX output."""

from __future__ import annotations

import argparse
import json
import sys
from collections import OrderedDict
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
WEB_UI = ROOT / "web_ui"
SCRIPTS = ROOT / "scripts"
for path in (WEB_UI, SCRIPTS):
    if str(path) not in sys.path:
        sys.path.insert(0, str(path))

from build_pptx_render_review_contact_sheet import convert_pptx_to_pdf, render_pdf_page  # noqa: E402
from build_visual_review_contact_sheet import draw_section, draw_tile, open_reference_samples  # noqa: E402
from services.renderers.html_dom_renderer import render_payload_to_html_dom  # noqa: E402
from services.renderers.reference_template_renderer import render_payload_to_reference_template  # noqa: E402


PLAN_JSON = ROOT / ".pytest_tmp/midea_social_render_plan.json"
OUT_DIR = ROOT / ".pytest_tmp/html_dom_vs_reference_review"
SELECTED_PLAN_JSON = OUT_DIR / "selected_render_plan.json"
HTML_DOM_PPTX = OUT_DIR / "stylemind_html_dom_selected.pptx"
HTML_DOM_REPORT = OUT_DIR / "stylemind_html_dom_selected.report.json"
HTML_DOM_DECK_DIR = OUT_DIR / "html_dom_deck"
REFERENCE_TEMPLATE_PPTX = OUT_DIR / "stylemind_reference_template_selected.pptx"
REFERENCE_TEMPLATE_REPORT = OUT_DIR / "stylemind_reference_template_selected.report.json"
DEFAULT_OUTPUT = OUT_DIR / "html_dom_vs_reference_contact_sheet.jpg"


def load_plan_payload(path: Path) -> dict:
    if not path.exists():
        raise SystemExit(f"RenderPlan JSON not found: {path}. Run scripts/verify_docx_outline_render_chain.py first.")
    payload = json.loads(path.read_text(encoding="utf-8"))
    if payload.get("schema") != "stylemind.render_plan.v1":
        raise SystemExit(f"Unsupported RenderPlan schema: {payload.get('schema')}")
    if not isinstance(payload.get("plans"), list) or not payload["plans"]:
        raise SystemExit(f"No plans found in {path}")
    return payload


def select_one_per_archetype(payload: dict, max_pages: int) -> tuple[dict, list[dict]]:
    grouped: OrderedDict[str, dict] = OrderedDict()
    for plan in payload.get("plans", []):
        visual = plan.get("visual_profile") or {}
        archetype = visual.get("archetype") or "unknown"
        grouped.setdefault(archetype, plan)
        if len(grouped) >= max_pages:
            break

    selected_plans = list(grouped.values())
    selected_payload = {
        "schema": "stylemind.render_plan.v1",
        "page_count": len(selected_plans),
        "plans": selected_plans,
    }
    records = []
    for rendered_slide, plan in enumerate(selected_plans, start=1):
        intent = plan.get("intent") or {}
        visual = plan.get("visual_profile") or {}
        records.append(
            {
                "rendered_slide": rendered_slide,
                "original_page": int(plan.get("index") or rendered_slide),
                "role": intent.get("page_role") or "",
                "archetype": visual.get("archetype") or "unknown",
                "title": plan.get("title") or "",
            }
        )
    return selected_payload, records


def inspect_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    text_shapes = 0
    picture_shapes = 0
    auto_shapes = 0
    full_slide_pictures = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
                if shape.width >= prs.slide_width * 0.9 and shape.height >= prs.slide_height * 0.9:
                    full_slide_pictures += 1
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                auto_shapes += 1
    return {
        "slides": len(prs.slides),
        "text_shapes": text_shapes,
        "picture_shapes": picture_shapes,
        "auto_shapes": auto_shapes,
        "full_slide_pictures": full_slide_pictures,
        "size_kb": round(path.stat().st_size / 1024, 1),
    }


def render_pptx_section(renderer_label: str, pptx_path: Path, records: list[dict]) -> list[dict]:
    pdf_path = convert_pptx_to_pdf(pptx_path)
    items = []
    safe_name = renderer_label.lower().replace(" ", "_").replace("-", "_")
    for record in records:
        slide = record["rendered_slide"]
        image_path = render_pdf_page(pdf_path, slide, OUT_DIR / "rendered_pages" / safe_name / f"{safe_name}_s{slide:03d}")
        items.append(
            {
                "kind": safe_name,
                "label": f"{renderer_label} S{slide:02d} / P{record['original_page']:02d} {record['role']} / {record['archetype']}",
                "image": Image.open(image_path),
            }
        )
    return items


def build_sheet(sections: list[tuple[str, list[dict]]], output: Path) -> dict:
    tile_w, tile_h = 320, 220
    margin = 28
    gap = 18
    cols = 3
    section_h = 48
    rows = [max(1, (len(items) + cols - 1) // cols) for _, items in sections]
    width = margin * 2 + cols * tile_w + (cols - 1) * gap
    height = margin * 2 + sum(section_h + row_count * tile_h + row_count * gap for row_count in rows)
    canvas = Image.new("RGB", (width, height), (245, 247, 251))
    draw = ImageDraw.Draw(canvas)

    y = margin
    for (title, items), row_count in zip(sections, rows):
        draw_section(draw, margin, y, title)
        y += section_h
        for idx, item in enumerate(items):
            row, col = divmod(idx, cols)
            draw_tile(draw, canvas, item, margin + col * (tile_w + gap), y + row * (tile_h + gap), tile_w, tile_h)
        y += row_count * (tile_h + gap)

    output.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(output, quality=92)
    return {"output": str(output), "width": width, "height": height}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=PLAN_JSON)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--max-pages", type=int, default=9)
    parser.add_argument("--max-reference", type=int, default=9)
    args = parser.parse_args()

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    payload = load_plan_payload(args.input)
    selected_payload, records = select_one_per_archetype(payload, args.max_pages)
    SELECTED_PLAN_JSON.write_text(json.dumps(selected_payload, ensure_ascii=False, indent=2), encoding="utf-8")

    html_dom_result = render_payload_to_html_dom(
        selected_payload,
        HTML_DOM_PPTX,
        render_plan_path=OUT_DIR / "html_dom_render_plan.json",
        report_path=HTML_DOM_REPORT,
        deck_dir=HTML_DOM_DECK_DIR,
        timeout=240,
    )
    reference_result = render_payload_to_reference_template(
        selected_payload,
        REFERENCE_TEMPLATE_PPTX,
        render_plan_path=OUT_DIR / "reference_template_render_plan.json",
        report_path=REFERENCE_TEMPLATE_REPORT,
        timeout=300,
    )

    reference_samples = open_reference_samples(args.max_reference)
    html_dom_items = render_pptx_section("HTML DOM", HTML_DOM_PPTX, records)
    reference_template_items = render_pptx_section("Reference Template", REFERENCE_TEMPLATE_PPTX, records)
    sheet_stats = build_sheet(
        [
            ("Original Feibo reference samples", reference_samples),
            ("Same selected pages: HTML DOM transcription", html_dom_items),
            ("Same selected pages: reference-template transcription", reference_template_items),
        ],
        args.output,
    )

    html_dom_report = json.loads(Path(html_dom_result.report_path).read_text(encoding="utf-8"))
    reference_report = json.loads(Path(reference_result.report_path).read_text(encoding="utf-8"))
    summary = {
        "status": "ok",
        **sheet_stats,
        "selected_render_plan": str(SELECTED_PLAN_JSON),
        "selected_pages": records,
        "html_dom": {
            "pptx": str(HTML_DOM_PPTX),
            "report": str(HTML_DOM_REPORT),
            "deck_dir": str(HTML_DOM_DECK_DIR),
            **inspect_pptx(HTML_DOM_PPTX),
            "export_report": {
                "slideCount": html_dom_report.get("slideCount"),
                "textObjects": html_dom_report.get("textObjects"),
                "shapeObjects": html_dom_report.get("shapeObjects"),
                "imageObjects": html_dom_report.get("imageObjects"),
                "warnings": len(html_dom_report.get("warnings") or []),
            },
        },
        "reference_template": {
            "pptx": str(REFERENCE_TEMPLATE_PPTX),
            "report": str(REFERENCE_TEMPLATE_REPORT),
            **inspect_pptx(REFERENCE_TEMPLATE_PPTX),
            "selected_count": len(reference_report.get("selected") or []),
            "picture_replacement": reference_report.get("pictureReplacementSummary") or {},
        },
    }
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
