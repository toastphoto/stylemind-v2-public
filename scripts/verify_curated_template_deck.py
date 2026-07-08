#!/usr/bin/env python3
"""Verify the curated named-placeholder template deck for pptx-automizer."""

from __future__ import annotations

import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "reference_samples" / "curated_templates" / "stylemind_named_placeholder_templates.pptx"
PROBE_OUTPUT = ROOT / ".pytest_tmp" / "curated_template_probe" / "stylemind_curated_template_probe.pptx"
PROBE_IMAGE = ROOT / ".pytest_tmp" / "curated_template_probe" / "curated-probe-image.png"
EXPECTED_NAMES = {
    "title",
    "body",
    "hero_image",
    "metric_1_label",
    "metric_1_value",
    "metric_2_label",
    "metric_2_value",
    "metric_3_label",
    "metric_3_value",
    "metric_4_label",
    "metric_4_value",
    "evidence_image_1",
    "evidence_image_2",
    "evidence_image_3",
    "evidence_image_4",
}


def build_template() -> None:
    subprocess.run(["node", "scripts/build_curated_template_deck.mjs"], cwd=ROOT, check=True)
    if not TEMPLATE.exists():
        raise SystemExit(f"template not found: {TEMPLATE}")


def inspect_template() -> dict:
    prs = Presentation(str(TEMPLATE))
    names = set()
    text_shapes = 0
    picture_shapes = 0
    sample_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name:
                names.add(shape.name)
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
                if len(sample_text) < 10:
                    sample_text.append(shape.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1

    missing = sorted(EXPECTED_NAMES - names)
    if len(prs.slides) != 3:
        raise SystemExit(f"expected 3 curated template slides, got {len(prs.slides)}")
    if missing:
        raise SystemExit(f"missing named placeholders: {missing}")
    if picture_shapes < 5:
        raise SystemExit(f"expected at least 5 image placeholders, got {picture_shapes}")

    return {
        "slides": len(prs.slides),
        "text_shapes": text_shapes,
        "picture_shapes": picture_shapes,
        "placeholder_count": len(names & EXPECTED_NAMES),
        "sample_text": sample_text,
    }


def run_automizer_probe() -> dict:
    result = subprocess.run(
        ["node", "scripts/run_curated_template_automizer_probe.mjs"],
        cwd=ROOT,
        check=True,
        stdout=subprocess.PIPE,
        text=True,
    )
    payload = json.loads(result.stdout)
    if not PROBE_OUTPUT.exists():
        raise SystemExit(f"curated automizer output not found: {PROBE_OUTPUT}")
    if not PROBE_IMAGE.exists():
        raise SystemExit(f"curated automizer probe image not found: {PROBE_IMAGE}")

    prs = Presentation(str(PROBE_OUTPUT))
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    probe_blob = PROBE_IMAGE.read_bytes()
    probe_image_found = any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.image.blob == probe_blob
        for slide in prs.slides
        for shape in slide.shapes
    )
    if payload["probeText"]["title"] not in text:
        raise SystemExit("curated automizer output missing replaced title")
    if payload["probeText"]["body"] not in text:
        raise SystemExit("curated automizer output missing replaced body")
    if not probe_image_found:
        raise SystemExit("curated automizer output missing replaced hero image")

    return {
        "output": str(PROBE_OUTPUT),
        "slides": len(prs.slides),
        "replaced_title_found": payload["probeText"]["title"] in text,
        "replaced_body_found": payload["probeText"]["body"] in text,
        "probe_image_found": probe_image_found,
    }


def main() -> int:
    build_template()
    stats = inspect_template()
    automizer_probe = run_automizer_probe()
    print(json.dumps({"status": "ok", "template": str(TEMPLATE), **stats, "automizer_probe": automizer_probe}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
