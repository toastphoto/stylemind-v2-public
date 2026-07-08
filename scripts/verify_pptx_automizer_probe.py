#!/usr/bin/env python3
"""Verify pptx-automizer as an isolated template-reuse probe."""

from __future__ import annotations

import json
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / ".pytest_tmp" / "automizer_probe"
REPORT = OUT_DIR / "pptx_automizer_probe_report.json"


def run_node_probe() -> dict:
    subprocess.run(
        ["node", "scripts/run_pptx_automizer_probe.mjs"],
        cwd=ROOT,
        check=True,
    )
    if not REPORT.exists():
        raise SystemExit(f"automizer report not found: {REPORT}")
    return json.loads(REPORT.read_text(encoding="utf-8"))


def inspect_output(report: dict) -> dict:
    output = Path(report["probe"]["output"])
    probe_image = Path(report["probe"]["probeImage"])
    probe_text = report["probe"]["probeText"]
    if not output.exists():
        raise SystemExit(f"automizer output not found: {output}")
    if not probe_image.exists():
        raise SystemExit(f"probe image not found: {probe_image}")

    prs = Presentation(str(output))
    text_shapes = 0
    picture_shapes = 0
    replaced_text_found = False
    probe_image_blob = probe_image.read_bytes()
    probe_image_found = False
    sample_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
                if len(sample_text) < 8:
                    sample_text.append(shape.text.strip()[:120])
                if probe_text.splitlines()[0] in shape.text:
                    replaced_text_found = True
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
                if shape.image.blob == probe_image_blob:
                    probe_image_found = True

    if len(prs.slides) != 1:
        raise SystemExit(f"expected 1 automizer output slide, got {len(prs.slides)}")
    if not replaced_text_found:
        raise SystemExit("automizer output did not contain the replaced probe text")
    if not probe_image_found:
        raise SystemExit("automizer output did not contain the replaced probe image")
    if report["selected"]["placeholderQuality"] != "weak":
        raise SystemExit(f"expected current source template placeholder quality to be weak, got {report['selected']['placeholderQuality']}")

    return {
        "slides": len(prs.slides),
        "text_shapes": text_shapes,
        "picture_shapes": picture_shapes,
        "replaced_text_found": replaced_text_found,
        "probe_image_found": probe_image_found,
        "sample_text": sample_text,
    }


def main() -> int:
    report = run_node_probe()
    stats = inspect_output(report)
    print(
        json.dumps(
            {
                "status": "ok",
                "library": report["library"],
                "libraryVersion": report["libraryVersion"],
                "selected": report["selected"],
                "output": report["probe"]["output"],
                "stats": stats,
                "template_readiness": [item["readiness"] for item in report["templates"]],
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
