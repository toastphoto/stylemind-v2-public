#!/usr/bin/env python3
"""Verify that PptxGenJS can render the shared StyleMind RenderPlan contract."""

from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
WEB_UI = ROOT / "web_ui"
SCRIPTS = ROOT / "scripts"
for path in (WEB_UI, SCRIPTS):
    if str(path) not in sys.path:
        sys.path.insert(0, str(path))

from export_render_plan_fixture import sample_outline  # noqa: E402
from services.renderers.pptxgenjs_renderer import render_payload_to_pptxgenjs  # noqa: E402
from services.slide_render_plan import deck_render_plan_to_dict  # noqa: E402

PLAN_JSON = ROOT / ".pytest_tmp" / "stylemind_render_plan_fixture.json"
PPTX_OUT = ROOT / ".pytest_tmp" / "stylemind_pptxgenjs_spike.pptx"


def inspect_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    text_shapes = 0
    picture_shapes = 0
    auto_shapes = 0
    full_slide_pictures = 0
    sample_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
                if len(sample_text) < 6:
                    sample_text.append(shape.text.strip()[:80])
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shapes += 1
                if shape.width >= prs.slide_width * 0.9 and shape.height >= prs.slide_height * 0.9:
                    full_slide_pictures += 1
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                auto_shapes += 1
    return {
        "slides": len(prs.slides),
        "text_shapes": text_shapes,
        "picture_shapes": picture_shapes,
        "auto_shapes": auto_shapes,
        "full_slide_pictures": full_slide_pictures,
        "size_kb": round(path.stat().st_size / 1024, 1),
        "sample_text": sample_text,
    }


def main() -> int:
    payload = deck_render_plan_to_dict(sample_outline())
    render_payload_to_pptxgenjs(payload, PPTX_OUT, render_plan_path=PLAN_JSON)

    stats = inspect_pptx(PPTX_OUT)
    if stats["slides"] != 5:
        raise SystemExit(f"expected 5 slides, got {stats['slides']}")
    if stats["text_shapes"] < 20:
        raise SystemExit(f"expected editable text objects, got {stats['text_shapes']}")
    if stats["picture_shapes"] < 2:
        raise SystemExit(f"expected replaceable image objects, got {stats['picture_shapes']}")
    if stats["full_slide_pictures"]:
        raise SystemExit(f"unexpected full-slide picture count: {stats['full_slide_pictures']}")

    print(json.dumps({"status": "ok", "renderer": "pptxgenjs", "output": str(PPTX_OUT), **stats}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
