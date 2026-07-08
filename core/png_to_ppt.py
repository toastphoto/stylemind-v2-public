"""PNG转可编辑PPT"""
import os
from typing import Optional
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from .api_client import APIClient


class PNGToPPT:
    """PNG转可编辑PPT"""

    def __init__(self, api_client: APIClient):
        """
        初始化PNG转PPT转换器

        Args:
            api_client: API客户端（用于调用视觉模型）
        """
        self.api_client = api_client

    def convert(self, png_path: str, model: str = "gpt-4o") -> str:
        """
        将PNG图片转换为可编辑PPT

        流程：
        1. 使用视觉模型识别图片中的元素
        2. 提取文字、形状、位置
        3. 用python-pptx重建PPT

        Args:
            png_path: PNG图片路径
            model: 用于视觉识别的模型

        Returns:
            生成的PPT文件路径

        Raises:
            Exception: 转换失败时抛出异常
        """
        try:
            # 1. 读取图片
            if not os.path.exists(png_path):
                raise Exception(f"图片文件不存在: {png_path}")

            image = Image.open(png_path)
            width, height = image.size

            # 2. 使用视觉模型识别元素
            elements = self._extract_elements(png_path, model)

            if not elements:
                # 如果识别失败，创建空白PPT
                elements = self._create_blank_slide()

            # 3. 创建PPT
            pptx_path = self._create_pptx(elements, width, height, png_path)

            return pptx_path

        except Exception as e:
            raise Exception(f"PNG转PPT失败: {str(e)}")

    def _extract_elements(self, png_path: str, model: str) -> list[dict]:
        """
        使用视觉模型提取图片元素

        Args:
            png_path: 图片路径
            model: 模型

        Returns:
            元素列表
        """
        try:
            # 构建提示词
            prompt = """请分析这张PPT截图，提取所有元素并以JSON格式返回。

要求：
1. 识别所有文本内容和位置
2. 识别主要的内容块和分区
3. 识别图片、图表等视觉元素
4. 估计每个元素的位置（百分比坐标）

返回格式：
{
    "elements": [
        {
            "type": "text",
            "content": "文本内容",
            "bbox": {"x": 0.1, "y": 0.2, "width": 0.8, "height": 0.1},
            "style": {"font_size": "large", "bold": true, "color": "dark"}
        },
        {
            "type": "image",
            "bbox": {"x": 0.5, "y": 0.3, "width": 0.4, "height": 0.4},
            "description": "图片描述"
        }
    ],
    "layout": "描述整体布局",
    "colors": ["主要颜色列表"]
}

请确保bbox坐标是相对于图片尺寸的百分比（0-1之间）。"""

            # 读取图片为base64
            import base64
            with open(png_path, "rb") as f:
                image_data = base64.b64encode(f.read()).decode()

            # 调用视觉模型
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_data}"
                            }
                        }
                    ]
                }
            ]

            response = self.api_client.chat(
                model=model,
                messages=messages,
                max_tokens=2000
            )

            # 解析响应
            import json
            import re

            json_match = re.search(r'\{[\s\S]*\}', response["content"])
            if json_match:
                data = json.loads(json_match.group())
                return data.get("elements", [])

            return []

        except Exception as e:
            print(f"元素提取失败: {e}")
            return []

    def _create_blank_slide(self) -> list[dict]:
        """
        创建空白幻灯片元素

        Returns:
            空白幻灯片元素
        """
        return [
            {
                "type": "text",
                "content": "标题",
                "bbox": {"x": 0.1, "y": 0.1, "width": 0.8, "height": 0.15},
                "style": {"font_size": "large", "bold": true, "color": "dark"}
            },
            {
                "type": "text",
                "content": "内容",
                "bbox": {"x": 0.1, "y": 0.3, "width": 0.8, "height": 0.6},
                "style": {"font_size": "medium", "bold": false, "color": "dark"}
            }
        ]

    def _create_pptx(
        self,
        elements: list[dict],
        img_width: int,
        img_height: int,
        png_path: str
    ) -> str:
        """
        创建PPTX文件

        Args:
            elements: 元素列表
            img_width: 图片宽度
            img_height: 图片高度
            png_path: 原始PNG路径

        Returns:
            PPT文件路径
        """
        # 创建演示文稿（16:9比例）
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 宽度
        prs.slide_height = Inches(7.5)    # 16:9 高度

        # 添加幻灯片
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加背景图片
        self._add_background_image(slide, png_path, prs.slide_width, prs.slide_height)

        # 转换比例（图片到PPT）
        ratio_x = prs.slide_width / img_width
        ratio_y = prs.slide_height / img_height

        # 添加元素
        for element in elements:
            try:
                if element.get("type") == "text":
                    self._add_text_element(slide, element, ratio_x, ratio_y, prs)
                elif element.get("type") == "image":
                    # 图片元素暂时保留为背景
                    pass
            except Exception as e:
                print(f"添加元素失败: {e}")
                continue

        # 保存文件
        output_dir = os.path.dirname(png_path) or "outputs"
        base_name = os.path.splitext(os.path.basename(png_path))[0]
        pptx_path = os.path.join(output_dir, f"{base_name}.pptx")

        prs.save(pptx_path)

        return pptx_path

    def _add_background_image(self, slide, png_path: str, width, height):
        """
        添加背景图片

        Args:
            slide: 幻灯片对象
            png_path: 图片路径
            width: 幻灯片宽度
            height: 幻灯片高度
        """
        try:
            # 添加图片作为背景
            slide.shapes.add_picture(
                png_path,
                0,  # left
                0,  # top
                width=width,
                height=height
            )
        except Exception as e:
            print(f"添加背景图片失败: {e}")

    def _add_text_element(
        self,
        slide,
        element: dict,
        ratio_x: float,
        ratio_y: float,
        prs: Presentation
    ):
        """
        添加文本元素

        Args:
            slide: 幻灯片对象
            element: 元素数据
            ratio_x: X轴比例
            ratio_y: Y轴比例
            prs: 演示文稿
        """
        bbox = element.get("bbox", {})
        content = element.get("content", "")
        style = element.get("style", {})

        # 计算位置和大小
        x = Emu(int(bbox.get("x", 0) * prs.slide_width))
        y = Emu(int(bbox.get("y", 0) * prs.slide_height))
        width = Emu(int(bbox.get("width", 0.8) * prs.slide_width))
        height = Emu(int(bbox.get("height", 0.1) * prs.slide_height))

        # 添加文本框
        textbox = slide.shapes.add_textbox(x, y, width, height)

        # 设置文本
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content

        # 设置样式
        font_size = style.get("font_size", "medium")
        font_sizes = {
            "large": Pt(44),
            "medium": Pt(24),
            "small": Pt(16)
        }
        p.font.size = font_sizes.get(font_size, Pt(24))
        p.font.bold = style.get("bold", False)

        # 颜色
        color = style.get("color", "dark")
        colors = {
            "dark": RGBColor(0, 0, 0),
            "white": RGBColor(255, 255, 255),
            "blue": RGBColor(37, 99, 235)
        }
        p.font.color.rgb = colors.get(color, RGBColor(0, 0, 0))

    def convert_batch(
        self,
        png_paths: list[str],
        output_pptx: str,
        model: str = "gpt-4o"
    ) -> str:
        """
        批量转换PNG到单个PPT

        Args:
            png_paths: PNG文件路径列表
            output_pptx: 输出PPT路径
            model: 模型

        Returns:
            生成的PPT路径
        """
        try:
            # 创建演示文稿
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # 转换每一页
            for i, png_path in enumerate(png_paths):
                if not os.path.exists(png_path):
                    print(f"跳过不存在的文件: {png_path}")
                    continue

                try:
                    elements = self._extract_elements(png_path, model)
                    image = Image.open(png_path)
                    width, height = image.size

                    # 添加幻灯片
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)

                    # 添加背景
                    self._add_background_image(
                        slide,
                        png_path,
                        prs.slide_width,
                        prs.slide_height
                    )

                    # 转换比例
                    ratio_x = prs.slide_width / width
                    ratio_y = prs.slide_height / height

                    # 添加文本元素
                    for element in elements:
                        if element.get("type") == "text":
                            try:
                                self._add_text_element(
                                    slide, element, ratio_x, ratio_y, prs
                                )
                            except Exception as e:
                                print(f"第{i+1}页添加文本失败: {e}")

                    print(f"第{i+1}页处理完成")

                except Exception as e:
                    print(f"第{i+1}页处理失败: {e}")
                    continue

            # 保存
            prs.save(output_pptx)
            print(f"PPT已保存: {output_pptx}")

            return output_pptx

        except Exception as e:
            raise Exception(f"批量转换失败: {str(e)}")
