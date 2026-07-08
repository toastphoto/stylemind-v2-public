"""RAG知识库 - 沉淀PPT和PDF的逻辑和风格"""
import os
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from storage.vector_store import VectorStore
from storage.database import Database
from .api_client import APIClient
from .pdf_parser import PDFParser
from .pdf_splitter import split_large_pdf


class RAGKnowledge:
    """RAG知识库 - 沉淀PPT和PDF的逻辑和风格"""

    def __init__(self, vector_store: VectorStore, db: Database, api_client: Optional[APIClient] = None):
        """
        初始化RAG知识库

        Args:
            vector_store: 向量存储实例
            db: 数据库实例
            api_client: 可选的API客户端用于LLM分析
        """
        self.vector_store = vector_store
        self.db = db
        self.api_client = api_client
        self.pdf_parser = PDFParser()

    def ingest_file(self, file_path: str, model: str = "gpt-4o", progress_callback=None) -> dict:
        """
        导入文件到知识库（自动识别 PPT 或 PDF）

        大PDF会自动拆分处理

        Args:
            file_path: 文件路径
            model: 用于分析的模型
            progress_callback: 进度回调函数

        Returns:
            导入结果
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self.ingest_pdf_with_split(file_path, model)
        elif ext in [".pptx", ".ppt"]:
            return self.ingest_ppt(file_path, model)
        else:
            return {"status": "error", "message": f"不支持的文件格式: {ext}"}

    def ingest_pdf_with_split(self, pdf_path: str, model: str = "gpt-4o") -> dict:
        """
        导入PDF到知识库（自动拆分大文件）

        Args:
            pdf_path: PDF文件路径
            model: 用于分析的模型

        Returns:
            导入结果，包含进度信息
        """
        import time
        start_time = time.time()

        # 进度信息
        progress = {
            "stage": "准备",
            "current": 0,
            "total": 100,
            "message": "正在初始化..."
        }

        try:
            # 1. 检查并拆分大文件
            progress["stage"] = "拆分"
            progress["message"] = "正在分析文件大小..."
            split_files, split_info = split_large_pdf(pdf_path, max_size_mb=40, max_pages=50)

            all_pages = []
            all_insights = []
            total_pages = 0

            # 2. 逐个处理拆分后的文件
            for idx, split_file in enumerate(split_files):
                try:
                    # 解析PDF
                    pages, metadata = self.pdf_parser.parse(split_file)

                    if pages:
                        all_pages.extend(pages)
                        total_pages += len(pages)

                        # LLM分析（API必须可用，否则直接报错停止）
                        if self.api_client:
                            insights = self._analyze_pdf_with_llm(pages, model)
                            all_insights.extend(insights)
                except Exception as e:
                    print(f"处理分片 {idx+1} 失败: {e}")
                    continue

            if not all_pages:
                return {"status": "error", "message": "PDF解析失败或为空"}

            # 3. 存储到向量数据库
            progress["stage"] = "存储"
            progress["message"] = f"正在存储 {len(all_pages)} 页到知识库..."

            for i, page in enumerate(all_pages):
                progress["current"] = int((i / len(all_pages)) * 100)
                progress["message"] = f"存储第 {i+1}/{len(all_pages)} 页..."

                content = f"{page.title}\n{page.content}"

                # 获取排版信息
                layout_data = {}
                if hasattr(page, 'layout') and page.layout:
                    layout_data = page.layout.to_dict()

                self.vector_store.add(
                    text=content,
                    metadata={
                        "source": os.path.basename(pdf_path),
                        "page_index": page.page_index,
                        "layout_type": layout_data.get("layout_type", "content"),
                        "layout_columns": layout_data.get("columns", 1),
                        "layout_alignment": layout_data.get("text_alignment", "left"),
                        "layout_density": layout_data.get("content_density", "normal"),
                        "has_tables": len(page.tables) > 0,
                        "has_images": len(page.images) > 0,
                        "title_hierarchy": layout_data.get("title_hierarchy", []),
                        "file_type": "pdf",
                        "insights": all_insights[i] if i < len(all_insights) else "",
                    }
                )

                # 记录到数据库（包含排版信息）
                import json
                self.db.add_knowledge(
                    source_file=os.path.basename(pdf_path),
                    page_index=page.page_index,
                    content=content,
                    insights=all_insights[i] if i < len(all_insights) else ""
                )

            progress["stage"] = "完成"
            progress["current"] = 100
            progress["message"] = "导入完成！"

            elapsed = time.time() - start_time

            return {
                "status": "success",
                "pages": len(all_pages),
                "total_pages": total_pages,
                "split_info": split_info,
                "metadata": {"parser": "PyPDF2"},
                "elapsed_seconds": int(elapsed),
                "progress": progress,
            }
        except Exception as e:
            return {"status": "error", "message": str(e)}

    def ingest_pdf(self, pdf_path: str, model: str = "gpt-4o", progress_callback=None) -> dict:
        """
        导入PDF到知识库

        流程：
        1. 解析PDF提取每页内容
        2. LLM分析：提炼内容模板、分页规律、排版逻辑、风格特征
        3. 存入向量数据库

        Args:
            pdf_path: PDF文件路径
            model: 用于分析的模型
            progress_callback: 进度回调函数 (current, total, message)

        Returns:
            导入结果
        """
        import time
        start_time = time.time()

        try:
            if progress_callback:
                progress_callback(0, 100, "正在解析PDF...")

            # 1. 解析PDF（带超时保护）
            pages, metadata = self.pdf_parser.parse(pdf_path)

            if not pages:
                return {"status": "error", "message": "PDF解析失败或为空"}

            total_pages = len(pages)

            # 大文件限制：超过50页只分析前50页
            if total_pages > 50:
                pages = pages[:50]
                if progress_callback:
                    progress_callback(10, 100, f"PDF共{total_pages}页，只分析前50页...")
            else:
                if progress_callback:
                    progress_callback(10, 100, f"PDF解析完成，共{total_pages}页...")

            # 2. LLM分析（如果提供了API客户端）
            insights = []
            if self.api_client:
                insights = self._analyze_pdf_with_llm(pages, model, progress_callback)

            # 3. 存储到向量数据库
            for i, page in enumerate(pages):
                if progress_callback:
                    progress = 50 + int((i / len(pages)) * 50)
                    progress_callback(progress, 100, f"正在存储第{i+1}/{len(pages)}页...")

                content = f"{page.title}\n{page.content}"

                self.vector_store.add(
                    text=content,
                    metadata={
                        "source": os.path.basename(pdf_path),
                        "page_index": page.page_index,
                        "layout_type": page.layout_type,
                        "has_tables": len(page.tables) > 0,
                        "has_images": len(page.images) > 0,
                        "insights": insights[page.page_index] if page.page_index < len(insights) else "",
                        "file_type": "pdf",
                    }
                )

                # 记录到数据库
                self.db.add_knowledge(
                    source_file=os.path.basename(pdf_path),
                    page_index=page.page_index,
                    content=content,
                    insights=insights[page.page_index] if page.page_index < len(insights) else "",
                    metadata={
                        "layout_type": page.layout_type,
                        "tables": page.tables,
                        "file_type": "pdf",
                    }
                )

            elapsed = time.time() - start_time

            return {
                "status": "success",
                "pages": len(pages),
                "total_pages": total_pages,
                "insights": insights,
                "metadata": metadata,
                "elapsed_seconds": int(elapsed),
            }
        except Exception as e:
            return {"status": "error", "message": str(e)}

    def ingest_ppt(self, ppt_path: str, model: str = "gpt-4o") -> dict:
        """
        导入PPT到知识库

        流程：
        1. 解析PPT提取每页内容
        2. LLM分析：提炼内容模板、分页规律、排版逻辑、风格特征
        3. 存入向量数据库

        Args:
            ppt_path: PPT文件路径
            model: 用于分析的模型

        Returns:
            导入结果 {"status": "success", "pages": int, "insights": str}
        """
        try:
            # 1. 解析PPT
            slides_content = self._extract_slides_content(ppt_path)

            if not slides_content:
                return {"status": "error", "message": "PPT解析失败或为空"}

            # 2. LLM分析（API必须可用，否则直接报错停止）
            insights = []
            if self.api_client:
                insights = self._analyze_with_llm(slides_content, model)

            # 3. 存储到向量数据库
            for idx, content in enumerate(slides_content):
                self.vector_store.add(
                    text=content,
                    metadata={
                        "source": os.path.basename(ppt_path),
                        "page_index": idx,
                        "insights": insights[idx] if idx < len(insights) else "",
                        "file_type": "ppt",
                        "layout_type": "content",
                    }
                )

                # 记录到数据库
                self.db.add_knowledge(
                    source_file=os.path.basename(ppt_path),
                    page_index=idx,
                    content=content,
                    insights=insights[idx] if idx < len(insights) else "",
                )

            return {
                "status": "success",
                "pages": len(slides_content),
                "insights": insights,
            }
        except Exception as e:
            return {"status": "error", "message": str(e)}

    def _extract_slides_content(self, ppt_path: str) -> list[str]:
        """
        提取PPT每页的文本内容

        Args:
            ppt_path: PPT文件路径

        Returns:
            每页内容的列表
        """
        contents = []
        try:
            prs = Presentation(ppt_path)
            for idx, slide in enumerate(prs.slides):
                slide_text = []

                # 提取形状文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                # 提取备注
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    slide_text.append(f"[备注] {slide.notes_slide.notes_text_frame.text}")

                content = "\n".join(slide_text)
                contents.append(content)

        except Exception as e:
            print(f"PPT解析错误: {e}")

        return contents

    def _analyze_with_llm(self, slides_content: list[str], model: str) -> list[str]:
        """
        使用LLM分析PPT内容

        Args:
            slides_content: 幻灯片内容列表
            model: 模型名称

        Returns:
            每页的分析洞察
        """
        if not self.api_client:
            return []

        insights = []
        consecutive_failures = 0
        max_failures = 2

        for idx, content in enumerate(slides_content):
            try:
                messages = [
                    {
                        "role": "system",
                        "content": """你是一个专业的PPT设计专家。请分析以下PPT页面内容，提炼：
1. 内容模板（这页属于什么类型：标题页、目录页、内容页、总结页等）
2. 分页规律（这页是如何组织的：上下结构、左右结构、卡片式等）
3. 排版逻辑（文字和图片如何布局）
4. 风格特征（配色、字体、视觉元素）

请用简洁的语言总结这些特征。"""
                    },
                    {
                        "role": "user",
                        "content": f"第{idx + 1}页内容：\n{content}"
                    }
                ]

                response = self.api_client.chat(model=model, messages=messages, max_tokens=500, timeout=10)
                insights.append(response["content"])
                consecutive_failures = 0

            except Exception as e:
                consecutive_failures += 1
                print(f"分析第{idx + 1}页失败: {e}")
                insights.append("")

                if consecutive_failures >= max_failures:
                    raise RuntimeError(
                        f"API 连续失败 {max_failures} 次，请检查 API 配置后重试！\n"
                        f"错误信息: {e}"
                    )

        return insights

    def _analyze_pdf_with_llm(self, pages: list, model: str, progress_callback=None) -> list[str]:
        """
        使用LLM深度分析PDF排版和内容

        Args:
            pages: PDF页面列表
            model: 模型名称
            progress_callback: 进度回调函数

        Returns:
            每页的分析洞察
        """
        import time
        import json

        if not self.api_client:
            return []

        insights = []
        consecutive_failures = 0
        max_consecutive_failures = 5  # 连续失败5次才停止
        max_retries = 3  # 每页最多重试3次
        sensitive_keywords = ['sensitive_words', 'sensitive words', 'content_filter']

        for i, page in enumerate(pages):
            success = False

            for retry in range(max_retries):
                try:
                    if progress_callback:
                        progress = 20 + int((i / len(pages)) * 30)
                        retry_info = f" (重试{retry+1})" if retry > 0 else ""
                        progress_callback(progress, 100, f"正在分析第{i+1}/{len(pages)}页...{retry_info}")

                    # 构建排版信息摘要
                    layout_summary = ""
                    if hasattr(page, 'get_layout_summary'):
                        layout_summary = page.get_layout_summary()
                    elif hasattr(page, 'layout') and page.layout:
                        layout = page.layout
                        layout_summary = f"布局类型: {layout.layout_type}\n"
                        layout_summary += f"分栏数: {layout.columns}\n"
                        layout_summary += f"文本对齐: {layout.text_alignment}\n"
                        layout_summary += f"内容密度: {layout.content_density}\n"
                        layout_summary += f"图片数: {len(page.images)}\n"
                        layout_summary += f"表格数: {len(page.tables)}\n"

                    content = page.content[:2000] if page.content else ""

                    messages = [
                        {
                            "role": "system",
                            "content": """你是一个专业的文档设计专家。请深度分析以下PDF页面，提炼排版设计模式：

请按以下格式输出（JSON）：
{
  "page_type": "封面/目录/正文/图表页/总结页",
  "layout_pattern": "上下结构/左右分栏/图文混排/卡片式/自由布局",
  "title_style": {
    "font_size_ratio": "标题与正文的字号比",
    "position": "居中/左对齐/右对齐",
    "style": "加粗/常规/斜体"
  },
  "content_structure": "单栏/双栏/三栏/混合",
  "visual_elements": ["图片位置", "图标使用", "线条装饰"],
  "spacing": "宽松/适中/紧凑",
  "color_tone": "商务蓝/简约白/科技感/暖色调",
  "design_rules": ["规则1", "规则2", "规则3"],
  "ppt_suggestion": "如果将此页转为PPT，建议的排版方式"
}

请确保输出是有效的JSON格式。"""
                        },
                        {
                            "role": "user",
                            "content": f"第{page.page_index + 1}页排版信息：\n{layout_summary}\n\n文字内容（前2000字）：\n{content}"
                        }
                    ]

                    # 带超时保护的API调用（30秒超时）
                    response = self.api_client.chat(model=model, messages=messages, max_tokens=800, timeout=30)
                    insights.append(response["content"])
                    consecutive_failures = 0  # 成功则重置
                    success = True
                    break  # 成功，跳出重试循环

                except Exception as e:
                    error_str = str(e)

                    # 敏感词错误 → 跳过此页，不重试
                    if any(kw in error_str.lower() for kw in sensitive_keywords):
                        print(f"分析第{page.page_index + 1}页: 触发敏感词，跳过")
                        insights.append("")
                        consecutive_failures = 0  # 敏感词不算连续失败
                        success = True
                        break

                    # 其他错误 → 重试
                    print(f"分析第{page.page_index + 1}页失败 (重试{retry+1}/{max_retries}): {error_str[:100]}")
                    if retry < max_retries - 1:
                        time.sleep(2 * (retry + 1))  # 递增等待

            if not success:
                consecutive_failures += 1
                print(f"分析第{page.page_index + 1}页最终失败，跳过")
                insights.append("")

                # 连续失败达到上限，停止分析（但保留已分析的结果）
                if consecutive_failures >= max_consecutive_failures:
                    print(f"⚠️ API 连续失败 {max_consecutive_failures} 次，停止 LLM 分析")
                    # 填充剩余页面为空
                    for j in range(i + 1, len(pages)):
                        insights.append("")
                    break

            # 每页分析间隔（避免频率限制）
            if i < len(pages) - 1:
                time.sleep(1)

        success_count = sum(1 for ins in insights if ins and len(ins) > 10)
        print(f"[INFO] LLM 分析完成: {success_count}/{len(pages)} 页成功")

        return insights

    def query(self, user_query: str, k: int = 5) -> list[dict]:
        """
        检索相关知识

        Args:
            user_query: 用户查询
            k: 返回数量

        Returns:
            相关知识列表 [{"text": "...", "metadata": {...}, "score": float}]
        """
        results = self.vector_store.search(query=user_query, k=k)
        return results

    def get_all_knowledge(self) -> list[dict]:
        """
        获取所有知识库内容

        Returns:
            知识列表
        """
        return self.db.get_all_knowledge()

    def clear(self) -> dict:
        """
        清空知识库

        Returns:
            结果
        """
        self.vector_store.clear()
        self.db.clear_knowledge()
        return {"status": "success", "message": "知识库已清空"}
